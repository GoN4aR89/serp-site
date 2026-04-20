#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Класс для сравнения результатов поисковой выдачи (SERP) с анализом тональности
Версия для веб-приложения (убраны telegram-зависимости)

СТРУКТУРА ФАЙЛА:
- ЯДРО СРАВНЕНИЯ (секция 6) - критичные методы, требуют тестирования
- Остальные секции - вспомогательные методы, безопасно изменять

Версия: 3.20.1
"""

import os
import logging
import pandas as pd
import re
import matplotlib.pyplot as plt
import numpy as np
import math
from io import BytesIO
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
import matplotlib
from datetime import datetime
import traceback
from matplotlib.ticker import MultipleLocator
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw
from pptx.enum.text import PP_ALIGN

def add_rounded_corners_to_image(image_path, radius=50):
    """Adds rounded corners to an image and returns it as a BytesIO buffer"""
    img = Image.open(image_path)
    img = img.convert("RGBA")
    
    # Create a mask with rounded corners
    width, height = img.size
    mask = Image.new("L", (width, height), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle([(0, 0), (width, height)], radius=radius, fill=255)
    
    # Create a new image with the background color
    background_color = (245, 246, 248, 255)  # #f5f6f8 with alpha
    result = Image.new("RGBA", (width, height), background_color)
    
    # Paste the original image using the mask
    result.paste(img, (0, 0), mask)
    
    # Convert back to RGB and save to buffer
    result = result.convert("RGB")
    buffer = BytesIO()
    result.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from constants import (
    SENTIMENT_COLOR_VARIANTS,
    SENTIMENT_CHART_COLORS,
    SENTIMENT_NAMES_RU,
    SENTIMENT_EMOJI
)

matplotlib.use('Agg')  # Для работы без GUI

logger = logging.getLogger(__name__)

class SERPComparatorWithSentiment:
    """Класс для сравнения результатов поисковой выдачи с анализом тональности"""

    def __init__(self):
        # Используем константы из constants.py
        self.sentiment_colors = SENTIMENT_COLOR_VARIANTS
        self.sentiment_chart_colors = SENTIMENT_CHART_COLORS
        self.sentiment_names = SENTIMENT_NAMES_RU
        self.sentiment_emoji = SENTIMENT_EMOJI

    # ============================================
    # 🔧 СЕКЦИЯ 1: ВСПОМОГАТЕЛЬНЫЕ МЕТОДЫ
    # Безопасно изменять - форматирование, валидация
    # ============================================

    def is_valid_url(self, text):
        """Упрощенная проверка URL - принимаем почти все что похоже на URL"""
        if not text or not isinstance(text, str):
            return False

        text = text.strip()
        if len(text) < 5:
            return False

        text_lower = text.lower()

        patterns = [
            r'https?://',  # http/https
            r'www\.',      # www
            r'\.\w{2,}',   # доменная зона
            r'xn--',       # punycode
            r'/[a-zA-Z0-9]', # слэш с текстом после
        ]

        for pattern in patterns:
            if re.search(pattern, text_lower):
                return True

        if re.search(r'[a-zA-Z0-9-]+\.[a-zA-Z]{2,}', text_lower):
            return True

        if '/' in text and '.' in text:
            return True

        special_patterns = [
            'yandex.ru', 'cian.ru', 'novostroy-m.ru', '2gis.ru',
            'pronovostroy.ru', 'mskguru.ru', 'realty.yandex.ru',
            'whitewill.ru', 'domclick.ru', 'avaho.ru', 'zoon.ru',
            'stroiki.ru', 'otzovik.com', 'avito.ru', 'gdeetotdom.ru',
            'mangazeya.ru', 'novostroev.ru', 'msk.restate.ru',
            'domkad.ru', 'archi.ru', 'cian.ru', 'novostroy.ru',
            'xn--p1ai'  # punycode для .рф
        ]

        for pattern in special_patterns:
            if pattern in text_lower:
                return True

        return False

    def get_cell_color(self, cell, workbook):
        """Получение цвета заливки ячейки"""
        try:
            if cell.fill and cell.fill.start_color:
                color_obj = cell.fill.start_color

                if color_obj.type == 'rgb' and color_obj.rgb and color_obj.rgb != '00000000':
                    color_hex = color_obj.rgb
                else:
                    color_hex = None

                if (not color_hex or color_hex == '00000000') and cell.fill.fgColor and cell.fill.fgColor.rgb and cell.fill.fgColor.rgb != '00000000':
                    color_hex = cell.fill.fgColor.rgb

                if color_hex:
                    color_str = str(color_hex).upper()
                    if len(color_str) == 8:  # AARRGGBB
                        return color_str[2:]
                    elif len(color_str) == 6:  # RRGGBB
                        return color_str
        except Exception as e:
            logger.debug(f"Ошибка получения цвета ячейки {cell.coordinate}: {e}")

        return None

    def determine_sentiment(self, color):
        """Определение тональности по цвету"""
        if not color:
            return 'unknown'

        color = color.upper()

        sentiment_priority = ['client_site', 'positive', 'negative', 'neutral', 'irrelevant']
        for sentiment_type in sentiment_priority:
            if sentiment_type in self.sentiment_colors:
                for sentiment_color in self.sentiment_colors[sentiment_type]:
                    if color == sentiment_color.upper():
                        return sentiment_type

        try:
            r = int(color[0:2], 16)
            g = int(color[2:4], 16)
            b = int(color[4:6], 16)

            if r > 200 and g > 200 and b < 100:
                return 'client_site'
            elif r > 180 and g > 180 and b < 150 and (r - b > 50) and (g - b > 50):
                return 'client_site'

            if g > r and g > b and g > 150 and (g - r > 50) and (g - b > 50):
                return 'positive'
            elif r > g and r > b and r > 150 and (r - g > 50) and (r - b > 50):
                return 'negative'
            elif b > r and b > g and b > 150 and (b - r > 50) and (b - g > 50):
                return 'neutral'
            elif b > 100 and g > 100 and r < 100 and abs(b - g) < 50:
                return 'neutral'
            elif abs(r - g) < 30 and abs(g - b) < 30 and abs(r - b) < 30 and r > 100:
                return 'irrelevant'
        except ValueError:
            pass

        return 'unknown'

    # ============================================
    # 📊 СЕКЦИЯ 2: ДИАГРАММЫ И ВИЗУАЛИЗАЦИЯ
    # Безопасно изменять - внешний вид графиков
    # ============================================

    def create_comparison_chart(self, stats1, stats2, label1, label2, title):
        """
        Создаёт сравнительную столбчатую диаграмму тональности для двух файлов.
        """
        try:
            plt.rcParams['font.family'] = 'DejaVu Sans'
            plt.rcParams['font.size'] = 10

            categories = ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']
            cat_names = [self.sentiment_names.get(c, c) for c in categories]

            percentages1 = [stats1.get(c, {}).get('percentage', 0) for c in categories]
            percentages2 = [stats2.get(c, {}).get('percentage', 0) for c in categories]
            counts1 = [stats1.get(c, {}).get('count', 0) for c in categories]
            counts2 = [stats2.get(c, {}).get('count', 0) for c in categories]

            if sum(percentages1) == 0 and sum(percentages2) == 0:
                return None

            x = np.arange(len(categories))
            width = 0.35

            max_percent = max(max(percentages1), max(percentages2))
            max_rounded = math.ceil(max_percent / 10) * 10
            if max_rounded < 10:
                max_rounded = 10

            fig, ax = plt.subplots(figsize=(8, 4))

            category_colors = [self.sentiment_chart_colors.get(c, '#808080') for c in categories]

            # Первый файл – сплошные столбцы
            bars1 = ax.bar(x - width/2, percentages1, width, label=label1,
                           color=category_colors, edgecolor='white', linewidth=1, alpha=0.9)

            # Второй файл – усиленная штриховка
            bars2 = ax.bar(x + width/2, percentages2, width, label=label2,
                           color=category_colors, edgecolor='white', linewidth=1, alpha=0.9, hatch='///')

            # Подписи значений
            for bars, percentages, counts in [(bars1, percentages1, counts1), (bars2, percentages2, counts2)]:
                for bar, perc, cnt in zip(bars, percentages, counts):
                    height = bar.get_height()
                    if height > 0:
                        ax.annotate(f'{perc}%\n({cnt})',
                                    xy=(bar.get_x() + bar.get_width() / 2, height),
                                    xytext=(0, 3),
                                    textcoords="offset points",
                                    ha='center', va='bottom', fontsize=7, fontweight='bold',
                                    bbox=dict(boxstyle="round,pad=0.2", facecolor='white', alpha=0.7))

            ax.set_ylabel('Процент URL', fontsize=9)
            ax.set_xticks(x)
            ax.set_xticklabels(cat_names, rotation=45, ha='right', fontsize=8)

            # Легенда с явным различием
            from matplotlib.patches import Patch
            legend1 = Patch(facecolor='#6fa8dc', edgecolor='black', label=label1)
            legend2 = Patch(facecolor='#6fa8dc', hatch='///', edgecolor='black', label=label2)
            ax.legend(handles=[legend1, legend2], loc='upper center', bbox_to_anchor=(0.5, 1.25),
                      ncol=2, frameon=True, fancybox=True, shadow=True, fontsize=7)

            ax.set_ylim(0, max_rounded)
            ax.yaxis.set_major_locator(MultipleLocator(10))
            ax.yaxis.set_minor_locator(MultipleLocator(2))
            ax.grid(axis='y', which='major', linestyle='--', alpha=0.5)
            ax.grid(axis='y', which='minor', linestyle=':', alpha=0.2)
            ax.set_axisbelow(True)

            for spine in ['top', 'right']:
                ax.spines[spine].set_visible(False)

            plt.tight_layout()
            plt.subplots_adjust(top=0.80)  # освобождаем место для легенды

            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buffer.seek(0)
            plt.close()
            return buffer

        except Exception as e:
            logger.error(f"Ошибка создания сравнительной диаграммы: {e}")
            return None

    def create_horizontal_chart(self, stats1, stats2, label1, label2, title, stats_start=None, label_start='Старт'):
        """
        Создаёт горизонтальную сравнительную столбчатую диаграмму тональности.
        Используется для общей (Яндекс+Google) диаграммы.
        Если stats_start передан, создаёт три столбца (старт + файл1 + файл2).
        """
        try:
            plt.rcParams['font.family'] = 'DejaVu Sans'
            plt.rcParams['font.size'] = 10

            categories = ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']
            cat_names = [self.sentiment_names.get(c, c) for c in categories]

            percentages1 = [stats1.get(c, {}).get('percentage', 0) for c in categories]
            percentages2 = [stats2.get(c, {}).get('percentage', 0) for c in categories]
            counts1 = [stats1.get(c, {}).get('count', 0) for c in categories]
            counts2 = [stats2.get(c, {}).get('count', 0) for c in categories]

            # Данные стартовых метрик (если есть)
            if stats_start:
                percentages_start = [stats_start.get(c, {}).get('percentage', 0) for c in categories]
                counts_start = [stats_start.get(c, {}).get('count', 0) for c in categories]
            else:
                percentages_start = [0] * len(categories)
                counts_start = [0] * len(categories)

            if sum(percentages1) == 0 and sum(percentages2) == 0 and sum(percentages_start) == 0:
                return None

            y = np.arange(len(categories))
            has_start = stats_start is not None and sum(percentages_start) > 0
            height = 0.25 if has_start else 0.35

            all_percentages = [p for p in (percentages1 + percentages2 + percentages_start) if p > 0]
            max_percent = max(all_percentages) if all_percentages else 0
            max_rounded = math.ceil(max_percent / 10) * 10
            if max_rounded < 10:
                max_rounded = 10

            fig, ax = plt.subplots(figsize=(8, 4))

            category_colors = [self.sentiment_chart_colors.get(c, '#808080') for c in categories]

            if has_start:
                # Три столбца: старт, файл1, файл2 (сверху вниз после invert_yaxis)
                # Старт (y - height) -> после invert_yaxis будет сверху
                # Файл 1 (y) -> после invert_yaxis будет в центре
                # Файл 2 (y + height) -> после invert_yaxis будет снизу
                bars_start = ax.barh(y - height, percentages_start, height, label=label_start,
                                    color=category_colors, edgecolor='white', linewidth=1, alpha=0.5)
                bars1 = ax.barh(y, percentages1, height, label=label1,
                               color=category_colors, edgecolor='white', linewidth=1, alpha=0.9)
                bars2 = ax.barh(y + height, percentages2, height, label=label2,
                               color=category_colors, edgecolor='white', linewidth=1, alpha=0.9, hatch='///')

                # Подписи для всех трех
                for bars, percentages, counts in [(bars_start, percentages_start, counts_start),
                                                  (bars1, percentages1, counts1),
                                                  (bars2, percentages2, counts2)]:
                    for bar, perc, cnt in zip(bars, percentages, counts):
                        width = bar.get_width()
                        if width > 0:
                            ax.annotate(f'{perc}% ({cnt})',
                                        xy=(width, bar.get_y() + bar.get_height() / 2),
                                        xytext=(5, 0),
                                        textcoords="offset points",
                                        ha='left', va='center', fontsize=6, fontweight='bold',
                                        bbox=dict(boxstyle="round,pad=0.2", facecolor='white', alpha=0.7))

                # Легенда для трех столбцов
                from matplotlib.patches import Patch
                legend_start = Patch(facecolor='#6fa8dc', edgecolor='black', alpha=0.5, label=label_start)
                legend1 = Patch(facecolor='#6fa8dc', edgecolor='black', label=label1)
                legend2 = Patch(facecolor='#6fa8dc', hatch='///', edgecolor='black', label=label2)
                ax.legend(handles=[legend_start, legend1, legend2], loc='upper center', bbox_to_anchor=(0.5, 1.15),
                          ncol=3, frameon=True, fancybox=True, shadow=True, fontsize=7)
            else:
                # Два столбца как раньше
                bars1 = ax.barh(y - height/2, percentages1, height, label=label1,
                               color=category_colors, edgecolor='white', linewidth=1, alpha=0.9)
                bars2 = ax.barh(y + height/2, percentages2, height, label=label2,
                               color=category_colors, edgecolor='white', linewidth=1, alpha=0.9, hatch='///')

                # Подписи для двух
                for bars, percentages, counts in [(bars1, percentages1, counts1), (bars2, percentages2, counts2)]:
                    for bar, perc, cnt in zip(bars, percentages, counts):
                        width = bar.get_width()
                        if width > 0:
                            ax.annotate(f'{perc}% ({cnt})',
                                        xy=(width, bar.get_y() + bar.get_height() / 2),
                                        xytext=(5, 0),
                                        textcoords="offset points",
                                        ha='left', va='center', fontsize=7, fontweight='bold',
                                        bbox=dict(boxstyle="round,pad=0.2", facecolor='white', alpha=0.7))

                # Легенда для двух столбцов
                from matplotlib.patches import Patch
                legend1 = Patch(facecolor='#6fa8dc', edgecolor='black', label=label1)
                legend2 = Patch(facecolor='#6fa8dc', hatch='///', edgecolor='black', label=label2)
                ax.legend(handles=[legend1, legend2], loc='upper center', bbox_to_anchor=(0.5, 1.15),
                          ncol=2, frameon=True, fancybox=True, shadow=True, fontsize=7)

            ax.set_xlabel('Процент URL', fontsize=9)
            ax.set_yticks(y)
            ax.set_yticklabels(cat_names, fontsize=8)
            ax.invert_yaxis()  # Домашний сайт сверху

            ax.set_xlim(0, max_rounded)
            ax.xaxis.set_major_locator(MultipleLocator(10))
            ax.xaxis.set_minor_locator(MultipleLocator(2))
            ax.grid(axis='x', which='major', linestyle='--', alpha=0.5)
            ax.grid(axis='x', which='minor', linestyle=':', alpha=0.2)
            ax.set_axisbelow(True)

            for spine in ['top', 'right']:
                ax.spines[spine].set_visible(False)

            plt.tight_layout()
            plt.subplots_adjust(top=0.85)

            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buffer.seek(0)
            plt.close()
            return buffer

        except Exception as e:
            logger.error(f"Ошибка создания горизонтальной диаграммы: {e}")
            return None

    def create_three_column_chart(self, stats_start, stats1, stats2, label_start='Старт', label1='Файл 1', label2='Файл 2', title=''):
        """
        Создаёт вертикальную столбчатую диаграмму с тремя столбцами (старт + 2 файла).
        """
        try:
            plt.rcParams['font.family'] = 'DejaVu Sans'
            plt.rcParams['font.size'] = 10

            categories = ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']
            cat_names = [self.sentiment_names.get(c, c) for c in categories]

            percentages_start = [stats_start.get(c, {}).get('percentage', 0) for c in categories]
            percentages1 = [stats1.get(c, {}).get('percentage', 0) for c in categories]
            percentages2 = [stats2.get(c, {}).get('percentage', 0) for c in categories]
            counts_start = [stats_start.get(c, {}).get('count', 0) for c in categories]
            counts1 = [stats1.get(c, {}).get('count', 0) for c in categories]
            counts2 = [stats2.get(c, {}).get('count', 0) for c in categories]

            if sum(percentages_start) == 0 and sum(percentages1) == 0 and sum(percentages2) == 0:
                return None

            x = np.arange(len(categories))
            width = 0.25

            max_percent = max(max(percentages_start), max(percentages1), max(percentages2))
            max_rounded = math.ceil(max_percent / 10) * 10
            if max_rounded < 10:
                max_rounded = 10

            fig, ax = plt.subplots(figsize=(8, 4))

            category_colors = [self.sentiment_chart_colors.get(c, '#808080') for c in categories]

            # Стартовые данные (слева) - полупрозрачные
            bars_start = ax.bar(x - width, percentages_start, width, label=label_start,
                               color=category_colors, edgecolor='white', linewidth=1, alpha=0.5)

            # Первый файл (центр) - сплошные
            bars1 = ax.bar(x, percentages1, width, label=label1,
                          color=category_colors, edgecolor='white', linewidth=1, alpha=0.9)

            # Второй файл (справа) - штриховка
            bars2 = ax.bar(x + width, percentages2, width, label=label2,
                          color=category_colors, edgecolor='white', linewidth=1, alpha=0.9, hatch='///')

            # Подписи значений
            for bars, percentages, counts in [(bars_start, percentages_start, counts_start),
                                            (bars1, percentages1, counts1),
                                            (bars2, percentages2, counts2)]:
                for bar, perc, cnt in zip(bars, percentages, counts):
                    height = bar.get_height()
                    if height > 0:
                        ax.annotate(f'{perc}%\n({cnt})',
                                    xy=(bar.get_x() + bar.get_width() / 2, height),
                                    xytext=(0, 3),
                                    textcoords="offset points",
                                    ha='center', va='bottom', fontsize=6, fontweight='bold',
                                    bbox=dict(boxstyle="round,pad=0.2", facecolor='white', alpha=0.7))

            ax.set_ylabel('Процент URL', fontsize=9)
            ax.set_xticks(x)
            ax.set_xticklabels(cat_names, rotation=45, ha='right', fontsize=8)
            ax.set_ylim(0, max_rounded)

            # Легенда
            from matplotlib.patches import Patch
            legend_start = Patch(facecolor='#6fa8dc', edgecolor='black', alpha=0.5, label=label_start)
            legend1 = Patch(facecolor='#6fa8dc', edgecolor='black', label=label1)
            legend2 = Patch(facecolor='#6fa8dc', edgecolor='black', hatch='///', label=label2)
            ax.legend(handles=[legend_start, legend1, legend2], loc='upper center', bbox_to_anchor=(0.5, 1.25),
                      ncol=3, frameon=True, fancybox=True, shadow=True, fontsize=7)

            ax.yaxis.set_major_locator(MultipleLocator(10))
            ax.yaxis.set_minor_locator(MultipleLocator(2))
            ax.grid(axis='y', which='major', linestyle='--', alpha=0.5)
            ax.grid(axis='y', which='minor', linestyle=':', alpha=0.2)
            ax.set_axisbelow(True)

            for spine in ['top', 'right']:
                ax.spines[spine].set_visible(False)

            if title:
                ax.set_title(title, fontsize=10, fontweight='bold', pad=10)

            plt.tight_layout()
            plt.subplots_adjust(top=0.80)

            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            buffer.seek(0)
            plt.close()
            return buffer

        except Exception as e:
            logger.error(f"Ошибка в create_three_column_chart: {e}")
            return None

    def create_baseline_chart(self, stats1, stats2, stats_start, title='', file1_name='Файл 1', file2_name='Файл 2'):
        """
        Создаёт baseline диаграмму с тремя столбцами (старт + 2 файла).
        """
        return self.create_three_column_chart(
            stats_start, stats1, stats2,
            label_start='Старт',
            label1=file1_name,
            label2=file2_name,
            title=title
        )

    # ============================================
    # 📝 СЕКЦИЯ 3: ГЕНЕРАЦИЯ ТЕКСТОВЫХ ОТЧЁТОВ
    # Безопасно изменять - форматирование, текст, эмодзи
    # ============================================

    def generate_sentiment_report(self, sentiment_stats, filename, total_urls):
        """Генерация текстового отчета о тональности"""
        report = f"📊 *Анализ тональности для файла:* `{filename}`\n\n"
        report += f"📈 *Всего URL проанализировано:* {total_urls}\n\n"

        for sentiment_type in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant', 'unknown']:
            if sentiment_type in sentiment_stats:
                stats = sentiment_stats[sentiment_type]
                emoji = self.sentiment_emoji.get(sentiment_type, '⚫')
                report += f"{emoji} *{stats['name']}:* {stats['count']} URL ({stats['percentage']}%)\n"

        return report

    # ============================================
    # 📋 СЕКЦИЯ 4: РАСШИРЕННЫЕ ОТЧЁТЫ И АНАЛИТИКА
    # Безопасно изменять - логика анализа, рекомендации
    # ============================================

    def generate_summary_report(self, comparison_df, label1, label2, top_n):
        """Генерация полного текстового отчета по улучшениям и ухудшениям (без аналитического вывода)"""
        try:
            data_df = comparison_df[~comparison_df['Запрос'].str.startswith('СТАТИСТИКА ТОНАЛЬНОСТИ', na=False)]
            data_df = data_df[data_df['Запрос'].notna() & (data_df['Запрос'] != '')]

            improvements = []
            deteriorations = []
            new_urls = []
            dropped_urls = []
            sentiment_changes = []
            
            # SERM-специфичные метрики
            negative_appeared = []  # Новый негатив в выдаче
            negative_disappeared = []  # Ушел негатив из выдачи
            negative_rose = []  # Негатив поднялся в позициях
            positive_disappeared = []  # Позитив ушел

            for _, row in data_df.iterrows():
                change = row['Изменение']
                query = row['Запрос']
                url = row['URL']
                sentiment1 = row[f'Тональность_{label1}']
                sentiment2 = row[f'Тональность_{label2}']
                
                if sentiment1 != sentiment2 and sentiment1 != "Неопределенная" and sentiment2 != "Неопределенная":
                    sentiment_changes.append({
                        'query': query,
                        'url': url,
                        'from': sentiment1,
                        'to': sentiment2
                    })
                
                if change == "Новый в ТОП":
                    new_urls.append(f"{query}: {url} ({sentiment2})")
                    if sentiment2 == 'Негативная':
                        negative_appeared.append(f"{query}: {url}")
                elif change == "Выпал из ТОП":
                    dropped_urls.append(f"{query}: {url} ({sentiment1})")
                    if sentiment1 == 'Негативная':
                        negative_disappeared.append(f"{query}: {url}")
                    elif sentiment1 == 'Позитивная':
                        positive_disappeared.append(f"{query}: {url}")
                elif isinstance(change, (int, float)):
                    if change > 0:
                        improvements.append(f"{query}: {url} (+{int(change)})")
                        # Негатив поднялся (хуже)
                        if sentiment2 == 'Негативная':
                            negative_rose.append(f"{query}: {url} (+{int(change)})")
                    elif change < 0:
                        deteriorations.append(f"{query}: {url} ({int(change)})")

            # Подсчитаем общее число изменений
            total_changes = len(improvements) + len(deteriorations) + len(new_urls) + len(dropped_urls)
            
            # Классифицируем по тональности для репутационного статуса
            new_negative = sum(1 for u in new_urls if 'Негативная' in u)
            dropped_negative = sum(1 for u in dropped_urls if 'Негативная' in u)
            
            report = f"📊 АНАЛИЗ РЕПУТАЦИИ В ВЫДАЧЕ (ТОП-{top_n})\n"
            report += f"{label1} → {label2}\n"
            report += "=" * 70 + "\n\n"

            # СТАТУС РЕПУТАЦИИ
            reputation_status = new_negative + len(negative_rose) - dropped_negative - len(negative_disappeared)
            if reputation_status > 0:
                report += f"⚠️  РЕПУТАЦИЯ: Возникли новые вызовы (+{reputation_status} негативных факторов)\n"
            elif reputation_status < 0:
                report += f"✅ РЕПУТАЦИЯ: Улучшается ({abs(reputation_status)} факторов разрешено)\n"
            else:
                report += f"⚖️  РЕПУТАЦИЯ: Стабильна, внимательно мониторьте\n"
            
            if total_changes == 0:
                report += "   Выдача не изменилась — продолжайте текущую стратегию\n"
            else:
                report += f"   Выявлено {total_changes} изменений в выдаче\n"
            report += "\n"
            
            # КЛЮЧЕВЫЕ СОБЫТИЯ
            if negative_appeared or negative_rose or positive_disappeared:
                report += "📍 КЛЮЧЕВЫЕ СОБЫТИЯ, ТРЕБУЮЩИЕ ВНИМАНИЯ:\n\n"
                
                if negative_appeared:
                    report += f"   ⚠️  Новый негатив ({len(negative_appeared)} сайт):\n"
                    for url in negative_appeared:
                        report += f"      • {url}\n"
                    report += f"      💡 Рекомендация: Ознакомьтесь с этими сайтами, поймите,\n"
                    report += f"         что они представляют. Обычно это отзывы, конкуренты или\n"
                    report += f"         форумы. Подготовьте контент с ответом на основные\n"
                    report += f"         критические точки.\n\n"
                
                if negative_rose:
                    report += f"   📈 Негатив поднялся в позициях ({len(negative_rose)} сайт):\n"
                    for url in negative_rose:
                        report += f"      • {url}\n"
                    report += f"      💡 Рекомендация: Это приоритет. Выпустите материал на\n"
                    report += f"         авторитетных источниках, чтобы оттеснить эти ссылки\n"
                    report += f"         вниз по выдаче. Работайте над PR и ссылочной массой.\n\n"
                
                if positive_disappeared:
                    report += f"   🔍 Ушел позитивный контент ({len(positive_disappeared)} сайт):\n"
                    for url in positive_disappeared:
                        report += f"      • {url}\n"
                    report += f"      💡 Рекомендация: Свяжитесь с владельцами этих сайтов,\n"
                    report += f"         выясните причину. Возможно, произошла техническая ошибка\n"
                    report += f"         или удаление. Восстановите контент или создайте новое\n"
                    report += f"         упоминание на авторитетном ресурсе.\n\n"
            
            if negative_disappeared:
                report += f"   ✅ Решено: Ушел негатив ({len(negative_disappeared)} сайт)\n"
                for url in negative_disappeared:
                    report += f"      • {url}\n\n"
            
            # ПОЛНЫЙ СПИСОК РОТАЦИИ
            report += "─" * 70 + "\n\n"
            report += f"📊 ДЕТАЛЬНЫЙ АНАЛИЗ РОТАЦИИ\n"
            
            # НОВЫЕ В ТОПЕ (развернутый список)
            if new_urls:
                report += f"\n🆕 НОВЫЕ САЙТЫ ({len(new_urls)}):\n"
                sentiment_groups = {}
                for url in new_urls:
                    sentiment = url.split('(')[-1].rstrip(')')
                    if sentiment not in sentiment_groups:
                        sentiment_groups[sentiment] = []
                    sentiment_groups[sentiment].append(url)
                
                for sentiment in ['Позитивная', 'Нейтральная', 'Негативная', 'Нерелевантная']:
                    if sentiment in sentiment_groups:
                        emoji = "🟢" if sentiment == "Позитивная" else "🔴" if sentiment == "Негативная" else "🔵"
                        report += f"\n   {emoji} {sentiment} ({len(sentiment_groups[sentiment])}):\n"
                        for url in sentiment_groups[sentiment]:
                            clean = url.rsplit(' (', 1)[0]
                            report += f"      • {clean}\n"
            else:
                report += f"\n🆕 НОВЫЕ САЙТЫ: Нет, выдача стабильна\n"
            
            # ВЫПАВШИЕ ИЗ ТОПа (развернутый список)
            if dropped_urls:
                report += f"\n❌ ВЫПАВШИЕ САЙТЫ ({len(dropped_urls)}):\n"
                sentiment_groups = {}
                for url in dropped_urls:
                    sentiment = url.split('(')[-1].rstrip(')')
                    if sentiment not in sentiment_groups:
                        sentiment_groups[sentiment] = []
                    sentiment_groups[sentiment].append(url)
                
                for sentiment in ['Позитивная', 'Нейтральная', 'Негативная', 'Нерелевантная']:
                    if sentiment in sentiment_groups:
                        emoji = "🟢" if sentiment == "Позитивная" else "🔴" if sentiment == "Негативная" else "🔵"
                        report += f"\n   {emoji} {sentiment} ({len(sentiment_groups[sentiment])}):\n"
                        for url in sentiment_groups[sentiment]:
                            clean = url.rsplit(' (', 1)[0]
                            report += f"      • {clean}\n"
            else:
                report += f"\n❌ ВЫПАВШИЕ САЙТЫ: Нет\n"
            
            # ПОЗИЦИОННЫЕ СДВИГИ
            report += f"\n"
            if improvements:
                report += f"📈 ПОДНЯЛИ ПОЗИЦИЮ: {len(improvements)} сайтов\n"
                for imp in improvements[:10]:  # показываем первые 10
                    report += f"   • {imp}\n"
                if len(improvements) > 10:
                    report += f"   ... и еще {len(improvements) - 10}\n"
            
            if deteriorations:
                report += f"\n📉 УПАЛИ В ПОЗИЦИЯХ: {len(deteriorations)} сайтов\n"
                for det in deteriorations[:10]:
                    report += f"   • {det}\n"
                if len(deteriorations) > 10:
                    report += f"   ... и еще {len(deteriorations) - 10}\n"
            
            if improvements or deteriorations:
                report += "\n"
            
            # ИЗМЕНЕНИЯ ТОНАЛЬНОСТИ
            if sentiment_changes:
                report += f"🎨 Сайты изменили тональность: {len(sentiment_changes)}\n"
                for change in sentiment_changes:
                    emoji_from = "🟢" if change['from'] == "Позитивная" else "🔴" if change['from'] == "Негативная" else "🔵"
                    emoji_to = "🟢" if change['to'] == "Позитивная" else "🔴" if change['to'] == "Негативная" else "🔵"
                    report += f"   {emoji_from}→{emoji_to}  {change['query']}: {change['url']}\n"
            else:
                report += "🎨 Тональность не изменилась\n"

            return report

        except Exception as e:
            logger.error(f"Ошибка генерации полного отчета: {e}")
            return f"Ошибка при генерации полного отчета: {str(e)}"

    def generate_analysis(self, improvements, deteriorations, new_urls, dropped_urls, sentiment_changes, top_n):
        """Анализ репутационных возможностей с рекомендательным тоном"""
        try:
            improvements = improvements or []
            deteriorations = deteriorations or []
            new_urls = new_urls or []
            dropped_urls = dropped_urls or []
            sentiment_changes = sentiment_changes or []

            analysis = []
            
            # Классифицируем по тональности
            new_negative = sum(1 for u in new_urls if 'Негативная' in u)
            new_positive = sum(1 for u in new_urls if 'Позитивная' in u)
            dropped_negative = sum(1 for u in dropped_urls if 'Негативная' in u)
            dropped_positive = sum(1 for u in dropped_urls if 'Позитивная' in u)
            
            sentiment_to_positive = sum(1 for c in sentiment_changes if '→ Позитивная' in str(c))
            sentiment_to_negative = sum(1 for c in sentiment_changes if '→ Негативная' in str(c))
            
            # Основной вопрос SERM: репутационный риск?
            reputation_risk = new_negative + sentiment_to_negative - dropped_negative - sentiment_to_positive
            
            analysis.append("🔎 АНАЛИЗ И РЕКОМЕНДАЦИИ\n")
            
            # 1. ИНТЕРПРЕТАЦИЯ СИТУАЦИИ
            analysis.append("📈 ОБЩАЯ СИТУАЦИЯ:\n")
            
            if reputation_risk > 2:
                analysis.append(f"• Репутация испытывает давление: обнаружено {reputation_risk} новых негативных факторов")
                analysis.append("• Это требует вашего внимания, но не критичность")
                analysis.append("• У вас есть время подготовить ответ\n")
            elif reputation_risk > 0:
                analysis.append(f"• Небольшое давление на репутацию (+{reputation_risk} негатива)")
                analysis.append("• Ничего страшного, включите профилактику")
                analysis.append("• Обычно это временное явление\n")
            elif reputation_risk < -2:
                analysis.append(f"• Репутация улучшается! Негативных факторов на {abs(reputation_risk)} меньше")
                analysis.append("• Ваши усилия дают результат")
                analysis.append("• Продолжайте текущую стратегию\n")
            elif reputation_risk < 0:
                analysis.append(f"• Небольшое улучшение репутации ({abs(reputation_risk)} факторов разрешено)")
                analysis.append("• Хороший знак, удерживайте курс\n")
            else:
                analysis.append("• Репутация в стабильном состоянии")
                analysis.append("• Никаких срочных проблем, но мониторьте ситуацию\n")
            
            # 2. АНАЛИЗ НОВЫХ ВОЗМОЖНОСТЕЙ И ВЫЗОВОВ
            analysis.append("\n💡 ЧТО ВЫ МОЖЕТЕ СДЕЛАТЬ:\n")
            
            # Новые позитивные
            if new_positive > 0:
                analysis.append(f"✅ Появились новые позитивные источники ({new_positive} сайтов)")
                analysis.append("   Предложение: используйте эти ресурсы для PR и партнерства")
                analysis.append("   • Посетите эти сайты, поймите их аудиторию")
                analysis.append("   • Разместите свою историю, если это уместно")
                analysis.append("   • Добавьте ссылки на эти сайты в свой контент\n")
            
            # Выпал негатив
            if dropped_negative > 0:
                analysis.append(f"✅ Негативные источники ушли из выдачи ({dropped_negative})")
                analysis.append("   Это хороший признак, репутация улучшается")
                analysis.append("   • Продолжайте то, что вы делали")
                analysis.append("   • Не ослабляйте качество контента\n")
            
            # Выпал позитив
            if dropped_positive > 0:
                analysis.append(f"⚠️  Позитивные источники исчезли ({dropped_positive})")
                analysis.append("   Возможно, временный баг или удаление")
                analysis.append("   • Проверьте эти сайты вручную")
                analysis.append("   • Если контент ценный, восстановите его")
                analysis.append("   • Создавайте резервные копии упоминаний\n")
            
            # Новый негатив
            if new_negative > 0:
                analysis.append(f"⚠️  Появились новые негативные источники ({new_negative})")
                analysis.append("   Это может быть отзыв, конкурент или форум")
                analysis.append("   • Тактика: не игнорируйте, но и не паникуйте")
                analysis.append("   • Ознакомьтесь с содержимым этих сайтов")
                analysis.append("   • Подготовьте контент с ответом на критику")
                analysis.append("   • Усильте присутствие позитивного контента\n")
            
            # Тональность
            if sentiment_to_positive > sentiment_to_negative + 1:
                analysis.append(f"📈 Тональность улучшается ({sentiment_to_positive} позитивных переходов)")
                analysis.append("   Это значит, что контент становится лучше восприниматься")
                analysis.append("   • Продолжайте улучшать качество\n")
            elif sentiment_to_negative > sentiment_to_positive + 1:
                analysis.append(f"📉 Тональность ухудшается ({sentiment_to_negative} негативных переходов)")
                analysis.append("   Проверьте, не произошло ли что-то негативное")
                analysis.append("   • Может быть, кампания конкурентов или PR")
                analysis.append("   • Выпустите позитивный пресс-релиз\n")
            
            # 3. ПРАКТИЧЕСКИЙ ПЛАН ДЕЙСТВИЙ
            analysis.append("\n🎯 ПРАКТИЧЕСКИЙ ПЛАН:\n")
            
            priority_actions = []
            
            if reputation_risk > 2:
                priority_actions.append("1. НЕМЕДЛЕННО (сегодня-завтра):")
                priority_actions.append("   • Посмотрите новые негативные сайты")
                priority_actions.append("   • Поймите, что они критикуют")
                priority_actions.append("   • Подумайте, есть ли обоснованные замечания")
                priority_actions.append("")
            
            if dropped_positive > 0 or new_negative > 1:
                priority_actions.append("2. НА НЕДЕЛЮ:")
                if dropped_positive > 0:
                    priority_actions.append("   • Выпустите пресс-релиз о вашем развитии")
                    priority_actions.append("   • Напишите статью для промежуточного сайта")
                if new_negative > 1:
                    priority_actions.append("   • Подготовьте ответ на критику в своем блоге")
                    priority_actions.append("   • Попросите позитивные отзывы от довольных клиентов")
                priority_actions.append("")
            
            if new_positive > 0 or reputation_risk <= 0:
                priority_actions.append("3. НА МЕСЯЦ:")
                priority_actions.append("   • Усильте PR и связи с медиа")
                priority_actions.append("   • Наращивайте упоминания на авторитетных источниках")
                priority_actions.append("   • Улучшайте качество контента и продукта")
                priority_actions.append("   • Мониторьте выдачу еженедельно")
                priority_actions.append("")
            
            for action in priority_actions:
                analysis.append(action)
            
            # 4. ИТОГОВАЯ ОЦЕНКА
            analysis.append("📊 ИТОГОВЫЙ ВЕРДИКТ:\n")
            
            if reputation_risk > 3:
                analysis.append("🟠 ПОВЫШЕННОЕ ВНИМАНИЕ: Есть вызовы, но ничего критичного")
                analysis.append("   Действуйте спокойно, методично, стратегически")
            elif reputation_risk > 0:
                analysis.append("🟡 НОРМАЛЬНО: Репутация требует профилактики")
                analysis.append("   Включите мониторинг, но не паникуйте")
            elif reputation_risk < -1:
                analysis.append("🟢 ХОРОШО: Репутация улучшается")
                analysis.append("   Удерживайте курс, не расслабляйтесь")
            else:
                analysis.append("⚪ СТАБИЛЬНО: Репутация в норме")
                analysis.append("   Продолжайте профилактику и мониторинг")
            
            return "\n".join(analysis)

        except Exception as e:
            logger.error(f"Ошибка SERM анализа: {e}", exc_info=True)
            return "Не удалось сгенерировать анализ репутации"

    # ============================================
    # 🔢 СЕКЦИЯ 5: ОБРАБОТКА ДАННЫХ И СТАТИСТИКА
    # Осторожно - влияет на расчёты метрик
    # ============================================

    def generate_summary_dataframe(self, comparison_df, label1, label2, top_n):
        """Генерация улучшенного Summary отчета в формате DataFrame для Excel"""
        try:
            if comparison_df.empty:
                return pd.DataFrame([{'Метрика': 'Нет данных', 'Значение': 'Нет данных для анализа', 'Детали': ''}]), {}
            
            data_df = comparison_df[~comparison_df['Запрос'].str.startswith('СТАТИСТИКА ТОНАЛЬНОСТИ', na=False)]
            data_df = data_df[data_df['Запрос'].notna() & (data_df['Запрос'] != '')]

            # Собираем данные по категориям
            improvements = []
            deteriorations = []
            new_urls = []
            dropped_urls = []
            sentiment_changes = []
            
            # SERM-метрики
            negative_new = 0
            negative_dropped = 0
            positive_dropped = 0
            
            for _, row in data_df.iterrows():
                change = row['Изменение']
                query = row['Запрос']
                url = row['URL']
                sentiment1 = row[f'Тональность_{label1}']
                sentiment2 = row[f'Тональность_{label2}']
                
                # Изменения тональности
                if sentiment1 != sentiment2 and sentiment1 != "Неопределенная" and sentiment2 != "Неопределенная":
                    sentiment_changes.append({
                        'Запрос': query, 'URL': url, 'Было': sentiment1, 'Стало': sentiment2
                    })
                
                # Новые в ТОП
                if change == "Новый в ТОП":
                    new_urls.append({'Запрос': query, 'URL': url, 'Тональность': sentiment2})
                    if sentiment2 == 'Негативная':
                        negative_new += 1
                
                # Выпавшие из ТОП
                elif change == "Выпал из ТОП":
                    dropped_urls.append({'Запрос': query, 'URL': url, 'Тональность': sentiment1})
                    if sentiment1 == 'Негативная':
                        negative_dropped += 1
                    elif sentiment1 == 'Позитивная':
                        positive_dropped += 1
                
                # Изменения позиций
                elif isinstance(change, (int, float)):
                    if change > 0:
                        improvements.append({
                            'Запрос': query, 'URL': url, 'Тональность': sentiment2, 
                            'Изменение': int(change)
                        })
                    elif change < 0:
                        deteriorations.append({
                            'Запрос': query, 'URL': url, 'Тональность': sentiment2, 
                            'Изменение': int(change)
                        })

            # Подсчитываем изменения позиций по тональности
            negative_improved = sum(1 for u in improvements if u['Тональность'] == 'Негативная')  # Негатив поднялся - плохо
            negative_worsened = sum(1 for u in deteriorations if u['Тональность'] == 'Негативная')  # Негатив опустился - хорошо
            positive_improved = sum(1 for u in improvements if u['Тональность'] == 'Позитивная')  # Позитив поднялся - хорошо
            positive_worsened = sum(1 for u in deteriorations if u['Тональность'] == 'Позитивная')  # Позитив опустился - плохо

            # SERM-индекс: только критичные изменения (негатив/позитив)
            # Ухудшает: новый негатив + негатив поднялся + позитив опустился + ушел позитив
            # Улучшает: негатив ушел + негатив опустился + позитив поднялся + пришел позитив
            serm_negative = negative_new + negative_improved + positive_worsened + positive_dropped
            serm_positive = negative_dropped + negative_worsened + positive_improved + sum(1 for u in new_urls if u['Тональность'] == 'Позитивная')
            serm_score = serm_negative - serm_positive
            
            if serm_score > 2:
                status = '⚠️ Ухудшается'
            elif serm_score > 0:
                status = '🟡 Требует внимания'
            elif serm_score < -1:
                status = '✅ Улучшается'
            else:
                status = '⚖️ Стабильна'

            # Формируем DataFrame
            rows = []
            
            # === СЕКЦИЯ 1: ОБЩАЯ СВОДКА ===
            rows.append({'Метрика': '📊 ОБЩАЯ СВОДКА', 'Значение': '', 'Детали': ''})
            rows.append({'Метрика': 'Статус репутации', 'Значение': status, 'Детали': f'SERM-индекс: {serm_score:+d}'})
            rows.append({'Метрика': 'Всего изменений', 'Значение': len(improvements) + len(deteriorations) + len(new_urls) + len(dropped_urls), 'Детали': f'Повышений: {len(improvements)}, Снижений: {len(deteriorations)}, Новых: {len(new_urls)}, Выпавших: {len(dropped_urls)}'})
            
            # Добавляем пояснение SERM-индекса
            rows.append({'Метрика': 'ℹ️ Как считается SERM-индекс?', 'Значение': '', 'Детали': 'Ухудшает: Новый_негатив + Негатив_поднялся + Позитив_опустился + Позитив_выпал. Улучшает: Негатив_выпал + Негатив_опустился + Позитив_поднялся + Позитив_новый. Итог = Ухудшает - Улучшает'})
            
            if new_urls:
                rows.append({'Метрика': '🆕 Новые URL в ТОП', 'Значение': len(new_urls), 'Детали': f'Негативных: {negative_new}, Позитивных: {sum(1 for u in new_urls if u["Тональность"] == "Позитивная")}, Нейтральных: {sum(1 for u in new_urls if u["Тональность"] == "Нейтральная")}, Нерелевантных: {sum(1 for u in new_urls if u["Тональность"] == "Нерелевантная")}, Домашних: {sum(1 for u in new_urls if u["Тональность"] == "Домашний сайт")}'})
            
            if dropped_urls:
                rows.append({'Метрика': '📤 Выпавшие URL из ТОП', 'Значение': len(dropped_urls), 'Детали': f'Негативных: {negative_dropped}, Позитивных: {positive_dropped}, Нейтральных: {sum(1 for u in dropped_urls if u["Тональность"] == "Нейтральная")}, Нерелевантных: {sum(1 for u in dropped_urls if u["Тональность"] == "Нерелевантная")}, Домашних: {sum(1 for u in dropped_urls if u["Тональность"] == "Домашний сайт")}'})
            
            if sentiment_changes:
                rows.append({'Метрика': '🔄 Изменения тональности', 'Значение': len(sentiment_changes), 'Детали': 'URL сменили тональность'})
            
            # === СЕКЦИЯ 2: ДЕТАЛИЗАЦИЯ ===
            if new_urls:
                rows.append({'Метрика': '', 'Значение': '', 'Детали': ''})  # Пустая строка-разделитель
                rows.append({'Метрика': '🔗 НОВЫЕ URL В ТОП', 'Значение': len(new_urls), 'Детали': ''})
                for u in new_urls:
                    emoji = {'Позитивная': '🟢', 'Негативная': '🔴', 'Нейтральная': '🔵', 'Нерелевантная': '⚪', 'Домашний сайт': '🟡'}.get(u['Тональность'], '⚫')
                    rows.append({'Метрика': '', 'Значение': f"{emoji} {u['Тональность']}", 'Детали': f"{u['Запрос']}: {u['URL']}"})
            
            if dropped_urls:
                rows.append({'Метрика': '', 'Значение': '', 'Детали': ''})
                rows.append({'Метрика': '📤 ВЫПАВШИЕ URL ИЗ ТОП', 'Значение': len(dropped_urls), 'Детали': ''})
                for u in dropped_urls:
                    emoji = {'Позитивная': '🟢', 'Негативная': '🔴', 'Нейтральная': '🔵', 'Нерелевантная': '⚪', 'Домашний сайт': '🟡'}.get(u['Тональность'], '⚫')
                    rows.append({'Метрика': '', 'Значение': f"{emoji} {u['Тональность']}", 'Детали': f"{u['Запрос']}: {u['URL']}"})
            
            if improvements:
                rows.append({'Метрика': '', 'Значение': '', 'Детали': ''})
                rows.append({'Метрика': '📈 ПОВЫШЕНИЕ ПОЗИЦИЙ', 'Значение': len(improvements), 'Детали': 'URL улучшили позиции'})
                # Сортируем по величине изменения (убывание)
                for u in sorted(improvements, key=lambda x: x['Изменение'], reverse=True):
                    emoji = {'Позитивная': '🟢', 'Негативная': '🔴', 'Нейтральная': '🔵', 'Домашний сайт': '🟡'}.get(u['Тональность'], '⚪')
                    rows.append({'Метрика': '', 'Значение': f"+{u['Изменение']} {emoji} {u['Тональность']}", 'Детали': f"{u['Запрос']}: {u['URL']}"})
            
            if deteriorations:
                rows.append({'Метрика': '', 'Значение': '', 'Детали': ''})
                rows.append({'Метрика': '📉 СНИЖЕНИЕ ПОЗИЦИЙ', 'Значение': len(deteriorations), 'Детали': 'URL ухудшили позиции'})
                # Сортируем по величине изменения (возрастание - самые большие потери первыми)
                for u in sorted(deteriorations, key=lambda x: x['Изменение']):
                    emoji = {'Позитивная': '🟢', 'Негативная': '🔴', 'Нейтральная': '🔵', 'Домашний сайт': '🟡'}.get(u['Тональность'], '⚪')
                    rows.append({'Метрика': '', 'Значение': f"{u['Изменение']} {emoji} {u['Тональность']}", 'Детали': f"{u['Запрос']}: {u['URL']}"})
            
            if sentiment_changes:
                rows.append({'Метрика': '', 'Значение': '', 'Детали': ''})
                rows.append({'Метрика': '🔄 СМЕНА ТОНАЛЬНОСТИ', 'Значение': len(sentiment_changes), 'Детали': ''})
                for u in sentiment_changes:
                    rows.append({'Метрика': '', 'Значение': f"{u['Было']} → {u['Стало']}", 'Детали': f"{u['Запрос']}: {u['URL']}"})

            summary_df = pd.DataFrame(rows)
            return summary_df, {}

        except Exception as e:
            logger.error(f"Ошибка генерации Summary DataFrame: {e}")
            return pd.DataFrame([{'Метрика': 'Ошибка', 'Значение': str(e), 'Детали': ''}]), {}

    # ============================================
    # ⚠️  СЕКЦИЯ 6: ЯДРО СРАВНЕНИЯ - КРИТИЧНО
    # Требует тестирования при любых изменениях
    # ============================================
    # Основные методы:
    # - compare_serp_data_with_sentiment: главная функция сравнения
    # - process_excel_file_with_sentiment: обработка входных файлов
    # ============================================

    def process_excel_file_with_sentiment(self, file_content, filename, sheet_index=0, top_n=20):
        """Обработка Excel файла с сохранением порядка и учетом всех URL (включая повторяющиеся)"""
        try:
            workbook = load_workbook(BytesIO(file_content))
            # Выбираем лист по индексу, если он существует, иначе первый
            if sheet_index < len(workbook.worksheets):
                worksheet = workbook.worksheets[sheet_index]
            else:
                worksheet = workbook.worksheets[0]

            sentiment_data = {}
            data_rows = []
            query_order = []

            start_row = 3

            for row_idx, row in enumerate(worksheet.iter_rows(min_row=start_row), start=start_row):
                if not any(cell.value for cell in row):
                    continue

                query = None
                urls_in_row = []

                for col_idx, cell in enumerate(row, start=1):
                    if cell.value:
                        cell_value = str(cell.value).strip()

                        if col_idx == 1 and cell_value.isdigit():
                            continue

                        if col_idx == 2 and cell_value and len(cell_value) > 2:
                            query = cell_value
                            if query not in query_order:
                                query_order.append(query)
                            continue

                        if col_idx >= 3 and col_idx < 3 + top_n and self.is_valid_url(cell_value):
                            position = col_idx - 2
                            url = cell_value
                            urls_in_row.append(url)

                            color = self.get_cell_color(cell, workbook)
                            sentiment = self.determine_sentiment(color)

                            if query:
                                key = f"{query}|{url}|{position}"
                                sentiment_data[key] = {
                                    'sentiment': sentiment,
                                    'color': color,
                                    'sentiment_name': self.sentiment_names.get(sentiment, 'Неопределенная'),
                                    'position': position,
                                    'url': url
                                }

                if query and urls_in_row:
                    data_rows.append([query] + urls_in_row)

            if data_rows:
                max_urls = min(max(len(row) - 1 for row in data_rows), top_n)
                columns = ['Запрос'] + [f'ТОП{i+1}' for i in range(max_urls)]

                ordered_data = []
                for query in query_order:
                    for row in data_rows:
                        if row[0] == query:
                            truncated_row = row[:1] + row[1:1+max_urls]
                            ordered_data.append(truncated_row)
                            break

                df = pd.DataFrame(ordered_data, columns=columns)
            else:
                df = pd.DataFrame(columns=['Запрос'])

            logger.info(f"Файл {filename} обработан: {len(df)} запросов, {len(sentiment_data)} URL с тональностью (ТОП-{top_n})")
            return df, sentiment_data, None

        except Exception as e:
            logger.error(f"Ошибка обработки файла {filename}: {str(e)}")
            logger.error(traceback.format_exc())
            return None, None, f"Ошибка обработки файла: {str(e)}"

    def get_urls_for_query(self, df, query, top_n=20):
        """Вспомогательная функция для извлечения URL для конкретного запроса (включая повторяющиеся)"""
        urls_ordered = []

        if df.empty or 'Запрос' not in df.columns:
            return urls_ordered

        query_rows = df[df['Запрос'] == query]
        if query_rows.empty:
            return urls_ordered

        for i, col in enumerate(df.columns[1:1+top_n]):
            if col in query_rows.columns:
                urls = query_rows[col].dropna().tolist()
                for url in urls:
                    if isinstance(url, str) and self.is_valid_url(url):
                        urls_ordered.append(url)

        return urls_ordered

    def calculate_sentiment_statistics(self, sentiment_data):
        """Расчет статистики тональности (учитываем все URL, включая повторяющиеся)"""
        if not sentiment_data:
            return {}

        total_urls = len(sentiment_data)
        sentiment_counts = {}

        for data in sentiment_data.values():
            sentiment = data['sentiment']
            sentiment_counts[sentiment] = sentiment_counts.get(sentiment, 0) + 1

        sentiment_percentages = {}
        for sentiment, count in sentiment_counts.items():
            percentage = (count / total_urls) * 100 if total_urls > 0 else 0
            sentiment_percentages[sentiment] = {
                'count': count,
                'percentage': round(percentage, 2),
                'name': self.sentiment_names.get(sentiment, 'Неопределенная')
            }

        return sentiment_percentages

    # ============================================
    # ⚠️  ЯДРО СРАВНЕНИЯ - ГЛАВНАЯ ФУНКЦИЯ
    # ============================================

    def compare_serp_data_with_sentiment(self, df1, df2, sentiment1, sentiment2, label1, label2, top_n=20):
        """Функция сравнения двух наборов данных SERP - ВСЕГДА показывает данные из обоих файлов"""
        try:
            results = []

            queries_order = []
            if not df1.empty and 'Запрос' in df1.columns:
                queries_order.extend(df1['Запрос'].tolist())
            if not df2.empty and 'Запрос' in df2.columns:
                for query in df2['Запрос'].tolist():
                    if query not in queries_order:
                        queries_order.append(query)

            for query in queries_order:
                urls_1_ordered = self.get_urls_for_query(df1, query, top_n)
                urls_2_ordered = self.get_urls_for_query(df2, query, top_n)

                logger.info(f"Запрос '{query}': {len(urls_1_ordered)} URL в файле1, {len(urls_2_ordered)} URL в файле2")

                all_urls = list(dict.fromkeys(urls_1_ordered + urls_2_ordered))

                for url in all_urls:
                    positions_1 = [i+1 for i, u in enumerate(urls_1_ordered) if u == url]
                    positions_2 = [i+1 for i, u in enumerate(urls_2_ordered) if u == url]

                    max_occurrences = max(len(positions_1), len(positions_2))

                    for i in range(max_occurrences):
                        pos_1 = positions_1[i] if i < len(positions_1) else None
                        pos_2 = positions_2[i] if i < len(positions_2) else None

                        change = None
                        if pos_1 is not None and pos_2 is not None:
                            change = pos_1 - pos_2
                        elif pos_1 is not None and pos_2 is None:
                            change = "Выпал из ТОП"
                        elif pos_1 is None and pos_2 is not None:
                            change = "Новый в ТОП"

                        sentiment_name_1 = "Неопределенная"
                        sentiment_name_2 = "Неопределенная"

                        if pos_1 is not None:
                            key_1 = f"{query}|{url}|{pos_1}"
                            if key_1 in sentiment1:
                                sentiment_name_1 = sentiment1[key_1]['sentiment_name']

                        if pos_2 is not None:
                            key_2 = f"{query}|{url}|{pos_2}"
                            if key_2 in sentiment2:
                                sentiment_name_2 = sentiment2[key_2]['sentiment_name']

                        results.append({
                            'Запрос': query,
                            'URL': url,
                            f'Позиция_{label1}': pos_1,
                            f'Позиция_{label2}': pos_2,
                            'Изменение': change,
                            f'Тональность_{label1}': sentiment_name_1,
                            f'Тональность_{label2}': sentiment_name_2
                        })

            if not results:
                return pd.DataFrame(), "Не найдено данных для сравнения", None, None

            final_df = pd.DataFrame(results)

            sentiment_stats1 = self.calculate_sentiment_statistics(sentiment1)
            sentiment_stats2 = self.calculate_sentiment_statistics(sentiment2)

            if sentiment_stats1 or sentiment_stats2:
                stats_rows = []

                if sentiment_stats1:
                    stats_rows.append({
                        'Запрос': f'СТАТИСТИКА ТОНАЛЬНОСТИ ({label1})', 'URL': '',
                        f'Позиция_{label1}': '', f'Позиция_{label2}': '', 'Изменение': '',
                        f'Тональность_{label1}': '', f'Тональность_{label2}': ''
                    })

                    for sentiment_type in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant', 'unknown']:
                        if sentiment_type in sentiment_stats1:
                            stats = sentiment_stats1[sentiment_type]
                            emoji = self.sentiment_emoji.get(sentiment_type, '⚫')
                            stats_rows.append({
                                'Запрос': '', 'URL': f"{emoji} {stats['name']}",
                                f'Позиция_{label1}': f"{stats['count']} URL", f'Позиция_{label2}': f"{stats['percentage']}%",
                                'Изменение': '', f'Тональность_{label1}': '', f'Тональность_{label2}': ''
                            })

                if sentiment_stats2:
                    stats_rows.append({
                        'Запрос': f'СТАТИСТИКА ТОНАЛЬНОСТИ ({label2})', 'URL': '',
                        f'Позиция_{label1}': '', f'Позиция_{label2}': '', 'Изменение': '',
                        f'Тональность_{label1}': '', f'Тональность_{label2}': ''
                    })

                    for sentiment_type in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant', 'unknown']:
                        if sentiment_type in sentiment_stats2:
                            stats = sentiment_stats2[sentiment_type]
                            emoji = self.sentiment_emoji.get(sentiment_type, '⚫')
                            stats_rows.append({
                                'Запрос': '', 'URL': f"{emoji} {stats['name']}",
                                f'Позиция_{label1}': f"{stats['count']} URL", f'Позиция_{label2}': f"{stats['percentage']}%",
                                'Изменение': '', f'Тональность_{label1}': '', f'Тональность_{label2}': ''
                            })

                if stats_rows:
                    stats_df = pd.DataFrame(stats_rows)
                    final_df = pd.concat([final_df, stats_df], ignore_index=True)

            logger.info(f"Сравнение завершено: {len(results)} записей (ТОП-{top_n})")
            return final_df, None, sentiment_stats1, sentiment_stats2

        except Exception as e:
            logger.error(f"Ошибка сравнения данных: {str(e)}")
            logger.error(traceback.format_exc())
            return None, f"Ошибка сравнения данных: {str(e)}", None, None

    # ============================================
    # 🔧 СЕКЦИЯ 8: PPTX ГЕНЕРАЦИЯ
    # ============================================

    def create_chart_pptx(self, chart_buffer, search_engine, top_n, label1, label2, analysis_text="", baseline_values=None):
        """Создает PowerPoint слайд с диаграммой"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Пустой слайд
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Устанавливаем фон слайда
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 246, 248)  # #f5f6f8

            # Заголовок слайда
            title_text = "Динамика тональности по месяцам"
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Уменьшаем шрифт, если заголовок слишком длинный
            if len(title_text) > 35:
                title_frame.paragraphs[0].font.size = Pt(18)
            else:
                title_frame.paragraphs[0].font.size = Pt(24)

            # Название поисковика в правом верхнем углу (лого)
            base_dir = os.path.dirname(os.path.abspath(__file__))
            if search_engine == "Яндекс":
                logo_path = os.path.join(base_dir, 'static', 'images', 'yandex_logo.png')
            elif search_engine == "Google":
                logo_path = os.path.join(base_dir, 'static', 'images', 'google_logo.png')
            else:
                logo_path = None
            
            if logo_path and os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(11.5), Inches(0.15), width=Inches(1.25))
            else:
                se_box = slide.shapes.add_textbox(Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.5))
                se_frame = se_box.text_frame
                se_frame.text = search_engine
                se_frame.paragraphs[0].font.size = Pt(18)
                se_frame.paragraphs[0].font.bold = True
                se_frame.paragraphs[0].font.color.rgb = RGBColor(251, 130, 39)  # Оранжевый цвет
                se_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            # Диаграмма слева
            if chart_buffer:
                # Сохраняем PNG во временный файл
                temp_chart_path = os.path.join(os.getcwd(), 'temp_chart.png')
                with open(temp_chart_path, 'wb') as f:
                    f.write(chart_buffer.getvalue())

                # Добавляем закругленные углы к изображению
                rounded_chart_buffer = add_rounded_corners_to_image(temp_chart_path, radius=30)
                
                # Добавляем изображение с закругленными углами
                slide.shapes.add_picture(rounded_chart_buffer, Inches(0.5), Inches(1.0),
                                        width=Inches(7.5), height=Inches(5.5))

                # Удаляем временный файл
                if os.path.exists(temp_chart_path):
                    os.remove(temp_chart_path)

            # Блок "вывод" справа с закругленными углами и оранжевой рамкой
            output_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.5), Inches(1.0),
                Inches(4.333), Inches(5.5)
            )
            output_box.fill.solid()
            output_box.fill.fore_color.rgb = RGBColor(254, 242, 233)  # #fef2e9
            output_box.line.color.rgb = RGBColor(251, 130, 39)  # #fb8227
            output_box.line.width = Pt(2)

            # Текст внутри блока "вывод"
            text_frame = output_box.text_frame
            text_frame.text = "Вывод"
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 107, 53)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

            # Добавляем стартовые значения если есть
            if baseline_values:
                p = text_frame.add_paragraph()
                p.text = "Стартовые значения:"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)

                baseline_text = f"Дом. сайт: {baseline_values.get('client_site', '0')}%, "
                baseline_text += f"Позитивная: {baseline_values.get('positive', '0')}%, "
                baseline_text += f"Нейтральная: {baseline_values.get('neutral', '0')}%, "
                baseline_text += f"Негативная: {baseline_values.get('negative', '0')}%, "
                baseline_text += f"Нерелев.: {baseline_values.get('irrelevant', '0')}%"
                p = text_frame.add_paragraph()
                p.text = baseline_text
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)

            # Добавляем текст анализа
            if analysis_text:
                p = text_frame.add_paragraph()
                p.text = analysis_text
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)

            # Сохраняем в BytesIO
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            return output

        except Exception as e:
            logger.error(f"Ошибка создания PPTX диаграммы: {str(e)}")
            logger.error(traceback.format_exc())
            return None