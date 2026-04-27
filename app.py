import os
from io import BytesIO
import uuid
import shutil
import logging
import secrets
import smtplib
import ssl
import json
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
from flask import Flask, render_template, request, redirect, url_for, session, send_file, flash, send_from_directory, jsonify
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from comparator import SERPComparatorWithSentiment, add_rounded_corners_to_image
import pandas as pd
import matplotlib.pyplot as plt
from openpyxl.styles import Alignment, PatternFill, Font
from openpyxl.utils import get_column_letter
from flask_sqlalchemy import SQLAlchemy
from dotenv import load_dotenv
from utils import validate_file_path, allowed_file, cleanup_old_files, format_excel_headers, apply_sentiment_coloring, set_column_widths, format_summary_sheet, is_real_query, is_valid_for_second_file
from constants import SENTIMENT_COLORS_RU, SENTIMENT_COLORS_EN, ALLOWED_EXTENSIONS

load_dotenv()

logger = logging.getLogger(__name__)

app = Flask(__name__)

# Загрузка конфигурации из переменных окружения
app.secret_key = os.getenv('SECRET_KEY', 'dev-key-change-in-production')
app.config['MAX_CONTENT_LENGTH'] = int(os.getenv('MAX_CONTENT_LENGTH', 16 * 1024 * 1024))
app.config['UPLOAD_FOLDER'] = os.getenv('UPLOAD_FOLDER', 'uploads')
app.config['PERMANENT_FOLDER'] = os.getenv('PERMANENT_FOLDER', 'user_data')
app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL', 'sqlite:///serp_comparator.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

ADMIN_EMAIL = os.getenv('ADMIN_EMAIL', 'admin@example.com')
FILE_CLEANUP_DAYS = int(os.getenv('FILE_CLEANUP_DAYS', 7))

# SMTP Configuration from environment
app.config['MAIL_SERVER'] = os.getenv('MAIL_SERVER', 'smtp.yandex.ru')
app.config['MAIL_PORT'] = int(os.getenv('MAIL_PORT', 465))
app.config['MAIL_USERNAME'] = os.getenv('MAIL_USERNAME')
app.config['MAIL_PASSWORD'] = os.getenv('MAIL_PASSWORD')
app.config['MAIL_USE_SSL'] = os.getenv('MAIL_USE_SSL', 'True') == 'True'
app.config['MAIL_DEFAULT_SENDER'] = os.getenv('MAIL_USERNAME')

# Контекстный процессор
@app.context_processor
def inject_common():
    try:
        with open('version.txt', 'r') as f:
            version = f.read().strip()
    except:
        version = '3.18.1'

    is_admin = False
    if 'user_id' in session:
        user = db.session.get(User, session['user_id'])
        if user:
            is_admin = user.is_admin

    return dict(version=version, is_admin=is_admin)

db = SQLAlchemy(app)

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['PERMANENT_FOLDER'], exist_ok=True)

# Декоратор для проверки авторизации
def login_required(f):
    from functools import wraps
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            flash('Пожалуйста, войдите в систему')
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

# Очистка старых файлов при запуске
cleanup_old_files(app.config['UPLOAD_FOLDER'], FILE_CLEANUP_DAYS)

comparator = SERPComparatorWithSentiment()
logger.info("SERP Comparator initialized")

# ---------- МОДЕЛИ ----------
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password_hash = db.Column(db.String(128), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.now)
    is_verified = db.Column(db.Boolean, default=False)
    is_admin = db.Column(db.Boolean, default=False)
    verification_code = db.Column(db.String(6), nullable=True)
    verification_code_expires = db.Column(db.DateTime, nullable=True)
    comparisons = db.relationship('Comparison', backref='user', lazy=True)

    def set_password(self, password):
        self.password_hash = generate_password_hash(password)

    def check_password(self, password):
        return check_password_hash(self.password_hash, password)

    def generate_verification_code(self):
        self.verification_code = ''.join(secrets.choice('0123456789') for _ in range(6))
        self.verification_code_expires = datetime.now() + timedelta(hours=1)
        db.session.commit()
        return self.verification_code

class Project(db.Model):
    """Модель проекта для группировки сравнений"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    name = db.Column(db.String(200), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.now)
    updated_at = db.Column(db.DateTime, default=datetime.now, onupdate=datetime.now)
    
    # === СТАРТОВЫЕ МЕТРИКИ ===
    # Яндекс ТОП-20 (5 тональностей + тип ввода + всего URL)
    y20_client_site = db.Column(db.Float, nullable=True)
    y20_positive = db.Column(db.Float, nullable=True)
    y20_neutral = db.Column(db.Float, nullable=True)
    y20_negative = db.Column(db.Float, nullable=True)
    y20_irrelevant = db.Column(db.Float, nullable=True)
    y20_input_type = db.Column(db.String(20), default='percentage')  # 'percentage' или 'count'
    y20_total_urls = db.Column(db.Integer, default=0)
    
    # Яндекс ТОП-10
    y10_client_site = db.Column(db.Float, nullable=True)
    y10_positive = db.Column(db.Float, nullable=True)
    y10_neutral = db.Column(db.Float, nullable=True)
    y10_negative = db.Column(db.Float, nullable=True)
    y10_irrelevant = db.Column(db.Float, nullable=True)
    y10_input_type = db.Column(db.String(20), default='percentage')
    y10_total_urls = db.Column(db.Integer, default=0)
    
    # Google ТОП-20
    g20_client_site = db.Column(db.Float, nullable=True)
    g20_positive = db.Column(db.Float, nullable=True)
    g20_neutral = db.Column(db.Float, nullable=True)
    g20_negative = db.Column(db.Float, nullable=True)
    g20_irrelevant = db.Column(db.Float, nullable=True)
    g20_input_type = db.Column(db.String(20), default='percentage')
    g20_total_urls = db.Column(db.Integer, default=0)
    
    # Google ТОП-10
    g10_client_site = db.Column(db.Float, nullable=True)
    g10_positive = db.Column(db.Float, nullable=True)
    g10_neutral = db.Column(db.Float, nullable=True)
    g10_negative = db.Column(db.Float, nullable=True)
    g10_irrelevant = db.Column(db.Float, nullable=True)
    g10_input_type = db.Column(db.String(20), default='percentage')
    g10_total_urls = db.Column(db.Integer, default=0)
    
    # Связи
    comparisons = db.relationship('Comparison', backref='project', lazy=True)
    
    def __repr__(self):
        return f'<Project {self.name}>'
    
    def has_baseline_metrics(self, search_engine=None, top_n=None):
        """Проверяет, есть ли у проекта стартовые метрики"""
        if search_engine == 'yandex' and top_n == 20:
            return any([self.y20_client_site, self.y20_positive, self.y20_neutral, self.y20_negative, self.y20_irrelevant])
        elif search_engine == 'yandex' and top_n == 10:
            return any([self.y10_client_site, self.y10_positive, self.y10_neutral, self.y10_negative, self.y10_irrelevant])
        elif search_engine == 'google' and top_n == 20:
            return any([self.g20_client_site, self.g20_positive, self.g20_neutral, self.g20_negative, self.g20_irrelevant])
        elif search_engine == 'google' and top_n == 10:
            return any([self.g10_client_site, self.g10_positive, self.g10_neutral, self.g10_negative, self.g10_irrelevant])
        else:
            # Проверяем любые метрики
            return any([
                self.y20_client_site, self.y20_positive, self.y20_neutral, self.y20_negative, self.y20_irrelevant,
                self.y10_client_site, self.y10_positive, self.y10_neutral, self.y10_negative, self.y10_irrelevant,
                self.g20_client_site, self.g20_positive, self.g20_neutral, self.g20_negative, self.g20_irrelevant,
                self.g10_client_site, self.g10_positive, self.g10_neutral, self.g10_negative, self.g10_irrelevant
            ])
    
    def get_baseline(self, search_engine, top_n):
        """Возвращает стартовые метрики для указанного поисковика и ТОП-N"""
        if search_engine == 'yandex' and top_n == 20:
            return {
                'client_site': self.y20_client_site or 0,
                'positive': self.y20_positive or 0,
                'neutral': self.y20_neutral or 0,
                'negative': self.y20_negative or 0,
                'irrelevant': self.y20_irrelevant or 0,
                'input_type': self.y20_input_type or 'percentage',
                'total_urls': self.y20_total_urls or 0
            }
        elif search_engine == 'yandex' and top_n == 10:
            return {
                'client_site': self.y10_client_site or 0,
                'positive': self.y10_positive or 0,
                'neutral': self.y10_neutral or 0,
                'negative': self.y10_negative or 0,
                'irrelevant': self.y10_irrelevant or 0,
                'input_type': self.y10_input_type or 'percentage',
                'total_urls': self.y10_total_urls or 0
            }
        elif search_engine == 'google' and top_n == 20:
            return {
                'client_site': self.g20_client_site or 0,
                'positive': self.g20_positive or 0,
                'neutral': self.g20_neutral or 0,
                'negative': self.g20_negative or 0,
                'irrelevant': self.g20_irrelevant or 0,
                'input_type': self.g20_input_type or 'percentage',
                'total_urls': self.g20_total_urls or 0
            }
        elif search_engine == 'google' and top_n == 10:
            return {
                'client_site': self.g10_client_site or 0,
                'positive': self.g10_positive or 0,
                'neutral': self.g10_neutral or 0,
                'negative': self.g10_negative or 0,
                'irrelevant': self.g10_irrelevant or 0,
                'input_type': self.g10_input_type or 'percentage',
                'total_urls': self.g10_total_urls or 0
            }
        return None
    
    def set_baseline(self, search_engine, top_n, metrics):
        """Устанавливает стартовые метрики"""
        if search_engine == 'yandex' and top_n == 20:
            self.y20_client_site = metrics.get('client_site')
            self.y20_positive = metrics.get('positive')
            self.y20_neutral = metrics.get('neutral')
            self.y20_negative = metrics.get('negative')
            self.y20_irrelevant = metrics.get('irrelevant')
            self.y20_input_type = metrics.get('input_type', 'percentage')
            self.y20_total_urls = metrics.get('total_urls', 0)
        elif search_engine == 'yandex' and top_n == 10:
            self.y10_client_site = metrics.get('client_site')
            self.y10_positive = metrics.get('positive')
            self.y10_neutral = metrics.get('neutral')
            self.y10_negative = metrics.get('negative')
            self.y10_irrelevant = metrics.get('irrelevant')
            self.y10_input_type = metrics.get('input_type', 'percentage')
            self.y10_total_urls = metrics.get('total_urls', 0)
        elif search_engine == 'google' and top_n == 20:
            self.g20_client_site = metrics.get('client_site')
            self.g20_positive = metrics.get('positive')
            self.g20_neutral = metrics.get('neutral')
            self.g20_negative = metrics.get('negative')
            self.g20_irrelevant = metrics.get('irrelevant')
            self.g20_input_type = metrics.get('input_type', 'percentage')
            self.g20_total_urls = metrics.get('total_urls', 0)
        elif search_engine == 'google' and top_n == 10:
            self.g10_client_site = metrics.get('client_site')
            self.g10_positive = metrics.get('positive')
            self.g10_neutral = metrics.get('neutral')
            self.g10_negative = metrics.get('negative')
            self.g10_irrelevant = metrics.get('irrelevant')
            self.g10_input_type = metrics.get('input_type', 'percentage')
            self.g10_total_urls = metrics.get('total_urls', 0)


class Comparison(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    project_id = db.Column(db.Integer, db.ForeignKey('project.id'), nullable=True)
    session_id = db.Column(db.String(36), nullable=True)
    date = db.Column(db.DateTime, default=datetime.now)
    file1_name = db.Column(db.String(200))
    file2_name = db.Column(db.String(200))
    excel_path = db.Column(db.String(500))
    yandex_excel_path = db.Column(db.String(500))
    google_excel_path = db.Column(db.String(500))
    summary_y20_path = db.Column(db.String(500))
    summary_y10_path = db.Column(db.String(500))
    summary_g20_path = db.Column(db.String(500))
    summary_g10_path = db.Column(db.String(500))
    analysis_y20_path = db.Column(db.String(500))
    analysis_g20_path = db.Column(db.String(500))
    analysis_total20_path = db.Column(db.String(500))
    analysis_y10_path = db.Column(db.String(500))
    analysis_g10_path = db.Column(db.String(500))
    analysis_total10_path = db.Column(db.String(500))
    url_y_excel_path = db.Column(db.String(500))
    url_g_excel_path = db.Column(db.String(500))
    url_combined_excel_path = db.Column(db.String(500))
    chart_y20_path = db.Column(db.String(500))
    chart_g20_path = db.Column(db.String(500))
    chart_total20_path = db.Column(db.String(500))
    chart_y10_path = db.Column(db.String(500))
    chart_g10_path = db.Column(db.String(500))
    chart_total10_path = db.Column(db.String(500))
    # Сетки ТОП-10 (file1 - старая дата, file2 - новая дата)
    grid_y10_path = db.Column(db.String(500))
    grid_y10_file2_path = db.Column(db.String(500))
    grid_g10_path = db.Column(db.String(500))
    grid_g10_file2_path = db.Column(db.String(500))
    # PowerPoint презентации сеток ТОП-10
    grid_y10_pptx_path = db.Column(db.String(500))
    grid_g10_pptx_path = db.Column(db.String(500))
    # Диаграммы со стартовыми показателями (ТОП-10 и ТОП-20)
    chart_start_y20_path = db.Column(db.String(500))
    chart_start_g20_path = db.Column(db.String(500))
    chart_start_y10_path = db.Column(db.String(500))
    chart_start_g10_path = db.Column(db.String(500))
    # PowerPoint презентации диаграмм
    chart_y20_pptx_path = db.Column(db.String(500))
    chart_g20_pptx_path = db.Column(db.String(500))
    chart_total20_pptx_path = db.Column(db.String(500))
    chart_y10_pptx_path = db.Column(db.String(500))
    chart_g10_pptx_path = db.Column(db.String(500))
    chart_total10_pptx_path = db.Column(db.String(500))


class BaselineMetrics(db.Model):
    """Модель для хранения стартовых показателей тональности"""
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    session_id = db.Column(db.String(36), nullable=True)
    comparison_id = db.Column(db.Integer, db.ForeignKey('comparison.id'), nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.now)
    
    # Тип: yandex_top20, google_top20, yandex_top10, google_top10
    metric_type = db.Column(db.String(50), nullable=False)
    
    # Тип ввода: percentage или count
    input_type = db.Column(db.String(20), default='percentage')
    
    # Данные тональности (можно хранить проценты или количество)
    client_site_value = db.Column(db.Float, default=0)
    client_site_is_count = db.Column(db.Boolean, default=False)  # True если это количество, False если процент
    
    positive_value = db.Column(db.Float, default=0)
    positive_is_count = db.Column(db.Boolean, default=False)
    
    neutral_value = db.Column(db.Float, default=0)
    neutral_is_count = db.Column(db.Boolean, default=False)
    
    negative_value = db.Column(db.Float, default=0)
    negative_is_count = db.Column(db.Boolean, default=False)
    
    irrelevant_value = db.Column(db.Float, default=0)
    irrelevant_is_count = db.Column(db.Boolean, default=False)
    
    # Общее количество URL (для пересчета процентов в количество и наоборот)
    total_urls = db.Column(db.Integer, default=0)


class Feedback(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    user_email = db.Column(db.String(120), nullable=False)
    message = db.Column(db.Text, nullable=False)
    status = db.Column(db.String(20), default='new')
    created_at = db.Column(db.DateTime, default=datetime.now)

# ---------- ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ----------
def clear_user_session():
    session_id = session.get('session_id')
    if session_id:
        user_folder = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        if os.path.exists(user_folder):
            shutil.rmtree(user_folder)
    keys_to_remove = ['session_id', 'file1', 'file1_name', 'file2', 'file2_name',
                      'report_excel', 'report_txt_y20', 'report_txt_y10', 'report_txt_g20', 'report_txt_g10',
                      'report_analysis_y20', 'report_analysis_g20', 'report_analysis_total20',
                      'report_analysis_y10', 'report_analysis_g10', 'report_analysis_total10',
                      'report_url_y_excel', 'report_url_g_excel', 'report_url_combined_excel',
                      'yandex_full_excel_path', 'google_full_excel_path',
                      'chart_yandex20_path', 'chart_google20_path', 'chart_total20_path',
                      'chart_yandex10_path', 'chart_google10_path', 'chart_total10_path',
                      'summary_y20_excel', 'summary_y10_excel', 'summary_g20_excel', 'summary_g10_excel']
    for key in keys_to_remove:
        session.pop(key, None)

def send_verification_email(email, code):
    """Отправка кода подтверждения на email (с выводом в консоль для отладки)"""
    # Для отладки выводим код в консоль
    print(f"[DEBUG] Код подтверждения для {email}: {code}")

    # Реальная отправка через SMTP
    try:
        msg = MIMEMultipart()
        msg['From'] = app.config['MAIL_DEFAULT_SENDER']
        msg['To'] = email
        msg['Subject'] = 'Подтверждение регистрации в SERP Comparator'

        body = f"""
        Здравствуйте!

        Вы зарегистрировались в сервисе SERP Comparator.
        Для завершения регистрации введите код подтверждения:

        {code}

        Код действителен в течение 1 часа.

        Если вы не регистрировались, проигнорируйте это письмо.

        С уважением,
        Команда SERP Comparator
        """
        msg.attach(MIMEText(body, 'plain'))

        context = ssl.create_default_context()
        with smtplib.SMTP_SSL(app.config['MAIL_SERVER'], app.config['MAIL_PORT'], context=context) as server:
            server.login(app.config['MAIL_USERNAME'], app.config['MAIL_PASSWORD'])
            server.send_message(msg)
        return True
    except Exception as e:
        logger.error(f"Ошибка отправки email: {e}")
        # Всё равно возвращаем True, чтобы регистрация прошла (код выведен в консоль)
        return True

# ---------- МАРШРУТЫ ----------
@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/user_data/<path:filename>')
def user_data_file(filename):
    return send_from_directory(app.config['PERMANENT_FOLDER'], filename)

@app.route('/')
def landing():
    """Новая главная страница (landing) с демо-сравнением"""
    if 'user_id' in session:
        # Если пользователь авторизован, перенаправляем на проекты
        return redirect(url_for('projects'))
    return render_template('landing.html')


@app.route('/demo-compare', methods=['POST'])
def demo_compare():
    """Демо-сравнение без регистрации"""
    session_id = session.get('session_id')
    if not session_id:
        session_id = str(uuid.uuid4())
        session['session_id'] = session_id
    
    session_path = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
    os.makedirs(session_path, exist_ok=True)
    
    file1 = request.files.get('file1')
    file2 = request.files.get('file2')
    
    if not file1 or not file2:
        return jsonify({'success': False, 'error': 'Необходимо загрузить оба файла'})
    
    # Сохраняем файлы во временную сессию
    file1_path = os.path.join(session_path, secure_filename(file1.filename))
    file2_path = os.path.join(session_path, secure_filename(file2.filename))
    file1.save(file1_path)
    file2.save(file2_path)
    
    session['file1'] = file1_path
    session['file2'] = file2_path
    session['file1_name'] = file1.filename
    session['file2_name'] = file2.filename
    
    return jsonify({'success': True, 'redirect': url_for('compare')})


@app.route('/upload', methods=['GET', 'POST'])
@app.route('/upload/<int:project_id>', methods=['GET', 'POST'])
@login_required
def index(project_id=None):
    """Страница загрузки файлов внутри проекта"""
    # Получаем проект
    project = None
    if project_id:
        project = Project.query.get_or_404(project_id)
        if project.user_id != session.get('user_id'):
            flash('Нет доступа к этому проекту')
            return redirect(url_for('projects'))
        session['current_project_id'] = project_id
    elif session.get('current_project_id'):
        project = Project.query.get(session['current_project_id'])
    else:
        flash('Выберите проект для работы')
        return redirect(url_for('projects'))
    
    if request.method == 'POST':
        session_id = session.get('session_id')
        if not session_id:
            session_id = str(uuid.uuid4())
            session['session_id'] = session_id
        session_path = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        os.makedirs(session_path, exist_ok=True)

        file1 = request.files.get('file1')
        file2 = request.files.get('file2')

        if not file1 or not file2:
            flash('Загрузите оба файла')
            return redirect(request.url)

        if not (allowed_file(file1.filename, ALLOWED_EXTENSIONS) and allowed_file(file2.filename, ALLOWED_EXTENSIONS)):
            flash('Оба файла должны быть в формате .xlsx')
            return redirect(request.url)

        filename1 = secure_filename(file1.filename)
        filename2 = secure_filename(file2.filename)
        filepath1 = os.path.join(session_path, filename1)
        filepath2 = os.path.join(session_path, filename2)
        file1.save(filepath1)
        file2.save(filepath2)

        session['file1'] = filepath1
        session['file1_name'] = filename1
        session['file2'] = filepath2
        session['file2_name'] = filename2

        with open(filepath1, 'rb') as f:
            content1 = f.read()
        with open(filepath2, 'rb') as f:
            content2 = f.read()

        df1_20, sentiment1_20, error1 = comparator.process_excel_file_with_sentiment(content1, filename1, sheet_index=0, top_n=20)
        df2_20, sentiment2_20, error2 = comparator.process_excel_file_with_sentiment(content2, filename2, sheet_index=0, top_n=20)
        if error1 or error2:
            flash(f'Ошибка обработки: {error1 or error2}')
            clear_user_session()
            return redirect(url_for('index'))

        stats1 = comparator.calculate_sentiment_statistics(sentiment1_20)
        stats2 = comparator.calculate_sentiment_statistics(sentiment2_20)
        report_text1 = comparator.generate_sentiment_report(stats1, filename1, len(sentiment1_20))
        report_text2 = comparator.generate_sentiment_report(stats2, filename2, len(sentiment2_20))

        return render_template('index.html',
                               file1_loaded=True,
                               file2_loaded=True,
                               file1_name=filename1,
                               file2_name=filename2,
                               report_text1=report_text1,
                               report_text2=report_text2,
                               project=project)

    else:
        file1_loaded = 'file1' in session
        file2_loaded = 'file2' in session
        context = {
            'file1_loaded': file1_loaded,
            'file2_loaded': file2_loaded,
            'file1_name': session.get('file1_name', ''),
            'file2_name': session.get('file2_name', '') if file2_loaded else '',
            'project': project
        }
        return render_template('index.html', **context)

@app.route('/clear_files')
def clear_files():
    """Очистка загруженных файлов из сессии"""
    # Удаляем файлы из сессии
    session.pop('file1', None)
    session.pop('file2', None)
    session.pop('file1_name', None)
    session.pop('file2_name', None)
    flash('Файлы успешно очищены', 'info')
    return redirect(url_for('index'))

@app.route('/register', methods=['GET', 'POST'])
def register():
    # Сбрасываем сессию при открытии страницы регистрации
    if request.method == 'GET':
        session.pop('user_id', None)
        session.pop('user_email', None)
        session.pop('pending_verification_email', None)
    
    if request.method == 'POST':
        email = request.form['email'].strip()
        password = request.form['password']
        if not email or not password:
            flash('Заполните все поля')
            return redirect(url_for('register'))
        if User.query.filter_by(email=email).first():
            flash('Пользователь с таким email уже существует')
            return redirect(url_for('register'))

        user = User(email=email)
        user.set_password(password)
        # Назначаем администратора если email совпадает с ADMIN_EMAIL
        if email == ADMIN_EMAIL:
            user.is_admin = True
        db.session.add(user)
        db.session.commit()
        code = user.generate_verification_code()
        send_verification_email(email, code)

        # Сохраняем email в сессии для последующего подтверждения
        session['pending_verification_email'] = email
        flash('На вашу почту отправлен код подтверждения. Введите его для завершения регистрации.', 'info')
        return redirect(url_for('verify'))

    return render_template('register.html')

@app.route('/verify', methods=['GET', 'POST'])
def verify():
    if 'pending_verification_email' not in session:
        flash('Пожалуйста, зарегистрируйтесь сначала', 'error')
        return redirect(url_for('register'))

    email = session['pending_verification_email']
    user = User.query.filter_by(email=email).first()
    if not user:
        session.pop('pending_verification_email', None)
        flash('Пользователь не найден', 'error')
        return redirect(url_for('register'))

    if user.is_verified:
        session.pop('pending_verification_email', None)
        flash('Аккаунт уже подтверждён. Войдите.', 'info')
        return redirect(url_for('login'))

    if request.method == 'POST':
        code = request.form.get('code', '').strip()
        if code == user.verification_code and user.verification_code_expires > datetime.now():
            user.is_verified = True
            user.verification_code = None
            user.verification_code_expires = None
            db.session.commit()
            session.pop('pending_verification_email', None)
            flash('Аккаунт успешно подтверждён! Теперь вы можете войти.', 'success')
            return redirect(url_for('login'))
        else:
            flash('Неверный или просроченный код', 'error')

    return render_template('verify.html', email=email)

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form['email'].strip()
        password = request.form['password']
        user = User.query.filter_by(email=email).first()
        if user and user.check_password(password):
            if not user.is_verified:
                flash('Аккаунт не подтверждён. Пожалуйста, проверьте почту и введите код.', 'error')
                session['pending_verification_email'] = email
                return redirect(url_for('verify'))
            session['user_id'] = user.id
            session['user_email'] = user.email
            flash('Вы успешно вошли')
            return redirect(url_for('index'))
        else:
            flash('Неверный email или пароль')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.pop('user_id', None)
    session.pop('user_email', None)
    flash('Вы вышли')
    return redirect(url_for('index'))

@app.route('/profile')
@app.route('/profile/project/<int:project_id>')
@login_required
def profile(project_id=None):
    """История сравнений - общая или для конкретного проекта"""
    user = db.session.get(User, session['user_id'])
    project = None
    if project_id:
        project = Project.query.get_or_404(project_id)
        if project.user_id != session.get('user_id'):
            flash('Нет доступа к этому проекту')
            return redirect(url_for('projects'))
        comparisons = Comparison.query.filter_by(
            user_id=session['user_id'], 
            project_id=project_id
        ).order_by(Comparison.date.desc()).all()
    else:
        # Показываем все сравнения пользователя
        comparisons = Comparison.query.filter_by(
            user_id=session['user_id']
        ).order_by(Comparison.date.desc()).all()
    
    def resolve_path(path):
        """Преобразует относительный путь из БД в абсолютный для проверки"""
        if not path:
            return None
        if os.path.isabs(path):
            return path
        # Пути в БД уже содержат 'user_data/' prefix, убираем его если PERMANENT_FOLDER тоже 'user_data'
        if path.startswith('user_data/'):
            path = path[10:]  # убираем 'user_data/'
        return os.path.join(app.config['PERMANENT_FOLDER'], path)
    
    for comp in comparisons:
        comp.files_exist = all([
            os.path.exists(resolve_path(comp.excel_path)) if comp.excel_path else False,
            os.path.exists(resolve_path(comp.yandex_excel_path)) if comp.yandex_excel_path else False,
            os.path.exists(resolve_path(comp.google_excel_path)) if comp.google_excel_path else False,
            os.path.exists(resolve_path(comp.analysis_y20_path)) if comp.analysis_y20_path else False,
            os.path.exists(resolve_path(comp.analysis_g20_path)) if comp.analysis_g20_path else False,
            os.path.exists(resolve_path(comp.analysis_total20_path)) if comp.analysis_total20_path else False,
            os.path.exists(resolve_path(comp.analysis_y10_path)) if comp.analysis_y10_path else False,
            os.path.exists(resolve_path(comp.analysis_g10_path)) if comp.analysis_g10_path else False,
            os.path.exists(resolve_path(comp.analysis_total10_path)) if comp.analysis_total10_path else False,
            os.path.exists(resolve_path(comp.url_y_excel_path)) if comp.url_y_excel_path else False,
            os.path.exists(resolve_path(comp.url_g_excel_path)) if comp.url_g_excel_path else False,
            os.path.exists(resolve_path(comp.url_combined_excel_path)) if comp.url_combined_excel_path else False,
            os.path.exists(resolve_path(comp.chart_y20_path)) if comp.chart_y20_path else False,
            os.path.exists(resolve_path(comp.chart_g20_path)) if comp.chart_g20_path else False,
            os.path.exists(resolve_path(comp.chart_total20_path)) if comp.chart_total20_path else False,
            os.path.exists(resolve_path(comp.chart_y10_path)) if comp.chart_y10_path else False,
            os.path.exists(resolve_path(comp.chart_g10_path)) if comp.chart_g10_path else False,
            os.path.exists(resolve_path(comp.chart_total10_path)) if comp.chart_total10_path else False,
            os.path.exists(resolve_path(comp.grid_y10_path)) if comp.grid_y10_path else False,
            os.path.exists(resolve_path(comp.grid_g10_path)) if comp.grid_g10_path else False,
            os.path.exists(resolve_path(comp.summary_y20_path)) if comp.summary_y20_path else False,
            os.path.exists(resolve_path(comp.summary_y10_path)) if comp.summary_y10_path else False,
            os.path.exists(resolve_path(comp.summary_g20_path)) if comp.summary_g20_path else False,
            os.path.exists(resolve_path(comp.summary_g10_path)) if comp.summary_g10_path else False,
        ])
    return render_template('profile.html', user=user, comparisons=comparisons, project=project)

@app.route('/view_comparison/<int:comp_id>')
def view_comparison(comp_id):
    if 'user_id' not in session:
        flash('Войдите, чтобы просмотреть сравнение')
        return redirect(url_for('login'))
    comp = Comparison.query.get_or_404(comp_id)
    if comp.user_id != session['user_id']:
        flash('Нет доступа к этому сравнению')
        return redirect(url_for('profile'))

    analysis_y20 = ''
    analysis_g20 = ''
    analysis_total20 = ''
    analysis_y10 = ''
    analysis_g10 = ''
    analysis_total10 = ''
    if comp.analysis_y20_path and os.path.exists(comp.analysis_y20_path):
        with open(comp.analysis_y20_path, 'r', encoding='utf-8') as f:
            analysis_y20 = f.read()
    if comp.analysis_g20_path and os.path.exists(comp.analysis_g20_path):
        with open(comp.analysis_g20_path, 'r', encoding='utf-8') as f:
            analysis_g20 = f.read()
    if comp.analysis_total20_path and os.path.exists(comp.analysis_total20_path):
        with open(comp.analysis_total20_path, 'r', encoding='utf-8') as f:
            analysis_total20 = f.read()
    if comp.analysis_y10_path and os.path.exists(comp.analysis_y10_path):
        with open(comp.analysis_y10_path, 'r', encoding='utf-8') as f:
            analysis_y10 = f.read()
    if comp.analysis_g10_path and os.path.exists(comp.analysis_g10_path):
        with open(comp.analysis_g10_path, 'r', encoding='utf-8') as f:
            analysis_g10 = f.read()
    if comp.analysis_total10_path and os.path.exists(comp.analysis_total10_path):
        with open(comp.analysis_total10_path, 'r', encoding='utf-8') as f:
            analysis_total10 = f.read()

    # Формируем URL для сеток ТОП-10
    grid_y10_file1_url = None
    grid_y10_file2_url = None
    grid_g10_file1_url = None
    grid_g10_file2_url = None
    
    if comp.grid_y10_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.grid_y10_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            grid_y10_file1_url = url_for('user_data_file', filename=comp.grid_y10_path.replace('user_data/', ''))
    
    if comp.grid_y10_file2_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.grid_y10_file2_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            grid_y10_file2_url = url_for('user_data_file', filename=comp.grid_y10_file2_path.replace('user_data/', ''))
    
    if comp.grid_g10_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.grid_g10_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            grid_g10_file1_url = url_for('user_data_file', filename=comp.grid_g10_path.replace('user_data/', ''))
    
    if comp.grid_g10_file2_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.grid_g10_file2_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            grid_g10_file2_url = url_for('user_data_file', filename=comp.grid_g10_file2_path.replace('user_data/', ''))

    # Получаем baseline метрики для этого сравнения
    baseline_metrics = {}
    metrics = BaselineMetrics.query.filter_by(comparison_id=comp.id).all()
    for metric in metrics:
        baseline_metrics[metric.metric_type] = {
            'input_type': metric.input_type,
            'total_urls': metric.total_urls,
            'client_site_value': metric.client_site_value,
            'positive_value': metric.positive_value,
            'neutral_value': metric.neutral_value,
            'negative_value': metric.negative_value,
            'irrelevant_value': metric.irrelevant_value
        }

    # Формируем URL для PPTX диаграмм
    chart_y20_pptx_url = None
    chart_g20_pptx_url = None
    chart_total20_pptx_url = None
    chart_y10_pptx_url = None
    chart_g10_pptx_url = None
    chart_total10_pptx_url = None

    if comp.chart_y20_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_y20_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_y20_pptx_url = url_for('user_data_file', filename=comp.chart_y20_pptx_path.replace('user_data/', ''))

    if comp.chart_g20_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_g20_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_g20_pptx_url = url_for('user_data_file', filename=comp.chart_g20_pptx_path.replace('user_data/', ''))

    if comp.chart_total20_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_total20_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_total20_pptx_url = url_for('user_data_file', filename=comp.chart_total20_pptx_path.replace('user_data/', ''))

    if comp.chart_y10_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_y10_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_y10_pptx_url = url_for('user_data_file', filename=comp.chart_y10_pptx_path.replace('user_data/', ''))

    if comp.chart_g10_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_g10_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_g10_pptx_url = url_for('user_data_file', filename=comp.chart_g10_pptx_path.replace('user_data/', ''))

    if comp.chart_total10_pptx_path:
        full_path = os.path.join(app.config['PERMANENT_FOLDER'], comp.chart_total10_pptx_path.replace('user_data/', ''))
        if os.path.exists(full_path):
            chart_total10_pptx_url = url_for('user_data_file', filename=comp.chart_total10_pptx_path.replace('user_data/', ''))

    return render_template('view_comparison.html',
                           comp=comp,
                           project=comp.project,
                           analysis_y20=analysis_y20,
                           analysis_g20=analysis_g20,
                           analysis_total20=analysis_total20,
                           analysis_y10=analysis_y10,
                           analysis_g10=analysis_g10,
                           analysis_total10=analysis_total10,
                           # Сетки ТОП-10
                           grid_y10_file1=grid_y10_file1_url,
                           grid_y10_file2=grid_y10_file2_url,
                           grid_g10_file1=grid_g10_file1_url,
                           grid_g10_file2=grid_g10_file2_url,
                           # PPTX диаграммы
                           chart_y20_pptx=chart_y20_pptx_url,
                           chart_g20_pptx=chart_g20_pptx_url,
                           chart_total20_pptx=chart_total20_pptx_url,
                           chart_y10_pptx=chart_y10_pptx_url,
                           chart_g10_pptx=chart_g10_pptx_url,
                           chart_total10_pptx=chart_total10_pptx_url,
                           # Baseline метрики
                           baseline_metrics=baseline_metrics)

@app.route('/info')
def info():
    return render_template('info.html')


# ---------- РОУТЫ ПРОЕКТОВ ----------
@app.route('/projects')
@login_required
def projects():
    """Страница со списком проектов пользователя (ЛК)"""
    user_projects = Project.query.filter_by(user_id=session['user_id']).order_by(Project.updated_at.desc()).all()
    return render_template('projects.html', projects=user_projects)


@app.route('/project/create', methods=['POST'])
@login_required
def create_project():
    """Создание нового проекта"""
    name = request.form.get('name', '').strip()
    if not name:
        flash('Введите название проекта', 'danger')
        return redirect(url_for('projects'))
    
    project = Project(
        user_id=session['user_id'],
        name=name
    )
    db.session.add(project)
    db.session.commit()
    
    flash(f'Проект "{name}" создан', 'success')
    return redirect(url_for('projects'))


@app.route('/project/<int:project_id>')
@login_required
def project_detail(project_id):
    """Страница проекта — перенаправляет на загрузку файлов"""
    project = Project.query.get_or_404(project_id)
    if project.user_id != session['user_id']:
        flash('Нет доступа к этому проекту', 'danger')
        return redirect(url_for('projects'))
    
    # Сохраняем текущий проект в сессии
    session['current_project_id'] = project_id
    
    # Очищаем загруженные файлы из предыдущего проекта
    # (только файлы для загрузки, не историю сравнений)
    if 'file1' in session:
        del session['file1']
    if 'file2' in session:
        del session['file2']
    if 'file1_name' in session:
        del session['file1_name']
    if 'file2_name' in session:
        del session['file2_name']
    
    # Перенаправляем на страницу загрузки
    return redirect(url_for('index', project_id=project_id))


@app.route('/project/<int:project_id>/rename', methods=['POST'])
@login_required
def rename_project(project_id):
    """Переименование проекта"""
    project = Project.query.get_or_404(project_id)
    if project.user_id != session['user_id']:
        flash('Нет доступа к этому проекту', 'danger')
        return redirect(url_for('projects'))
    
    name = request.form.get('name', '').strip()
    if not name:
        flash('Введите название проекта', 'danger')
        return redirect(url_for('projects'))
    
    project.name = name
    project.updated_at = datetime.now()
    db.session.commit()
    
    flash('Проект переименован', 'success')
    return redirect(url_for('projects'))


@app.route('/project/<int:project_id>/delete', methods=['POST'])
@login_required
def delete_project(project_id):
    """Удаление проекта и всех связанных сравнений"""
    project = Project.query.get_or_404(project_id)
    if project.user_id != session['user_id']:
        flash('Нет доступа к этому проекту', 'danger')
        return redirect(url_for('projects'))
    
    # Удаляем все сравнения проекта
    for comparison in project.comparisons:
        # Удаляем связанные baseline метрики
        BaselineMetrics.query.filter_by(comparison_id=comparison.id).delete()
        
        # Удаляем файлы сравнения
        for attr in ['excel_path', 'yandex_excel_path', 'google_excel_path',
                     'summary_y20_path', 'summary_y10_path', 'summary_g20_path',
                     'yandex_pptx_path', 'google_pptx_path',
                     'yandex_y10_grid_png', 'google_y10_grid_png']:
            path = getattr(comparison, attr, None)
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                except:
                    pass
        
        db.session.delete(comparison)
    
    db.session.delete(project)
    db.session.commit()
    
    # Очищаем текущий проект из сессии если это был он
    if session.get('current_project_id') == project_id:
        session.pop('current_project_id', None)
    
    flash('Проект удален', 'success')
    return redirect(url_for('projects'))


@app.route('/project/<int:project_id>/baseline', methods=['POST'])
@login_required
def save_baseline_metrics(project_id):
    """Сохранение стартовых метрик проекта - 5 тональностей для ТОП-10 и ТОП-20"""
    project = Project.query.get_or_404(project_id)
    if project.user_id != session['user_id']:
        flash('Нет доступа к этому проекту', 'danger')
        return redirect(url_for('projects'))
    
    try:
        # Сохраняем метрики для 4 групп: Яндекс ТОП-20, Яндекс ТОП-10, Google ТОП-20, Google ТОП-10
        groups = [
            ('yandex', 20, 'y20'),
            ('yandex', 10, 'y10'),
            ('google', 20, 'g20'),
            ('google', 10, 'g10')
        ]
        
        for search_engine, top_n, prefix in groups:
            # Получаем тип ввода и общее количество URL
            input_type = request.form.get(f'{prefix}_input_type', 'percentage')
            total_urls = int(request.form.get(f'{prefix}_total_urls', 0) or 0)
            
            # Получаем значения для 5 тональностей
            metrics = {
                'client_site': float(request.form.get(f'{prefix}_client_site', 0) or 0),
                'positive': float(request.form.get(f'{prefix}_positive', 0) or 0),
                'neutral': float(request.form.get(f'{prefix}_neutral', 0) or 0),
                'negative': float(request.form.get(f'{prefix}_negative', 0) or 0),
                'irrelevant': float(request.form.get(f'{prefix}_irrelevant', 0) or 0),
                'input_type': input_type,
                'total_urls': total_urls
            }
            
            # Сохраняем через метод модели
            project.set_baseline(search_engine, top_n, metrics)
        
        db.session.commit()
        flash('Стартовые метрики сохранены', 'success')
        
    except Exception as e:
        flash(f'Ошибка при сохранении метрик: {str(e)}', 'danger')
    
    return redirect(url_for('projects'))


@app.route('/project/<int:project_id>/baseline', methods=['GET'])
@login_required
def get_baseline_metrics(project_id):
    """Получение стартовых метрик проекта в формате JSON"""
    project = Project.query.get_or_404(project_id)
    if project.user_id != session['user_id']:
        return jsonify({'error': 'Нет доступа к этому проекту'}), 403
    
    groups = [
        ('yandex', 20, 'y20'),
        ('yandex', 10, 'y10'),
        ('google', 20, 'g20'),
        ('google', 10, 'g10')
    ]
    
    result = {}
    for search_engine, top_n, prefix in groups:
        baseline = project.get_baseline(search_engine, top_n)
        if baseline:
            result[prefix] = {
                'client_site': baseline.get('client_site') or 0,
                'positive': baseline.get('positive') or 0,
                'neutral': baseline.get('neutral') or 0,
                'negative': baseline.get('negative') or 0,
                'irrelevant': baseline.get('irrelevant') or 0,
                'input_type': baseline.get('input_type', 'percentage'),
                'total_urls': baseline.get('total_urls', 0)
            }
    
    return jsonify(result)


@app.route('/download_example')
def download_example():
    from openpyxl import Workbook
    import tempfile

    wb = Workbook()
    ws = wb.active
    ws.title = "Пример"

    headers = ["№", "Поисковый запрос", "ТОП1", "ТОП2", "ТОП3", "ТОП4", "ТОП5", "ТОП6", "ТОП7", "ТОП8", "ТОП9", "ТОП10"]
    for col, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True)
        cell.alignment = Alignment(horizontal='center')
    
    ws.cell(row=2, column=2, value="Данные начинаются с 3-й строки").font = Font(italic=True)

    example_data = [
        [1, "купить квартиру в москве", "https://cian.ru/...", "https://avito.ru/...", "https://domclick.ru/...", "", "", "", "", "", "", ""],
        [2, "ремонт квартир цена", "https://profi.ru/...", "https://youdo.com/...", "", "", "", "", "", "", "", ""],
        [3, "строительство домов", "https://example.com/...", "https://another.com/...", "", "", "", "", "", "", "", ""],
    ]
    for row_idx, row_data in enumerate(example_data, start=3):
        for col_idx, value in enumerate(row_data, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    fills = [
        (3, 3, 'client_site'),
        (3, 4, 'positive'),
        (4, 3, 'neutral'),
        (4, 4, 'negative'),
        (5, 3, 'irrelevant'),
    ]
    for row, col, sent in fills:
        cell = ws.cell(row=row, column=col)
        fill = PatternFill(start_color=SENTIMENT_COLORS_EN[sent], end_color=SENTIMENT_COLORS_EN[sent], fill_type="solid")
        cell.fill = fill

    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        wb.save(tmp.name)
        tmp_path = tmp.name

    return send_file(tmp_path, as_attachment=True, download_name='example_serp_file.xlsx')

@app.route('/compare')
def compare():
    if 'file1' not in session or 'file2' not in session:
        flash('Не хватает файлов для сравнения')
        return redirect(url_for('index'))
    
    # Получаем текущий проект из сессии
    project = None
    project_id = session.get('current_project_id')
    if project_id and 'user_id' in session:
        project = Project.query.get(project_id)
        if project and project.user_id != session['user_id']:
            project = None

    session_id = session.get('session_id')
    if not session_id:
        flash('Сессия не найдена')
        return redirect(url_for('index'))
    session_path = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
    if not os.path.exists(session_path):
        flash('Данные сессии устарели, загрузите файлы заново')
        return redirect(url_for('index'))

    file1_path = session['file1']
    file2_path = session['file2']
    file1_name = session['file1_name']
    file2_name = session['file2_name']
    label1 = os.path.splitext(file1_name)[0]
    label2 = os.path.splitext(file2_name)[0]

    with open(file1_path, 'rb') as f1, open(file2_path, 'rb') as f2:
        content1 = f1.read()
        content2 = f2.read()

    def process_sheet(content, filename, sheet_index, top_n):
        try:
            return comparator.process_excel_file_with_sentiment(content, filename, sheet_index=sheet_index, top_n=top_n)
        except Exception as e:
            return None, None, str(e)

    df1_y_20, sent1_y_20, err1_y = process_sheet(content1, file1_name, 0, 20)
    df1_y_10, sent1_y_10, _ = process_sheet(content1, file1_name, 0, 10)
    df2_y_20, sent2_y_20, err2_y = process_sheet(content2, file2_name, 0, 20)
    df2_y_10, sent2_y_10, _ = process_sheet(content2, file2_name, 0, 10)

    df1_g_20, sent1_g_20, err1_g = process_sheet(content1, file1_name, 1, 20)
    df1_g_10, sent1_g_10, _ = process_sheet(content1, file1_name, 1, 10)
    df2_g_20, sent2_g_20, err2_g = process_sheet(content2, file2_name, 1, 20)
    df2_g_10, sent2_g_10, _ = process_sheet(content2, file2_name, 1, 10)

    if (err1_y and err1_g) or (err2_y and err2_g):
        flash('Не удалось обработать листы в одном из файлов')
        clear_user_session()
        return redirect(url_for('index'))

    comp_y_20, err_y20, stats1_y20, stats2_y20 = comparator.compare_serp_data_with_sentiment(
        df1_y_20, df2_y_20, sent1_y_20, sent2_y_20, label1, label2, top_n=20) if df1_y_20 is not None and df2_y_20 is not None else (None, "Нет данных", None, None)
    comp_y_10, err_y10, stats1_y10, stats2_y10 = comparator.compare_serp_data_with_sentiment(
        df1_y_10, df2_y_10, sent1_y_10, sent2_y_10, label1, label2, top_n=10) if df1_y_10 is not None and df2_y_10 is not None else (None, "Нет данных", None, None)

    comp_g_20, err_g20, stats1_g20, stats2_g20 = comparator.compare_serp_data_with_sentiment(
        df1_g_20, df2_g_20, sent1_g_20, sent2_g_20, label1, label2, top_n=20) if df1_g_20 is not None and df2_g_20 is not None else (None, "Нет данных", None, None)
    comp_g_10, err_g10, stats1_g10, stats2_g10 = comparator.compare_serp_data_with_sentiment(
        df1_g_10, df2_g_10, sent1_g_10, sent2_g_10, label1, label2, top_n=10) if df1_g_10 is not None and df2_g_10 is not None else (None, "Нет данных", None, None)

    # Сохраняем stats для стартовых диаграмм
    session['stats_y20_1'] = stats1_y20
    session['stats_y20_2'] = stats2_y20
    session['stats_y10_1'] = stats1_y10
    session['stats_y10_2'] = stats2_y10
    session['stats_g20_1'] = stats1_g20
    session['stats_g20_2'] = stats2_g20
    session['stats_g10_1'] = stats1_g10
    session['stats_g10_2'] = stats2_g10

    def extract_lists(comp_df, label1, label2):
        improvements = []
        deteriorations = []
        new_urls = []
        dropped_urls = []
        sentiment_changes = []
        if comp_df is None or comp_df.empty:
            return improvements, deteriorations, new_urls, dropped_urls, sentiment_changes
        data_df = comp_df[~comp_df['Запрос'].str.startswith('СТАТИСТИКА', na=False)]
        data_df = data_df[data_df['Запрос'].notna() & (data_df['Запрос'] != '')]
        for _, row in data_df.iterrows():
            change = row['Изменение']
            query = row['Запрос']
            url = row['URL']
            sentiment1 = row[f'Тональность_{label1}']
            sentiment2 = row[f'Тональность_{label2}']
            if sentiment1 != sentiment2 and sentiment1 != "Неопределенная" and sentiment2 != "Неопределенная":
                sentiment_changes.append(f"{query}: {url}\n  {sentiment1} → {sentiment2}")
            if change == "Новый в ТОП":
                new_urls.append({
                    'query': query,
                    'url': url,
                    'sentiment': sentiment2,
                    'display': f"{query}: {url}"
                })
            elif change == "Выпал из ТОП":
                dropped_urls.append({
                    'query': query,
                    'url': url,
                    'sentiment': sentiment1,
                    'display': f"{query}: {url}"
                })
            elif isinstance(change, (int, float)):
                if change > 0:
                    improvements.append({
                        'query': query,
                        'url': url,
                        'change': f"+{change}",
                        'sentiment': sentiment2,
                        'display': f"{query}: {url} (+{change})"
                    })
                elif change < 0:
                    deteriorations.append({
                        'query': query,
                        'url': url,
                        'change': str(change),
                        'sentiment': sentiment2,
                        'display': f"{query}: {url} ({change})"
                    })
        return improvements, deteriorations, new_urls, dropped_urls, sentiment_changes

    imp_y20, det_y20, new_y20, drop_y20, sent_y20 = extract_lists(comp_y_20, label1, label2)
    imp_g20, det_g20, new_g20, drop_g20, sent_g20 = extract_lists(comp_g_20, label1, label2)
    imp_y10, det_y10, new_y10, drop_y10, sent_y10 = extract_lists(comp_y_10, label1, label2)
    imp_g10, det_g10, new_g10, drop_g10, sent_g10 = extract_lists(comp_g_10, label1, label2)

    imp_total20 = imp_y20 + imp_g20
    det_total20 = det_y20 + det_g20
    new_total20 = new_y20 + new_g20
    drop_total20 = drop_y20 + drop_g20
    sent_total20 = sent_y20 + sent_g20
    imp_total10 = imp_y10 + imp_g10
    det_total10 = det_y10 + det_g10
    new_total10 = new_y10 + new_g10
    drop_total10 = drop_y10 + drop_g10
    sent_total10 = sent_y10 + sent_g10

    analysis_y20 = comparator.generate_analysis(imp_y20, det_y20, new_y20, drop_y20, sent_y20, 20)
    analysis_g20 = comparator.generate_analysis(imp_g20, det_g20, new_g20, drop_g20, sent_g20, 20)
    analysis_total20 = comparator.generate_analysis(imp_total20, det_total20, new_total20, drop_total20, sent_total20, 20)
    analysis_y10 = comparator.generate_analysis(imp_y10, det_y10, new_y10, drop_y10, sent_y10, 10)
    analysis_g10 = comparator.generate_analysis(imp_g10, det_g10, new_g10, drop_g10, sent_g10, 10)
    analysis_total10 = comparator.generate_analysis(imp_total10, det_total10, new_total10, drop_total10, sent_total10, 10)

    # Генерируем Summary в формате DataFrame
    summary_y_20_df, summary_y_20_details = comparator.generate_summary_dataframe(comp_y_20, label1, label2, 20) if comp_y_20 is not None else (pd.DataFrame([{'Метрика': 'Нет данных', 'Значение': 'Нет данных для Яндекса ТОП-20', 'Детали': ''}]), {})
    summary_y_10_df, summary_y_10_details = comparator.generate_summary_dataframe(comp_y_10, label1, label2, 10) if comp_y_10 is not None else (pd.DataFrame([{'Метрика': 'Нет данных', 'Значение': 'Нет данных для Яндекса ТОП-10', 'Детали': ''}]), {})
    summary_g_20_df, summary_g_20_details = comparator.generate_summary_dataframe(comp_g_20, label1, label2, 20) if comp_g_20 is not None else (pd.DataFrame([{'Метрика': 'Нет данных', 'Значение': 'Нет данных для Google ТОП-20', 'Детали': ''}]), {})
    summary_g_10_df, summary_g_10_details = comparator.generate_summary_dataframe(comp_g_10, label1, label2, 10) if comp_g_10 is not None else (pd.DataFrame([{'Метрика': 'Нет данных', 'Значение': 'Нет данных для Google ТОП-10', 'Детали': ''}]), {})

    def sum_stats(stats_list):
        total = {}
        total_count = 0
        for stats in stats_list:
            if stats:
                for k, v in stats.items():
                    if k not in total:
                        total[k] = {'count': 0, 'name': v['name']}
                    total[k]['count'] += v['count']
        total_count = sum(v['count'] for v in total.values())
        for k in total:
            total[k]['percentage'] = round(total[k]['count'] / total_count * 100, 2) if total_count else 0
        return total, total_count

    total_stats1_20, total_count1_20 = sum_stats([stats1_y20, stats1_g20])
    total_stats2_20, total_count2_20 = sum_stats([stats2_y20, stats2_g20])
    total_stats1_10, total_count1_10 = sum_stats([stats1_y10, stats1_g10])
    total_stats2_10, total_count2_10 = sum_stats([stats2_y10, stats2_g10])

    # === СТАРТОВЫЕ МЕТРИКИ ПРОЕКТА ===
    baseline_stats = {}

    def get_total_urls_from_stats(stats):
        """Получает общее количество URL из stats файла"""
        if not stats:
            return 0
        total = 0
        for cat in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
            if cat in stats and isinstance(stats[cat], dict):
                total += stats[cat].get('count', 0)
        return total

    def create_baseline_chart_data(baseline, stats_from_file):
        """Преобразует baseline метрики в формат для диаграмм с учетом реальных URL"""
        if not baseline:
            return None

        # Получаем общее количество URL из файла сравнения
        total_urls = get_total_urls_from_stats(stats_from_file)

        # Если total_urls = 0, используем сохраненное значение из baseline
        if total_urls == 0:
            total_urls = baseline.get('total_urls', 0) or 120  # fallback

        input_type = baseline.get('input_type', 'percentage')

        def calc_values(val):
            """Рассчитывает percentage и count в зависимости от типа ввода"""
            if val is None:
                val = 0

            if input_type == 'count':
                # Введено количество — считаем процент
                count = int(val)
                percentage = round((count / total_urls * 100), 2) if total_urls > 0 else 0
            else:
                # Введены проценты — считаем количество
                percentage = val
                count = int(total_urls * percentage / 100) if total_urls > 0 else 0

            return {'percentage': percentage, 'count': count}

        return {
            'client_site': calc_values(baseline['client_site']),
            'positive': calc_values(baseline['positive']),
            'neutral': calc_values(baseline['neutral']),
            'negative': calc_values(baseline['negative']),
            'irrelevant': calc_values(baseline['irrelevant'])
        }

    if project:
        # Получаем все baseline метрики проекта с учетом реальных URL из файлов
        baseline_stats_map = {
            'y20': (stats1_y20, 'yandex', 20),
            'y10': (stats1_y10, 'yandex', 10),
            'g20': (stats1_g20, 'google', 20),
            'g10': (stats1_g10, 'google', 10)
        }

        for key, (file_stats, engine, top_n) in baseline_stats_map.items():
            baseline = project.get_baseline(engine, top_n)
            if baseline and any([baseline['client_site'], baseline['positive'],
                                baseline['neutral'], baseline['negative'], baseline['irrelevant']]):
                baseline_stats[key] = create_baseline_chart_data(baseline, file_stats)

        # Суммируем стартовые метрики Яндекс + Google для общих диаграмм
        def sum_baseline_stats(stats_list):
            """Суммирует несколько baseline stats"""
            if not stats_list or all(s is None for s in stats_list):
                return None
            categories = ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']
            result = {}
            for cat in categories:
                total_count = sum(s.get(cat, {}).get('count', 0) for s in stats_list if s)
                total_percentage = sum(s.get(cat, {}).get('percentage', 0) for s in stats_list if s)
                result[cat] = {'count': total_count, 'percentage': round(total_percentage, 2)}
            return result

        # Суммируем для ТОП-20 и ТОП-10
        if baseline_stats.get('y20') or baseline_stats.get('g20'):
            baseline_stats['total20'] = sum_baseline_stats([baseline_stats.get('y20'), baseline_stats.get('g20')])

        if baseline_stats.get('y10') or baseline_stats.get('g10'):
            baseline_stats['total10'] = sum_baseline_stats([baseline_stats.get('y10'), baseline_stats.get('g10')])

    def create_and_save_chart(stats1, stats2, label1, label2, title, filename_prefix, top_n, stats_start=None):
        """Создает диаграмму с 2 или 3 столбцами (если есть стартовые метрики)"""
        # Для total диаграмм (Яндекс+Google) используем горизонтальную ориентацию
        if filename_prefix.startswith('chart_total'):
            chart_buffer = comparator.create_horizontal_chart(
                stats1, stats2, label1, label2, title,
                stats_start=stats_start, label_start='Старт'
            )
        elif stats_start:
            # Для обычных диаграмм со стартовыми метриками — вертикальная трехстолбцовая
            chart_buffer = comparator.create_three_column_chart(
                stats_start, stats1, stats2,
                label_start='Старт', label1=label1, label2=label2,
                title=''  # без заголовка
            )
        else:
            # Обычная двухстолбцовая диаграмма
            chart_buffer = comparator.create_comparison_chart(stats1, stats2, label1, label2, title)

        if chart_buffer:
            chart_filename = f"{filename_prefix}.png"
            chart_path = os.path.join(session_path, chart_filename)
            with open(chart_path, 'wb') as f:
                f.write(chart_buffer.getvalue())
            logger.info(f"Диаграмма сохранена: {chart_path}")
            return url_for('uploaded_file', filename=f"{session_id}/{chart_filename}")
        else:
            logger.warning(f"Диаграмма {filename_prefix} не создана (нет данных)")
            return None

    # Создаем диаграммы с учетом стартовых метрик
    chart_y20_url = create_and_save_chart(stats1_y20, stats2_y20, label1, label2,
                                          "Сравнение тональности (Яндекс, ТОП-20)", "chart_yandex20", 20,
                                          baseline_stats.get('y20'))
    chart_g20_url = create_and_save_chart(stats1_g20, stats2_g20, label1, label2,
                                          "Сравнение тональности (Google, ТОП-20)", "chart_google20", 20,
                                          baseline_stats.get('g20'))
    chart_total20_url = create_and_save_chart(total_stats1_20, total_stats2_20, label1, label2,
                                              "Общая тональность (Яндекс+Google, ТОП-20)", "chart_total20", 20,
                                              baseline_stats.get('total20'))

    chart_y10_url = create_and_save_chart(stats1_y10, stats2_y10, label1, label2,
                                          "Сравнение тональности (Яндекс, ТОП-10)", "chart_yandex10", 10,
                                          baseline_stats.get('y10'))
    chart_g10_url = create_and_save_chart(stats1_g10, stats2_g10, label1, label2,
                                          "Сравнение тональности (Google, ТОП-10)", "chart_google10", 10,
                                          baseline_stats.get('g10'))
    chart_total10_url = create_and_save_chart(total_stats1_10, total_stats2_10, label1, label2,
                                              "Общая тональность (Яндекс+Google, ТОП-10)", "chart_total10", 10,
                                              baseline_stats.get('total10'))

    # Сохраняем пути диаграмм в сессии
    if chart_y20_url:
        session['chart_yandex20_path'] = os.path.join(session_path, 'chart_yandex20.png')
    if chart_g20_url:
        session['chart_google20_path'] = os.path.join(session_path, 'chart_google20.png')
    if chart_total20_url:
        session['chart_total20_path'] = os.path.join(session_path, 'chart_total20.png')
    if chart_y10_url:
        session['chart_yandex10_path'] = os.path.join(session_path, 'chart_yandex10.png')
    if chart_g10_url:
        session['chart_google10_path'] = os.path.join(session_path, 'chart_google10.png')
    if chart_total10_url:
        session['chart_total10_path'] = os.path.join(session_path, 'chart_total10.png')

    # === СЕТКА ТОП-10 ===
    def create_top10_grid_png(comp_df, label1, label2, search_engine, filename_prefix):
        """Создание двух PNG сеток ТОП-10 - отдельно для каждого файла"""
        if comp_df is None or comp_df.empty:
            return None, None
        
        # Фильтруем только реальные запросы
        real_df = comp_df[comp_df.apply(is_real_query, axis=1)].copy()
        if real_df.empty:
            return None, None
        
        # Цвета тональности - умеренно яркие
        sentiment_colors = {
            'client_site': '#F9CB9C',   # Домашний сайт - оранжевый
            'positive': '#B5D7A8',      # Позитивная - зеленый
            'negative': '#EA9999',      # Негативная - красный
            'irrelevant': '#B7B7B7',    # Нерелевантная - темно-серый
            'neutral': '#A0C5E8',       # Нейтральная - голубой
            'unknown': '#808080'        # Неопределенная - темно-серый
        }
        
        # Маппинг русских названий на английские ключи
        sentiment_mapping = {
            'Домашний сайт': 'client_site',
            'Позитивная': 'positive',
            'Негативная': 'negative',
            'Нерелевантная': 'irrelevant',
            'Нейтральная': 'neutral',
            'Неопределенная': 'unknown'
        }
        
        queries = real_df['Запрос'].unique()
        
        def create_single_grid(label, neg_col_prefix, pos_col_prefix, file_num):
            """Создаем одну сетку для одного файла"""
            table_data = []
            total_negative = 0
            
            neg_col = f'Тональность_{label}'
            pos_col = f'Позиция_{label}'
            
            for query in queries:
                query_data = real_df[real_df['Запрос'] == query]
                
                # Считаем негативные только в позициях 1-10
                if neg_col in query_data.columns and pos_col in query_data.columns:
                    # DEBUG: показываем уникальные значения тональности
                    unique_sentiments = query_data[neg_col].unique()
                    pos_values = pd.to_numeric(query_data[pos_col], errors='coerce')
                    negative_in_top10 = query_data[(query_data[neg_col] == 'Негативная') & 
                                                   (pos_values >= 1) & (pos_values <= 10)]
                    negative_count = len(negative_in_top10)
                    # DEBUG
                    all_negative = query_data[query_data[neg_col] == 'Негативная']
                    print(f"DEBUG query={query[:30]}: sentiments={unique_sentiments}, all_neg={len(all_negative)}, in_top10={negative_count}", flush=True)
                else:
                    negative_count = 0
                    print(f"DEBUG query={query[:30]}: columns not found", flush=True)
                total_negative += negative_count
                
                row_data = [query[:40], f"{negative_count}/10"]
                row_colors = ['white', 'white']
                
                # Позиции 1-10
                for pos in range(1, 11):
                    if pos_col in query_data.columns:
                        pos_values = pd.to_numeric(query_data[pos_col], errors='coerce')
                        pos_data = query_data[pos_values == pos]
                        if not pos_data.empty:
                            sentiment_ru = pos_data.iloc[0].get(neg_col, 'unknown')
                            sentiment_en = sentiment_mapping.get(sentiment_ru, sentiment_ru)
                            color = sentiment_colors.get(sentiment_en, '#808080')
                            row_data.append('')
                            row_colors.append(color)
                        else:
                            row_data.append('')
                            row_colors.append('#f0f0f0')
                    else:
                        row_data.append('')
                        row_colors.append('#f0f0f0')
                
                table_data.append((row_data, row_colors))
            
            # Создаем изображение с динамической высотой
            # При большем масштабе нужна больше высота
            # Баланс: мало запросов - компактно, 4+ - одинаковые квадратные ячейки
            num_queries = len(queries)
            if num_queries <= 3:
                height_scale = 0.5  # мало запросов - компактно
            else:
                height_scale = 0.8  # 4+ запросов - одинаковые квадратные ячейки
            fig_height = max(3.5, num_queries * 0.45 * height_scale + 1.5)
            fig, ax = plt.subplots(figsize=(14, fig_height))
            ax.axis('off')
            
            # Заголовок - заменяем точки на похожий символ чтобы не было URL
            safe_label = label.replace('.', '·')
            title = f'Сетка тональности {search_engine} ТОП-10\n{file_num}: {safe_label}'
            ax.set_title(title, fontsize=14, fontweight='bold', color='#ff6b35', pad=15)
            
            # Таблица
            headers = ['Запрос', 'Негативные', '1', '2', '3', '4', '5', '6', '7', '8', '9', '10']
            cell_data = [row[0] for row in table_data]
            cell_colors = [row[1] for row in table_data]
            
            # Настройка ширины колонок - Запрос шире
            col_widths = [0.25, 0.08] + [0.055] * 10  # Запрос 25%, Негативные 8%, остальные ~5.5%
            
            table = ax.table(cellText=cell_data, colLabels=headers, cellLoc='center',
                            loc='center', bbox=[0, 0, 1, 1],
                            colWidths=col_widths)
            
            table.auto_set_font_size(False)
            table.set_fontsize(12)  # увеличили шрифт для читаемости
            # Компенсируем увеличение шрифта уменьшением масштаба
            table.scale(1, height_scale * 0.75)
            
            # Цвета ячеек и границы
            for i, row_colors in enumerate(cell_colors):
                for j, color in enumerate(row_colors):
                    cell = table[(i+1, j)]
                    cell.set_facecolor(color)
                    cell.set_edgecolor('black')  # Чёрные границы
                    cell.set_linewidth(0.5)
                    # Для колонки Запрос добавляем перенос текста
                    if j == 0:
                        cell.set_text_props(wrap=True, ha='left')
                    if j == 1:
                        cell.set_text_props(color='#cc0000', fontweight='bold')
            
            # Заголовок
            for j in range(len(headers)):
                cell = table[(0, j)]
                cell.set_facecolor('#ff6b35')
                cell.set_text_props(color='white', fontweight='bold')
            
            # Легенда с цветными квадратами matplotlib
            from matplotlib.patches import Rectangle
            legend_items = [
                ('#F9CB9C', 'Домашний'),
                ('#B5D7A8', 'Позитивный'),
                ('#EA9999', 'Негативный'),
                ('#B7B7B7', 'Нерелевантный'),
                ('#A0C5E8', 'Нейтральный')
            ]
            
            # Создаем легенду вручную
            legend_x = 0.12
            legend_y = 0.02
            box_height = 0.025
            box_width = 0.025
            spacing = 0.18
            
            for i, (color, label) in enumerate(legend_items):
                x_pos = legend_x + i * spacing
                # Цветной квадрат
                rect = Rectangle((x_pos, legend_y), box_width, box_height, 
                                facecolor=color, edgecolor='#333', linewidth=0.5,
                                transform=fig.transFigure, clip_on=False)
                fig.patches.append(rect)
                # Текст
                fig.text(x_pos + box_width + 0.01, legend_y + box_height/2, label, 
                        ha='left', va='center', fontsize=8, transform=fig.transFigure)
            
            plt.tight_layout()
            plt.subplots_adjust(bottom=0.08)
            
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            buffer.seek(0)
            plt.close()
            
            return buffer, total_negative
        
        # Создаем две сетки
        buffer1, neg1 = create_single_grid(label1, f'Тональность_{label1}', f'Позиция_{label1}', "1")
        buffer2, neg2 = create_single_grid(label2, f'Тональность_{label2}', f'Позиция_{label2}', "2")
        
        # Сохраняем файл 1
        png1_filename = f"{filename_prefix}_file1.png"
        png1_path = os.path.join(session_path, png1_filename)
        with open(png1_path, 'wb') as f:
            f.write(buffer1.getvalue())
        
        # Сохраняем файл 2
        png2_filename = f"{filename_prefix}_file2.png"
        png2_path = os.path.join(session_path, png2_filename)
        with open(png2_path, 'wb') as f:
            f.write(buffer2.getvalue())
        
        # URL
        png1_url = url_for('uploaded_file', filename=f"{session_id}/{png1_filename}")
        png2_url = url_for('uploaded_file', filename=f"{session_id}/{png2_filename}")
        
        stats = {
            'total_queries': len(queries),
            'total_negative_file1': neg1,
            'total_negative_file2': neg2
        }
        
        return {'file1': png1_url, 'file2': png2_url}, stats

    def create_top10_grid_pptx(grid_png_path1, grid_png_path2, search_engine, session_path, num_queries, label1, label2):
        """
        Создает PowerPoint презентацию с сетками ТОП-10.
        Если < 7 запросов — обе сетки на одном слайде (file1 сверху, file2 снизу).
        Если >= 7 запросов — на разных слайдах.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        if num_queries < 7:
            # Один слайд с двумя изображениями
            blank_slide_layout = prs.slide_layouts[6]  # Пустой слайд
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Устанавливаем фон слайда
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 246, 248)  # #f5f6f8
            
            # Заголовок слайда
            title_text = f"Аналитика запросов {label1} vs {label2}"
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Уменьшаем шрифт, если заголовок слишком длинный
            if len(title_text) > 50:
                title_frame.paragraphs[0].font.size = Pt(16)
            elif len(title_text) > 40:
                title_frame.paragraphs[0].font.size = Pt(18)
            else:
                title_frame.paragraphs[0].font.size = Pt(24)
            
            # Название поисковика в правом верхнем углу (лого)
            base_dir = os.path.dirname(os.path.abspath(__file__))
            if search_engine == "Яндекс":
                logo_path = os.path.join(base_dir, 'static', 'images', 'yandex_logo.png')
            elif search_engine == "Google":
                logo_path = os.path.join(base_dir, 'static', 'images', 'google_logo.png')
            else:
                logo_path = None
            
            if logo_path and os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(11.5), Inches(0.15), width=Inches(1.25))
            else:
                se_box = slide.shapes.add_textbox(Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.5))
                se_frame = se_box.text_frame
                se_frame.text = search_engine
                se_frame.paragraphs[0].font.size = Pt(18)
                se_frame.paragraphs[0].font.bold = True
                se_frame.paragraphs[0].font.color.rgb = RGBColor(251, 130, 39)  # Оранжевый цвет
                se_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Подпись файла 1 (сверху)
            label1_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.333), Inches(0.3))
            label1_frame = label1_box.text_frame
            label1_frame.text = f"Файл 1: {label1}"
            label1_frame.paragraphs[0].font.size = Pt(12)
            label1_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Сетка файла 1 (сверху, ~45% высоты)
            if grid_png_path1 and os.path.exists(grid_png_path1):
                # Добавляем закругленные углы к изображению
                rounded_grid1_buffer = add_rounded_corners_to_image(grid_png_path1, radius=30)
                slide.shapes.add_picture(rounded_grid1_buffer, Inches(0.5), Inches(1.1), 
                                        width=Inches(12.333), height=Inches(2.8))
            
            # Подпись файла 2 (после первой сетки)
            label2_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.3))
            label2_frame = label2_box.text_frame
            label2_frame.text = f"Файл 2: {label2}"
            label2_frame.paragraphs[0].font.size = Pt(12)
            label2_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Сетка файла 2 (снизу)
            if grid_png_path2 and os.path.exists(grid_png_path2):
                # Добавляем закругленные углы к изображению
                rounded_grid2_buffer = add_rounded_corners_to_image(grid_png_path2, radius=30)
                slide.shapes.add_picture(rounded_grid2_buffer, Inches(0.5), Inches(4.35), 
                                        width=Inches(12.333), height=Inches(2.8))
        else:
            # Два слайда — отдельно для каждого файла
            # Слайд 1 — Файл 1
            blank_slide_layout = prs.slide_layouts[6]
            slide1 = prs.slides.add_slide(blank_slide_layout)
            
            # Устанавливаем фон слайда
            background1 = slide1.background
            fill1 = background1.fill
            fill1.solid()
            fill1.fore_color.rgb = RGBColor(245, 246, 248)  # #f5f6f8
            
            title1_text = f"Аналитика запросов {label1} vs {label2}"
            title1_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.5))
            title1_frame = title1_box.text_frame
            title1_frame.text = title1_text
            title1_frame.paragraphs[0].font.bold = True
            title1_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Уменьшаем шрифт, если заголовок слишком длинный
            if len(title1_text) > 50:
                title1_frame.paragraphs[0].font.size = Pt(16)
            elif len(title1_text) > 40:
                title1_frame.paragraphs[0].font.size = Pt(18)
            else:
                title1_frame.paragraphs[0].font.size = Pt(24)
            
            # Название поисковика в правом верхнем углу (лого)
            base_dir = os.path.dirname(os.path.abspath(__file__))
            if search_engine == "Яндекс":
                logo_path = os.path.join(base_dir, 'static', 'images', 'yandex_logo.png')
            elif search_engine == "Google":
                logo_path = os.path.join(base_dir, 'static', 'images', 'google_logo.png')
            else:
                logo_path = None
            
            if logo_path and os.path.exists(logo_path):
                slide1.shapes.add_picture(logo_path, Inches(11.5), Inches(0.15), width=Inches(1.25))
            else:
                se1_box = slide1.shapes.add_textbox(Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.5))
                se1_frame = se1_box.text_frame
                se1_frame.text = search_engine
                se1_frame.paragraphs[0].font.size = Pt(18)
                se1_frame.paragraphs[0].font.bold = True
                se1_frame.paragraphs[0].font.color.rgb = RGBColor(251, 130, 39)  # Оранжевый цвет
                se1_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            if grid_png_path1 and os.path.exists(grid_png_path1):
                # Добавляем закругленные углы к изображению
                rounded_grid1_buffer = add_rounded_corners_to_image(grid_png_path1, radius=30)
                slide1.shapes.add_picture(rounded_grid1_buffer, Inches(0.5), Inches(0.9), 
                                         width=Inches(12.333), height=Inches(6.0))
            
            # Слайд 2 — Файл 2
            slide2 = prs.slides.add_slide(blank_slide_layout)
            
            # Устанавливаем фон слайда
            background2 = slide2.background
            fill2 = background2.fill
            fill2.solid()
            fill2.fore_color.rgb = RGBColor(245, 246, 248)  # #f5f6f8
            
            title2_text = f"Аналитика запросов {label1} vs {label2}"
            title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.5))
            title2_frame = title2_box.text_frame
            title2_frame.text = title2_text
            title2_frame.paragraphs[0].font.bold = True
            title2_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            # Уменьшаем шрифт, если заголовок слишком длинный
            if len(title2_text) > 50:
                title2_frame.paragraphs[0].font.size = Pt(16)
            elif len(title2_text) > 40:
                title2_frame.paragraphs[0].font.size = Pt(18)
            else:
                title2_frame.paragraphs[0].font.size = Pt(24)
            
            # Название поисковика в правом верхнем углу (лого)
            base_dir = os.path.dirname(os.path.abspath(__file__))
            if search_engine == "Яндекс":
                logo_path = os.path.join(base_dir, 'static', 'images', 'yandex_logo.png')
            elif search_engine == "Google":
                logo_path = os.path.join(base_dir, 'static', 'images', 'google_logo.png')
            else:
                logo_path = None
            
            if logo_path and os.path.exists(logo_path):
                slide2.shapes.add_picture(logo_path, Inches(11.5), Inches(0.15), width=Inches(1.25))
            else:
                se2_box = slide2.shapes.add_textbox(Inches(10.0), Inches(0.2), Inches(3.0), Inches(0.5))
                se2_frame = se2_box.text_frame
                se2_frame.text = search_engine
                se2_frame.paragraphs[0].font.size = Pt(18)
                se2_frame.paragraphs[0].font.bold = True
                se2_frame.paragraphs[0].font.color.rgb = RGBColor(251, 130, 39)  # Оранжевый цвет
                se2_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            if grid_png_path2 and os.path.exists(grid_png_path2):
                # Добавляем закругленные углы к изображению
                rounded_grid2_buffer = add_rounded_corners_to_image(grid_png_path2, radius=30)
                slide2.shapes.add_picture(rounded_grid2_buffer, Inches(0.5), Inches(0.9), 
                                         width=Inches(12.333), height=Inches(6.0))
        
        # Сохраняем PPTX
        pptx_filename = f"grid_{search_engine.lower()}_top10.pptx"
        pptx_path = os.path.join(session_path, pptx_filename)
        prs.save(pptx_path)
        
        return pptx_path

    # Создаем сетки для Яндекс и Google TOP-10 (PNG)
    # DEBUG: Показываем колонки для отладки
    if comp_y_10 is not None:
        print(f"DEBUG comp_y_10 columns: {list(comp_y_10.columns)}", flush=True)
    # Используем label1 и label2 (имена файлов без расширения) для колонок
    grid_y10_url, grid_y10_stats = create_top10_grid_png(comp_y_10, label1, label2, "Яндекс", "grid_yandex10")
    grid_g10_url, grid_g10_stats = create_top10_grid_png(comp_g_10, label1, label2, "Google", "grid_google10")

    grid_y10_path_file1 = os.path.join(session_path, 'grid_yandex10_file1.png') if grid_y10_url else None
    grid_y10_path_file2 = os.path.join(session_path, 'grid_yandex10_file2.png') if grid_y10_url else None
    grid_g10_path_file1 = os.path.join(session_path, 'grid_google10_file1.png') if grid_g10_url else None
    grid_g10_path_file2 = os.path.join(session_path, 'grid_google10_file2.png') if grid_g10_url else None

    if grid_y10_url:
        session['grid_yandex10_file1'] = grid_y10_path_file1
        session['grid_yandex10_file2'] = grid_y10_path_file2
        session['grid_yandex10_stats'] = grid_y10_stats
        # Создаем PowerPoint для Яндекса
        num_queries_y = grid_y10_stats.get('total_queries', 0) if grid_y10_stats else 0
        grid_y10_pptx_path = create_top10_grid_pptx(
            grid_y10_path_file1, grid_y10_path_file2, "Яндекс", 
            session_path, num_queries_y, label1, label2
        )
        session['grid_yandex10_pptx'] = grid_y10_pptx_path
    if grid_g10_url:
        session['grid_google10_file1'] = grid_g10_path_file1
        session['grid_google10_file2'] = grid_g10_path_file2
        session['grid_google10_stats'] = grid_g10_stats
        # Создаем PowerPoint для Google
        num_queries_g = grid_g10_stats.get('total_queries', 0) if grid_g10_stats else 0
        grid_g10_pptx_path = create_top10_grid_pptx(
            grid_g10_path_file1, grid_g10_path_file2, "Google", 
            session_path, num_queries_g, label1, label2
        )
        session['grid_google10_pptx'] = grid_g10_pptx_path

    def create_url_stats_excel(df_20, df_10, label2, sheet_name):
        if df_20 is None or df_10 is None:
            return None

        real_20 = df_20[df_20.apply(is_real_query, axis=1)]
        real_10 = df_10[df_10.apply(is_real_query, axis=1)]
        filtered_20 = real_20[real_20.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
        filtered_10 = real_10[real_10.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]

        all_urls = set(filtered_20['URL'].unique()) | set(filtered_10['URL'].unique())
        data = []
        total_20 = len(filtered_20)
        total_10 = len(filtered_10)

        for url in all_urls:
            row_20 = filtered_20[filtered_20['URL'] == url]
            row_10 = filtered_10[filtered_10['URL'] == url]
            if not row_20.empty:
                sentiment = row_20.iloc[0][f'Тональность_{label2}']
            elif not row_10.empty:
                sentiment = row_10.iloc[0][f'Тональность_{label2}']
            else:
                continue
            count_20 = len(row_20)
            count_10 = len(row_10)
            data.append([url, sentiment, count_20, count_10])

        df_urls = pd.DataFrame(data, columns=["URL", "Тональность", "Count TOP-20", "Count TOP-10"])
        df_urls['% TOP-20'] = (df_urls['Count TOP-20'] / total_20 * 100).round(2) if total_20 else 0
        df_urls['% TOP-10'] = (df_urls['Count TOP-10'] / total_10 * 100).round(2) if total_10 else 0
        df_urls = df_urls[["URL", "Тональность", "Count TOP-20", "% TOP-20", "Count TOP-10", "% TOP-10"]]

        output = BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df_urls.to_excel(writer, index=False, sheet_name=sheet_name)
            workbook = writer.book
            ws = workbook[sheet_name]
            format_excel_headers(ws, 6)
            set_column_widths(ws, {'A': 100, 'B': 20, 'C': 15, 'D': 15, 'E': 15, 'F': 15})
            apply_sentiment_coloring(ws, sentiment_column=2, url_column=1)
        output.seek(0)
        return output

    def create_combined_url_stats_excel(comp_y_20, comp_y_10, comp_g_20, comp_g_10, label2, sheet_name):
        def extract_urls(df_20, df_10, label2):
            if df_20 is None or df_10 is None:
                return pd.DataFrame(), pd.DataFrame()

            real_20 = df_20[df_20.apply(is_real_query, axis=1)]
            real_10 = df_10[df_10.apply(is_real_query, axis=1)]
            filtered_20 = real_20[real_20.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
            filtered_10 = real_10[real_10.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
            return filtered_20, filtered_10

        y_20, y_10 = extract_urls(comp_y_20, comp_y_10, label2)
        g_20, g_10 = extract_urls(comp_g_20, comp_g_10, label2)

        all_urls = set(y_20['URL'].unique()) | set(y_10['URL'].unique()) | set(g_20['URL'].unique()) | set(g_10['URL'].unique())
        data = []
        total_20 = 0
        total_10 = 0

        for url in all_urls:
            y20_count = len(y_20[y_20['URL'] == url]) if not y_20.empty else 0
            y10_count = len(y_10[y_10['URL'] == url]) if not y_10.empty else 0
            g20_count = len(g_20[g_20['URL'] == url]) if not g_20.empty else 0
            g10_count = len(g_10[g_10['URL'] == url]) if not g_10.empty else 0
            total_20 += y20_count + g20_count
            total_10 += y10_count + g10_count

        for url in all_urls:
            y20_count = len(y_20[y_20['URL'] == url]) if not y_20.empty else 0
            y10_count = len(y_10[y_10['URL'] == url]) if not y_10.empty else 0
            g20_count = len(g_20[g_20['URL'] == url]) if not g_20.empty else 0
            g10_count = len(g_10[g_10['URL'] == url]) if not g_10.empty else 0

            sentiment = None
            if not y_20[y_20['URL'] == url].empty:
                sentiment = y_20[y_20['URL'] == url].iloc[0][f'Тональность_{label2}']
            elif not y_10[y_10['URL'] == url].empty:
                sentiment = y_10[y_10['URL'] == url].iloc[0][f'Тональность_{label2}']
            elif not g_20[g_20['URL'] == url].empty:
                sentiment = g_20[g_20['URL'] == url].iloc[0][f'Тональность_{label2}']
            elif not g_10[g_10['URL'] == url].empty:
                sentiment = g_10[g_10['URL'] == url].iloc[0][f'Тональность_{label2}']
            else:
                continue

            data.append([url, sentiment, y20_count + g20_count, y10_count + g10_count])

        df = pd.DataFrame(data, columns=["URL", "Тональность", "Count TOP-20", "Count TOP-10"])
        df['% TOP-20'] = (df['Count TOP-20'] / total_20 * 100).round(2) if total_20 else 0
        df['% TOP-10'] = (df['Count TOP-10'] / total_10 * 100).round(2) if total_10 else 0
        df = df[["URL", "Тональность", "Count TOP-20", "% TOP-20", "Count TOP-10", "% TOP-10"]]

        output = BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name=sheet_name)
            workbook = writer.book
            ws = workbook[sheet_name]
            format_excel_headers(ws, 6)
            set_column_widths(ws, {'A': 100, 'B': 20, 'C': 15, 'D': 15, 'E': 15, 'F': 15})
            apply_sentiment_coloring(ws, sentiment_column=2, url_column=1)
        output.seek(0)
        return output

    url_stats_y_excel = create_url_stats_excel(comp_y_20, comp_y_10, label2, "URL Stats Yandex")
    url_stats_g_excel = create_url_stats_excel(comp_g_20, comp_g_10, label2, "URL Stats Google")
    url_stats_combined_excel = create_combined_url_stats_excel(comp_y_20, comp_y_10, comp_g_20, comp_g_10, label2, "URL Stats Combined")

    def create_dual_sheet_excel(df_20, df_10, sheet_prefix, file_prefix):
        output = BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            if df_20 is not None and not df_20.empty:
                df_20.to_excel(writer, index=False, sheet_name=f"{sheet_prefix} TOP-20")
            if df_10 is not None and not df_10.empty:
                df_10.to_excel(writer, index=False, sheet_name=f"{sheet_prefix} TOP-10")
            workbook = writer.book
            for sheet_name in workbook.sheetnames:
                ws = workbook[sheet_name]
                format_excel_headers(ws, 7)
                # Динамическое определение ширины колонок
                lengths_a = [len(str(ws.cell(row=1, column=1).value))]
                for row in range(2, ws.max_row + 1):
                    val = ws.cell(row=row, column=1).value
                    if val and isinstance(val, str) and not val.startswith('СТАТИСТИКА') and not val.startswith('ОБЩАЯ'):
                        lengths_a.append(len(val))
                ws.column_dimensions['A'].width = max(lengths_a) + 2 if lengths_a else 15
                ws.column_dimensions['B'].width = 60
                for col_letter in ['C', 'D', 'E', 'F', 'G']:
                    col_idx = ord(col_letter) - 64
                    lengths = []
                    for row in range(1, ws.max_row + 1):
                        val = ws.cell(row=row, column=col_idx).value
                        if val is not None:
                            lengths.append(len(str(val)))
                    ws.column_dimensions[col_letter].width = max(lengths) + 2 if lengths else 10
                    for row in range(2, ws.max_row + 1):
                        ws.cell(row=row, column=col_idx).alignment = Alignment(horizontal='center', vertical='center')
                apply_sentiment_coloring(ws, sentiment_column=7, url_column=2)
        output.seek(0)
        return output

    yandex_full_excel = create_dual_sheet_excel(comp_y_20, comp_y_10, "Yandex", "yandex_full")
    google_full_excel = create_dual_sheet_excel(comp_g_20, comp_g_10, "Google", "google_full")

    url_stats_y_excel_path = os.path.join(session_path, 'url_stats_yandex.xlsx') if url_stats_y_excel else None
    url_stats_g_excel_path = os.path.join(session_path, 'url_stats_google.xlsx') if url_stats_g_excel else None
    url_stats_combined_excel_path = os.path.join(session_path, 'url_stats_combined.xlsx') if url_stats_combined_excel else None
    yandex_full_excel_path = os.path.join(session_path, 'yandex_full.xlsx') if yandex_full_excel else None
    google_full_excel_path = os.path.join(session_path, 'google_full.xlsx') if google_full_excel else None

    if url_stats_y_excel:
        with open(url_stats_y_excel_path, 'wb') as f:
            f.write(url_stats_y_excel.getvalue())
    if url_stats_g_excel:
        with open(url_stats_g_excel_path, 'wb') as f:
            f.write(url_stats_g_excel.getvalue())
    if url_stats_combined_excel:
        with open(url_stats_combined_excel_path, 'wb') as f:
            f.write(url_stats_combined_excel.getvalue())
    if yandex_full_excel:
        with open(yandex_full_excel_path, 'wb') as f:
            f.write(yandex_full_excel.getvalue())
    if google_full_excel:
        with open(google_full_excel_path, 'wb') as f:
            f.write(google_full_excel.getvalue())

    excel_buffer = BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
        def write_sheet(df, sheet_name):
            if df is not None and not df.empty:
                df.to_excel(writer, index=False, sheet_name=sheet_name)

        write_sheet(comp_y_20, "Yandex TOP-20")
        write_sheet(comp_y_10, "Yandex TOP-10")
        write_sheet(comp_g_20, "Google TOP-20")
        write_sheet(comp_g_10, "Google TOP-10")

        workbook = writer.book

        def append_total_stats(ws, total_stats1, total_stats2, label1, label2):
            if not total_stats1 and not total_stats2:
                return
            start_row = ws.max_row + 2
            ws.cell(row=start_row, column=1, value="ОБЩАЯ СТАТИСТИКА (Яндекс+Google)")
            start_row += 1
            if total_stats1:
                ws.cell(row=start_row, column=1, value=f"ДЛЯ ФАЙЛА: {label1}")
                start_row += 1
                for sentiment_type in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
                    if sentiment_type in total_stats1:
                        s = total_stats1[sentiment_type]
                        emoji = comparator.sentiment_emoji.get(sentiment_type, '⚫')
                        ws.cell(row=start_row, column=1, value=f"{emoji} {s['name']}")
                        ws.cell(row=start_row, column=2, value=f"{s['count']} URL")
                        ws.cell(row=start_row, column=3, value=f"{s['percentage']}%")
                        start_row += 1
                start_row += 1
            if total_stats2:
                ws.cell(row=start_row, column=1, value=f"ДЛЯ ФАЙЛА: {label2}")
                start_row += 1
                for sentiment_type in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
                    if sentiment_type in total_stats2:
                        s = total_stats2[sentiment_type]
                        emoji = comparator.sentiment_emoji.get(sentiment_type, '⚫')
                        ws.cell(row=start_row, column=1, value=f"{emoji} {s['name']}")
                        ws.cell(row=start_row, column=2, value=f"{s['count']} URL")
                        ws.cell(row=start_row, column=3, value=f"{s['percentage']}%")
                        start_row += 1

        if "Yandex TOP-20" in workbook.sheetnames:
            append_total_stats(workbook["Yandex TOP-20"], total_stats1_20, total_stats2_20, label1, label2)
        if "Yandex TOP-10" in workbook.sheetnames:
            append_total_stats(workbook["Yandex TOP-10"], total_stats1_10, total_stats2_10, label1, label2)
        if "Google TOP-20" in workbook.sheetnames:
            append_total_stats(workbook["Google TOP-20"], total_stats1_20, total_stats2_20, label1, label2)
        if "Google TOP-10" in workbook.sheetnames:
            append_total_stats(workbook["Google TOP-10"], total_stats1_10, total_stats2_10, label1, label2)

        def add_summary_sheet(summary_df, details_dfs, sheet_name):
            """Добавляет Summary лист с объединенными данными"""
            try:
                if summary_df is not None and not summary_df.empty:
                    # Создаем один лист с основными и детальными данными
                    summary_df.to_excel(writer, index=False, sheet_name=sheet_name)
                    ws = writer.book[sheet_name]
                    format_summary_sheet(ws)
            except Exception as e:
                logger.error(f"Ошибка при создании Summary листа {sheet_name}: {e}")
                # Создаем простой лист с информацией об ошибке
                error_df = pd.DataFrame([{'Ошибка': str(e)}])
                error_df.to_excel(writer, index=False, sheet_name=f"{sheet_name} - Ошибка")

        add_summary_sheet(summary_y_20_df, summary_y_20_details, "Summary Yandex TOP-20")
        add_summary_sheet(summary_y_10_df, summary_y_10_details, "Summary Yandex TOP-10")
        add_summary_sheet(summary_g_20_df, summary_g_20_details, "Summary Google TOP-20")
        add_summary_sheet(summary_g_10_df, summary_g_10_details, "Summary Google TOP-10")

        def add_url_stats_sheet(df_20, df_10, label2, sheet_name):
            if df_20 is None or df_10 is None:
                return

            real_20 = df_20[df_20.apply(is_real_query, axis=1)]
            real_10 = df_10[df_10.apply(is_real_query, axis=1)]
            filtered_20 = real_20[real_20.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
            filtered_10 = real_10[real_10.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]

            all_urls = set(filtered_20['URL'].unique()) | set(filtered_10['URL'].unique())
            data = []
            total_20 = len(filtered_20)
            total_10 = len(filtered_10)

            for url in all_urls:
                row_20 = filtered_20[filtered_20['URL'] == url]
                row_10 = filtered_10[filtered_10['URL'] == url]
                if not row_20.empty:
                    sentiment = row_20.iloc[0][f'Тональность_{label2}']
                elif not row_10.empty:
                    sentiment = row_10.iloc[0][f'Тональность_{label2}']
                else:
                    continue
                count_20 = len(row_20)
                count_10 = len(row_10)
                data.append([url, sentiment, count_20, count_10])

            df_urls = pd.DataFrame(data, columns=["URL", "Тональность", "Count TOP-20", "Count TOP-10"])
            df_urls['% TOP-20'] = (df_urls['Count TOP-20'] / total_20 * 100).round(2) if total_20 else 0
            df_urls['% TOP-10'] = (df_urls['Count TOP-10'] / total_10 * 100).round(2) if total_10 else 0
            df_urls = df_urls[["URL", "Тональность", "Count TOP-20", "% TOP-20", "Count TOP-10", "% TOP-10"]]

            df_urls.to_excel(writer, index=False, sheet_name=sheet_name)
            ws = writer.book[sheet_name]
            format_excel_headers(ws, 6)
            set_column_widths(ws, {'A': 100, 'B': 20, 'C': 15, 'D': 15, 'E': 15, 'F': 15})
            apply_sentiment_coloring(ws, sentiment_column=2, url_column=1)

        add_url_stats_sheet(comp_y_20, comp_y_10, label2, "URL Stats Yandex")
        add_url_stats_sheet(comp_g_20, comp_g_10, label2, "URL Stats Google")

        def add_combined_url_stats_sheet(comp_y_20, comp_y_10, comp_g_20, comp_g_10, label2, sheet_name):
            def extract_urls(df_20, df_10, label2):
                if df_20 is None or df_10 is None:
                    return pd.DataFrame(), pd.DataFrame()

                real_20 = df_20[df_20.apply(is_real_query, axis=1)]
                real_10 = df_10[df_10.apply(is_real_query, axis=1)]
                filtered_20 = real_20[real_20.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
                filtered_10 = real_10[real_10.apply(lambda row: is_valid_for_second_file(row, label2), axis=1)]
                return filtered_20, filtered_10

            y_20, y_10 = extract_urls(comp_y_20, comp_y_10, label2)
            g_20, g_10 = extract_urls(comp_g_20, comp_g_10, label2)

            all_urls = set(y_20['URL'].unique()) | set(y_10['URL'].unique()) | set(g_20['URL'].unique()) | set(g_10['URL'].unique())
            data = []
            total_20 = 0
            total_10 = 0

            for url in all_urls:
                y20_count = len(y_20[y_20['URL'] == url]) if not y_20.empty else 0
                y10_count = len(y_10[y_10['URL'] == url]) if not y_10.empty else 0
                g20_count = len(g_20[g_20['URL'] == url]) if not g_20.empty else 0
                g10_count = len(g_10[g_10['URL'] == url]) if not g_10.empty else 0
                total_20 += y20_count + g20_count
                total_10 += y10_count + g10_count

            for url in all_urls:
                y20_count = len(y_20[y_20['URL'] == url]) if not y_20.empty else 0
                y10_count = len(y_10[y_10['URL'] == url]) if not y_10.empty else 0
                g20_count = len(g_20[g_20['URL'] == url]) if not g_20.empty else 0
                g10_count = len(g_10[g_10['URL'] == url]) if not g_10.empty else 0

                sentiment = None
                if not y_20[y_20['URL'] == url].empty:
                    sentiment = y_20[y_20['URL'] == url].iloc[0][f'Тональность_{label2}']
                elif not y_10[y_10['URL'] == url].empty:
                    sentiment = y_10[y_10['URL'] == url].iloc[0][f'Тональность_{label2}']
                elif not g_20[g_20['URL'] == url].empty:
                    sentiment = g_20[g_20['URL'] == url].iloc[0][f'Тональность_{label2}']
                elif not g_10[g_10['URL'] == url].empty:
                    sentiment = g_10[g_10['URL'] == url].iloc[0][f'Тональность_{label2}']
                else:
                    continue

                data.append([url, sentiment, y20_count + g20_count, y10_count + g10_count])

            df = pd.DataFrame(data, columns=["URL", "Тональность", "Count TOP-20", "Count TOP-10"])
            df['% TOP-20'] = (df['Count TOP-20'] / total_20 * 100).round(2) if total_20 else 0
            df['% TOP-10'] = (df['Count TOP-10'] / total_10 * 100).round(2) if total_10 else 0
            df = df[["URL", "Тональность", "Count TOP-20", "% TOP-20", "Count TOP-10", "% TOP-10"]]

            df.to_excel(writer, index=False, sheet_name=sheet_name)
            ws = writer.book[sheet_name]
            format_excel_headers(ws, 6)
            set_column_widths(ws, {'A': 100, 'B': 20, 'C': 15, 'D': 15, 'E': 15, 'F': 15})
            apply_sentiment_coloring(ws, sentiment_column=2, url_column=1)

        add_combined_url_stats_sheet(comp_y_20, comp_y_10, comp_g_20, comp_g_10, label2, "URL Stats Яндекс + Google")

        for sheet_name in workbook.sheetnames:
            ws = workbook[sheet_name]
            if sheet_name.startswith(("Yandex", "Google")) and not sheet_name.startswith("URL Stats") and not sheet_name.startswith("Summary"):
                format_excel_headers(ws, 7)
                # Динамическое определение ширины колонок
                lengths_a = [len(str(ws.cell(row=1, column=1).value))]
                for row in range(2, ws.max_row + 1):
                    val = ws.cell(row=row, column=1).value
                    if val and isinstance(val, str) and not val.startswith('СТАТИСТИКА') and not val.startswith('ОБЩАЯ'):
                        lengths_a.append(len(val))
                ws.column_dimensions['A'].width = max(lengths_a) + 2 if lengths_a else 15
                ws.column_dimensions['B'].width = 60
                for col_letter in ['C', 'D', 'E', 'F', 'G']:
                    col_idx = ord(col_letter) - 64
                    lengths = []
                    for row in range(1, ws.max_row + 1):
                        val = ws.cell(row=row, column=col_idx).value
                        if val is not None:
                            lengths.append(len(str(val)))
                    ws.column_dimensions[col_letter].width = max(lengths) + 2 if lengths else 10
                    for row in range(2, ws.max_row + 1):
                        ws.cell(row=row, column=col_idx).alignment = Alignment(horizontal='center', vertical='center')
                apply_sentiment_coloring(ws, sentiment_column=7, url_column=2)

            elif sheet_name.startswith("URL Stats"):
                pass  # URL Stats уже форматированы ранее
            # Summary листы уже форматированы в add_summary_sheet

    excel_buffer.seek(0)
    analysis_y20_path = os.path.join(session_path, 'analysis_yandex20.txt')
    analysis_g20_path = os.path.join(session_path, 'analysis_google20.txt')
    analysis_total20_path = os.path.join(session_path, 'analysis_total20.txt')
    analysis_y10_path = os.path.join(session_path, 'analysis_yandex10.txt')
    analysis_g10_path = os.path.join(session_path, 'analysis_google10.txt')
    analysis_total10_path = os.path.join(session_path, 'analysis_total10.txt')
    excel_path = os.path.join(session_path, 'comparison.xlsx')

    with open(analysis_y20_path, 'w', encoding='utf-8') as f:
        f.write(analysis_y20)
    with open(analysis_g20_path, 'w', encoding='utf-8') as f:
        f.write(analysis_g20)
    with open(analysis_total20_path, 'w', encoding='utf-8') as f:
        f.write(analysis_total20)
    with open(analysis_y10_path, 'w', encoding='utf-8') as f:
        f.write(analysis_y10)
    with open(analysis_g10_path, 'w', encoding='utf-8') as f:
        f.write(analysis_g10)
    with open(analysis_total10_path, 'w', encoding='utf-8') as f:
        f.write(analysis_total10)
    with open(excel_path, 'wb') as f:
        f.write(excel_buffer.getvalue())

    # Создаем отдельные Summary Excel файлы для скачивания
    def create_summary_excel(summary_df, output_path):
        """Создает отдельный Excel файл с Summary"""
        if summary_df is None or summary_df.empty:
            return None
        output = BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            summary_df.to_excel(writer, index=False, sheet_name='Summary')
            workbook = writer.book
            ws = writer.book['Summary']
            format_summary_sheet(ws)
        output.seek(0)
        with open(output_path, 'wb') as f:
            f.write(output.getvalue())
        return output_path

    summary_y20_excel_path = os.path.join(session_path, 'summary_yandex20.xlsx')
    summary_y10_excel_path = os.path.join(session_path, 'summary_yandex10.xlsx')
    summary_g20_excel_path = os.path.join(session_path, 'summary_google20.xlsx')
    summary_g10_excel_path = os.path.join(session_path, 'summary_google10.xlsx')

    create_summary_excel(summary_y_20_df, summary_y20_excel_path)
    create_summary_excel(summary_y_10_df, summary_y10_excel_path)
    create_summary_excel(summary_g_20_df, summary_g20_excel_path)
    create_summary_excel(summary_g_10_df, summary_g10_excel_path)

    if 'user_id' in session:
        existing = Comparison.query.filter_by(user_id=session['user_id'], session_id=session['session_id'],
                                             file1_name=file1_name, file2_name=file2_name).first()
        if not existing:
            comp = Comparison(
                user_id=session['user_id'],
                project_id=session.get('current_project_id'),
                session_id=session['session_id'],
                file1_name=file1_name,
                file2_name=file2_name
            )
            db.session.add(comp)
            db.session.commit()

            user_perm_dir = os.path.join(app.config['PERMANENT_FOLDER'], str(session['user_id']))
            comp_perm_dir = os.path.join(user_perm_dir, str(comp.id))
            os.makedirs(comp_perm_dir, exist_ok=True)

            def copy_to_permanent(src, dest_filename):
                if src and os.path.exists(src):
                    dest = os.path.join(comp_perm_dir, dest_filename)
                    shutil.copy2(src, dest)
                    return os.path.join('user_data', str(session['user_id']), str(comp.id), dest_filename)
                return None

            comp.excel_path = copy_to_permanent(excel_path, 'comparison.xlsx')
            comp.yandex_excel_path = copy_to_permanent(yandex_full_excel_path, 'yandex_full.xlsx')
            comp.google_excel_path = copy_to_permanent(google_full_excel_path, 'google_full.xlsx')
            comp.analysis_y20_path = copy_to_permanent(analysis_y20_path, 'analysis_yandex20.txt')
            comp.analysis_g20_path = copy_to_permanent(analysis_g20_path, 'analysis_google20.txt')
            comp.analysis_total20_path = copy_to_permanent(analysis_total20_path, 'analysis_total20.txt')
            comp.analysis_y10_path = copy_to_permanent(analysis_y10_path, 'analysis_yandex10.txt')
            comp.analysis_g10_path = copy_to_permanent(analysis_g10_path, 'analysis_google10.txt')
            comp.analysis_total10_path = copy_to_permanent(analysis_total10_path, 'analysis_total10.txt')
            comp.url_y_excel_path = copy_to_permanent(url_stats_y_excel_path, 'url_stats_yandex.xlsx')
            comp.url_g_excel_path = copy_to_permanent(url_stats_g_excel_path, 'url_stats_google.xlsx')
            comp.url_combined_excel_path = copy_to_permanent(url_stats_combined_excel_path, 'url_stats_combined.xlsx')
            comp.chart_y20_path = copy_to_permanent(session.get('chart_yandex20_path'), 'chart_yandex20.png') if chart_y20_url else None
            comp.chart_g20_path = copy_to_permanent(session.get('chart_google20_path'), 'chart_google20.png') if chart_g20_url else None
            comp.chart_total20_path = copy_to_permanent(session.get('chart_total20_path'), 'chart_total20.png') if chart_total20_url else None
            comp.chart_y10_path = copy_to_permanent(session.get('chart_yandex10_path'), 'chart_yandex10.png') if chart_y10_url else None
            comp.chart_g10_path = copy_to_permanent(session.get('chart_google10_path'), 'chart_google10.png') if chart_g10_url else None
            comp.chart_total10_path = copy_to_permanent(session.get('chart_total10_path'), 'chart_total10.png') if chart_total10_url else None
            # Сохраняем стартовые диаграммы
            chart_start_y20_path = session.get('chart_start_y20_path')
            chart_start_g20_path = session.get('chart_start_g20_path')
            chart_start_y10_path = session.get('chart_start_y10_path')
            chart_start_g10_path = session.get('chart_start_g10_path')
            comp.chart_start_y20_path = copy_to_permanent(chart_start_y20_path, 'chart_start_y20.png') if chart_start_y20_path and os.path.exists(chart_start_y20_path) else None
            comp.chart_start_g20_path = copy_to_permanent(chart_start_g20_path, 'chart_start_g20.png') if chart_start_g20_path and os.path.exists(chart_start_g20_path) else None
            comp.chart_start_y10_path = copy_to_permanent(chart_start_y10_path, 'chart_start_y10.png') if chart_start_y10_path and os.path.exists(chart_start_y10_path) else None
            comp.chart_start_g10_path = copy_to_permanent(chart_start_g10_path, 'chart_start_g10.png') if chart_start_g10_path and os.path.exists(chart_start_g10_path) else None
            
            # === ГЕНЕРАЦИЯ PPTX ДИАГРАММ С БАЗЛАЙН ЗНАЧЕНИЯМИ ===
            def get_baseline_values(metric_type):
                baseline = BaselineMetrics.query.filter_by(
                    user_id=session['user_id'],
                    session_id=session['session_id'],
                    metric_type=metric_type
                ).first()
                if baseline:
                    return {
                        'client_site': str(baseline.client_site_value) if baseline.client_site_value else '0',
                        'positive': str(baseline.positive_value) if baseline.positive_value else '0',
                        'neutral': str(baseline.neutral_value) if baseline.neutral_value else '0',
                        'negative': str(baseline.negative_value) if baseline.negative_value else '0',
                        'irrelevant': str(baseline.irrelevant_value) if baseline.irrelevant_value else '0'
                    }
                return None
            
            baseline_y20 = get_baseline_values('yandex_top20')
            baseline_g20 = get_baseline_values('google_top20')
            baseline_y10 = get_baseline_values('yandex_top10')
            baseline_g10 = get_baseline_values('google_top10')
            
            def create_pptx_chart(chart_path, search_engine, top_n, label1, label2, analysis_text, filename_prefix, baseline_values=None):
                if not chart_path or not os.path.exists(chart_path):
                    return None
                # Read chart file as buffer
                with open(chart_path, 'rb') as f:
                    chart_buffer = BytesIO(f.read())
                pptx_buffer = comparator.create_chart_pptx(chart_buffer, search_engine, top_n, label1, label2, analysis_text, baseline_values)
                if pptx_buffer:
                    pptx_filename = f"{filename_prefix}.pptx"
                    pptx_path = os.path.join(session_path, pptx_filename)
                    with open(pptx_path, 'wb') as f:
                        f.write(pptx_buffer.getvalue())
                    return pptx_path
                return None
            
            # Генерируем PPTX для каждой диаграммы
            chart_y20_pptx_path = create_pptx_chart(
                session.get('chart_yandex20_path') if chart_y20_url else None,
                "Яндекс", 20, label1, label2, analysis_y20, "chart_yandex20", baseline_y20
            )
            chart_g20_pptx_path = create_pptx_chart(
                session.get('chart_google20_path') if chart_g20_url else None,
                "Google", 20, label1, label2, analysis_g20, "chart_google20", baseline_g20
            )
            chart_total20_pptx_path = create_pptx_chart(
                session.get('chart_total20_path') if chart_total20_url else None,
                "Яндекс+Google", 20, label1, label2, analysis_total20, "chart_total20", None
            )
            chart_y10_pptx_path = create_pptx_chart(
                session.get('chart_yandex10_path') if chart_y10_url else None,
                "Яндекс", 10, label1, label2, analysis_y10, "chart_yandex10", baseline_y10
            )
            chart_g10_pptx_path = create_pptx_chart(
                session.get('chart_google10_path') if chart_g10_url else None,
                "Google", 10, label1, label2, analysis_g10, "chart_google10", baseline_g10
            )
            chart_total10_pptx_path = create_pptx_chart(
                session.get('chart_total10_path') if chart_total10_url else None,
                "Яндекс+Google", 10, label1, label2, analysis_total10, "chart_total10", None
            )
            
            # Сохраняем пути PPTX в сессию
            if chart_y20_pptx_path:
                session['chart_yandex20_pptx_path'] = chart_y20_pptx_path
            if chart_g20_pptx_path:
                session['chart_google20_pptx_path'] = chart_g20_pptx_path
            if chart_total20_pptx_path:
                session['chart_total20_pptx_path'] = chart_total20_pptx_path
            if chart_y10_pptx_path:
                session['chart_yandex10_pptx_path'] = chart_y10_pptx_path
            if chart_g10_pptx_path:
                session['chart_google10_pptx_path'] = chart_g10_pptx_path
            if chart_total10_pptx_path:
                session['chart_total10_pptx_path'] = chart_total10_pptx_path
            
            # Сохраняем сетки ТОП-10 (file1 - старая дата, file2 - новая дата)
            comp.grid_y10_path = copy_to_permanent(grid_y10_path_file1, 'grid_yandex10_file1.png') if grid_y10_url else None
            comp.grid_y10_file2_path = copy_to_permanent(grid_y10_path_file2, 'grid_yandex10_file2.png') if grid_y10_url else None
            comp.grid_g10_path = copy_to_permanent(grid_g10_path_file1, 'grid_google10_file1.png') if grid_g10_url else None
            comp.grid_g10_file2_path = copy_to_permanent(grid_g10_path_file2, 'grid_google10_file2.png') if grid_g10_url else None
            # Сохраняем PowerPoint презентации сеток
            grid_y10_pptx_path = session.get('grid_yandex10_pptx')
            grid_g10_pptx_path = session.get('grid_google10_pptx')
            print(f"DEBUG PPTX paths from session: y={grid_y10_pptx_path}, g={grid_g10_pptx_path}", flush=True)
            print(f"DEBUG PPTX files exist: y={os.path.exists(grid_y10_pptx_path) if grid_y10_pptx_path else 'N/A'}, g={os.path.exists(grid_g10_pptx_path) if grid_g10_pptx_path else 'N/A'}", flush=True)
            comp.grid_y10_pptx_path = copy_to_permanent(grid_y10_pptx_path, 'grid_yandex10.pptx') if grid_y10_pptx_path and os.path.exists(grid_y10_pptx_path) else None
            comp.grid_g10_pptx_path = copy_to_permanent(grid_g10_pptx_path, 'grid_google10.pptx') if grid_g10_pptx_path and os.path.exists(grid_g10_pptx_path) else None
            print(f"DEBUG Saved PPTX paths to DB: y={comp.grid_y10_pptx_path}, g={comp.grid_g10_pptx_path}", flush=True)
            # Сохраняем Summary Excel файлы
            comp.summary_y20_path = copy_to_permanent(summary_y20_excel_path, 'summary_yandex20.xlsx') if os.path.exists(summary_y20_excel_path) else None
            comp.summary_y10_path = copy_to_permanent(summary_y10_excel_path, 'summary_yandex10.xlsx') if os.path.exists(summary_y10_excel_path) else None
            comp.summary_g20_path = copy_to_permanent(summary_g20_excel_path, 'summary_google20.xlsx') if os.path.exists(summary_g20_excel_path) else None
            comp.summary_g10_path = copy_to_permanent(summary_g10_excel_path, 'summary_google10.xlsx') if os.path.exists(summary_g10_excel_path) else None
            # Сохраняем PPTX диаграммы
            comp.chart_y20_pptx_path = copy_to_permanent(session.get('chart_yandex20_pptx_path'), 'chart_yandex20.pptx') if session.get('chart_yandex20_pptx_path') and os.path.exists(session.get('chart_yandex20_pptx_path')) else None
            comp.chart_g20_pptx_path = copy_to_permanent(session.get('chart_google20_pptx_path'), 'chart_google20.pptx') if session.get('chart_google20_pptx_path') and os.path.exists(session.get('chart_google20_pptx_path')) else None
            comp.chart_total20_pptx_path = copy_to_permanent(session.get('chart_total20_pptx_path'), 'chart_total20.pptx') if session.get('chart_total20_pptx_path') and os.path.exists(session.get('chart_total20_pptx_path')) else None
            comp.chart_y10_pptx_path = copy_to_permanent(session.get('chart_yandex10_pptx_path'), 'chart_yandex10.pptx') if session.get('chart_yandex10_pptx_path') and os.path.exists(session.get('chart_yandex10_pptx_path')) else None
            comp.chart_g10_pptx_path = copy_to_permanent(session.get('chart_google10_pptx_path'), 'chart_google10.pptx') if session.get('chart_google10_pptx_path') and os.path.exists(session.get('chart_google10_pptx_path')) else None
            comp.chart_total10_pptx_path = copy_to_permanent(session.get('chart_total10_pptx_path'), 'chart_total10.pptx') if session.get('chart_total10_pptx_path') and os.path.exists(session.get('chart_total10_pptx_path')) else None

            db.session.commit()

    def prepare_html(df, n=25):
        if df is not None and not df.empty:
            return df.head(n).to_html(classes='table table-striped', index=False, escape=False, float_format=lambda x: '%.0f' % x if pd.notna(x) else '')
        return "<p class='text-muted'>Нет данных</p>"

    table_y20_html = prepare_html(comp_y_20)
    table_y10_html = prepare_html(comp_y_10)
    table_g20_html = prepare_html(comp_g_20)
    table_g10_html = prepare_html(comp_g_10)

    full_table_y20_html = comp_y_20.to_html(classes='table table-striped', index=False, escape=False, float_format=lambda x: '%.0f' % x if pd.notna(x) else '') if comp_y_20 is not None else ""
    full_table_y10_html = comp_y_10.to_html(classes='table table-striped', index=False, escape=False, float_format=lambda x: '%.0f' % x if pd.notna(x) else '') if comp_y_10 is not None else ""
    full_table_g20_html = comp_g_20.to_html(classes='table table-striped', index=False, escape=False, float_format=lambda x: '%.0f' % x if pd.notna(x) else '') if comp_g_20 is not None else ""
    full_table_g10_html = comp_g_10.to_html(classes='table table-striped', index=False, escape=False, float_format=lambda x: '%.0f' % x if pd.notna(x) else '') if comp_g_10 is not None else ""

    def stats_html(stats1, stats2):
        if not stats1 or not stats2:
            return "<p>Нет данных</p>"
        html1 = '<br>'.join([f"{comparator.sentiment_emoji.get(k,'⚫')} {v['name']}: {v['count']} ({v['percentage']}%)" for k, v in stats1.items()])
        html2 = '<br>'.join([f"{comparator.sentiment_emoji.get(k,'⚫')} {v['name']}: {v['count']} ({v['percentage']}%)" for k, v in stats2.items()])
        return html1, html2

    stats_y1_html, stats_y2_html = stats_html(stats1_y20, stats2_y20)
    stats_g1_html, stats_g2_html = stats_html(stats1_g20, stats2_g20)

    session['yandex_full_excel_path'] = yandex_full_excel_path
    session['google_full_excel_path'] = google_full_excel_path
    session['report_excel'] = excel_path
    session['report_analysis_y20'] = analysis_y20_path
    session['report_analysis_g20'] = analysis_g20_path
    session['report_analysis_total20'] = analysis_total20_path
    session['report_analysis_y10'] = analysis_y10_path
    session['report_analysis_g10'] = analysis_g10_path
    session['report_analysis_total10'] = analysis_total10_path
    session['report_url_y_excel'] = url_stats_y_excel_path
    session['report_url_g_excel'] = url_stats_g_excel_path
    session['report_url_combined_excel'] = url_stats_combined_excel_path
    # Сохраняем пути к Summary Excel файлам
    session['summary_y20_excel'] = summary_y20_excel_path if os.path.exists(summary_y20_excel_path) else None
    session['summary_y10_excel'] = summary_y10_excel_path if os.path.exists(summary_y10_excel_path) else None
    session['summary_g20_excel'] = summary_g20_excel_path if os.path.exists(summary_g20_excel_path) else None
    session['summary_g10_excel'] = summary_g10_excel_path if os.path.exists(summary_g10_excel_path) else None
    # Сохраняем имена файлов для формирования названий отчетов
    session['file1_name'] = file1_name
    session['file2_name'] = file2_name

    return render_template('compare.html',
                           file1_name=file1_name,
                           file2_name=file2_name,
                           project=project,
                           table_y20_html=table_y20_html,
                           table_y10_html=table_y10_html,
                           table_g20_html=table_g20_html,
                           table_g10_html=table_g10_html,
                           full_table_y20_html=full_table_y20_html,
                           full_table_y10_html=full_table_y10_html,
                           full_table_g20_html=full_table_g20_html,
                           full_table_g10_html=full_table_g10_html,
                           stats_y1_html=stats_y1_html,
                           stats_y2_html=stats_y2_html,
                           stats_g1_html=stats_g1_html,
                           stats_g2_html=stats_g2_html,
                           chart_y20_url=chart_y20_url,
                           chart_g20_url=chart_g20_url,
                           chart_total20_url=chart_total20_url,
                           chart_y10_url=chart_y10_url,
                           chart_g10_url=chart_g10_url,
                           chart_total10_url=chart_total10_url,
                           # Сетки ТОП-10 (URL для отображения)
                           grid_y10_file1=url_for('uploaded_file', filename=f"{session_id}/grid_yandex10_file1.png") if grid_y10_url else None,
                           grid_y10_file2=url_for('uploaded_file', filename=f"{session_id}/grid_yandex10_file2.png") if grid_y10_url else None,
                           grid_g10_file1=url_for('uploaded_file', filename=f"{session_id}/grid_google10_file1.png") if grid_g10_url else None,
                           grid_g10_file2=url_for('uploaded_file', filename=f"{session_id}/grid_google10_file2.png") if grid_g10_url else None,
                           grid_y10_stats=grid_y10_stats,
                           grid_g10_stats=grid_g10_stats,
                           analysis_y20=analysis_y20,
                           analysis_g20=analysis_g20,
                           analysis_total20=analysis_total20,
                           analysis_y10=analysis_y10,
                           analysis_g10=analysis_g10,
                           analysis_total10=analysis_total10,
                           # Данные для сводки изменений
                           imp_y20=imp_y20,
                           det_y20=det_y20,
                           new_y20=new_y20,
                           drop_y20=drop_y20,
                           imp_y10=imp_y10,
                           det_y10=det_y10,
                           new_y10=new_y10,
                           drop_y10=drop_y10,
                           imp_g20=imp_g20,
                           det_g20=det_g20,
                           new_g20=new_g20,
                           drop_g20=drop_g20,
                           imp_g10=imp_g10,
                           det_g10=det_g10,
                           new_g10=new_g10,
                           drop_g10=drop_g10)

@app.route('/download/<report_type>')
def download(report_type):
    if 'session_id' not in session:
        flash('Сессия не найдена')
        return redirect(url_for('index'))

    # Получаем имена файлов для формирования названия отчета
    file1_name = session.get('file1_name', 'Файл1')
    file2_name = session.get('file2_name', 'Файл2')
    # Убираем расширения из имен файлов
    file1_base = os.path.splitext(file1_name)[0]
    file2_base = os.path.splitext(file2_name)[0]
    
    # Формируем названия отчетов с именами файлов
    report_map = {
        'excel': (f'Сравнение_{file1_base}_vs_{file2_base}.xlsx', session.get('report_excel')),
        'analysis_y20': (f'Анализ_Яндекс_ТОП20_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_y20')),
        'analysis_g20': (f'Анализ_Гугл_ТОП20_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_g20')),
        'analysis_total20': (f'Анализ_Итого_ТОП20_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_total20')),
        'analysis_y10': (f'Анализ_Яндекс_ТОП10_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_y10')),
        'analysis_g10': (f'Анализ_Гугл_ТОП10_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_g10')),
        'analysis_total10': (f'Анализ_Итого_ТОП10_{file1_base}_vs_{file2_base}.txt', session.get('report_analysis_total10')),
        'url_y_excel': (f'URL_статистика_Яндекс_{file1_base}_vs_{file2_base}.xlsx', session.get('report_url_y_excel')),
        'url_g_excel': (f'URL_статистика_Гугл_{file1_base}_vs_{file2_base}.xlsx', session.get('report_url_g_excel')),
        'url_combined_excel': (f'URL_статистика_Яндекс_Гугл_{file1_base}_vs_{file2_base}.xlsx', session.get('report_url_combined_excel')),
        'yandex_full_excel': (f'Яндекс_полный_{file1_base}_vs_{file2_base}.xlsx', session.get('yandex_full_excel_path')),
        'google_full_excel': (f'Гугл_полный_{file1_base}_vs_{file2_base}.xlsx', session.get('google_full_excel_path')),
        # Summary отчеты
        'summary_y20': (f'Summary_Яндекс_ТОП20_{file1_base}_vs_{file2_base}.xlsx', session.get('summary_y20_excel')),
        'summary_y10': (f'Summary_Яндекс_ТОП10_{file1_base}_vs_{file2_base}.xlsx', session.get('summary_y10_excel')),
        'summary_g20': (f'Summary_Гугл_ТОП20_{file1_base}_vs_{file2_base}.xlsx', session.get('summary_g20_excel')),
        'summary_g10': (f'Summary_Гугл_ТОП10_{file1_base}_vs_{file2_base}.xlsx', session.get('summary_g10_excel')),
        # PPTX сетки ТОП-10
        'grid_y10_pptx': (f'Сетка_Яндекс_ТОП10_{file1_base}_vs_{file2_base}.pptx', session.get('grid_yandex10_pptx')),
        'grid_g10_pptx': (f'Сетка_Гугл_ТОП10_{file1_base}_vs_{file2_base}.pptx', session.get('grid_google10_pptx')),
        # PPTX диаграммы
        'chart_y20_pptx': (f'Диаграмма_Яндекс_ТОП20_{file1_base}_vs_{file2_base}.pptx', session.get('chart_yandex20_pptx_path')),
        'chart_g20_pptx': (f'Диаграмма_Гугл_ТОП20_{file1_base}_vs_{file2_base}.pptx', session.get('chart_google20_pptx_path')),
        'chart_total20_pptx': (f'Диаграмма_Итого_ТОП20_{file1_base}_vs_{file2_base}.pptx', session.get('chart_total20_pptx_path')),
        'chart_y10_pptx': (f'Диаграмма_Яндекс_ТОП10_{file1_base}_vs_{file2_base}.pptx', session.get('chart_yandex10_pptx_path')),
        'chart_g10_pptx': (f'Диаграмма_Гугл_ТОП10_{file1_base}_vs_{file2_base}.pptx', session.get('chart_google10_pptx_path')),
        'chart_total10_pptx': (f'Диаграмма_Итого_ТОП10_{file1_base}_vs_{file2_base}.pptx', session.get('chart_total10_pptx_path'))
    }
    if report_type not in report_map:
        flash('Неверный тип отчёта')
        return redirect(url_for('index'))

    filename, filepath = report_map[report_type]
    if not filepath or not os.path.exists(filepath):
        flash('Файл отчёта не найден')
        return redirect(url_for('index'))

    # Определяем mimetype по расширению файла
    mimetype = None
    if filename.endswith('.pptx'):
        mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    
    return send_file(filepath, as_attachment=True, download_name=filename, mimetype=mimetype)

@app.route('/download_chart/<int:top>/<chart_type>')
def download_chart(top, chart_type):
    if 'session_id' not in session:
        flash('Сессия не найдена')
        return redirect(url_for('compare'))

    # Получаем имена файлов для формирования названия
    file1_name = session.get('file1_name', 'Файл1')
    file2_name = session.get('file2_name', 'Файл2')
    file1_base = os.path.splitext(file1_name)[0]
    file2_base = os.path.splitext(file2_name)[0]
    
    # Проверяем если это сетка ТОП-10
    if chart_type.startswith('grid_'):
        grid_type = chart_type.replace('grid_', '')  # yandex или google
        key = f'grid_{grid_type}10_path'
        filepath = session.get(key)
        if not filepath or not os.path.exists(filepath):
            flash('Файл сетки не найден')
            return redirect(url_for('compare'))
        
        search_engine = 'Яндекс' if grid_type == 'yandex' else 'Гугл'
        filename = f'Setka_{search_engine}_TOP10_{file1_base}_vs_{file2_base}.html'
        return send_file(filepath, as_attachment=True, download_name=filename, mimetype='text/html')
    
    # Обычные диаграммы
    if top not in (20, 10):
        flash('Неверный топ')
        return redirect(url_for('compare'))

    # Определяем название поисковика
    search_engine = 'Яндекс' if chart_type == 'yandex' else 'Гугл' if chart_type == 'google' else 'Итого'
    
    key = f'chart_{chart_type}{top}_path'
    filepath = session.get(key)
    if not filepath or not os.path.exists(filepath):
        flash('Файл диаграммы не найден')
        return redirect(url_for('compare'))

    # Формируем название: Диаграмма_Яндекс_ТОП20_имяфайла1_vs_имяфайла2.png
    filename = f'Диаграмма_{search_engine}_ТОП{top}_{file1_base}_vs_{file2_base}.png'
    return send_file(filepath, as_attachment=True, download_name=filename, mimetype='image/png')

@app.route('/download_baseline_chart/<chart_type>')
@login_required
def download_baseline_chart(chart_type):
    """Скачивание стартовых (baseline) диаграмм"""
    if 'session_id' not in session:
        flash('Сессия не найдена')
        return redirect(url_for('compare'))
    
    file1_name = session.get('file1_name', 'Файл1')
    file2_name = session.get('file2_name', 'Файл2')
    file1_base = os.path.splitext(file1_name)[0]
    file2_base = os.path.splitext(file2_name)[0]
    
    # Определяем тип и ключ сессии
    type_map = {
        'start_y20': ('Яндекс_Старт_ТОП20', 'chart_start_y20_path'),
        'start_g20': ('Google_Старт_ТОП20', 'chart_start_g20_path'),
        'start_y10': ('Яндекс_Старт_ТОП10', 'chart_start_y10_path'),
        'start_g10': ('Google_Старт_ТОП10', 'chart_start_g10_path')
    }
    
    if chart_type not in type_map:
        flash('Неверный тип диаграммы')
        return redirect(url_for('compare'))
    
    title, key = type_map[chart_type]
    filepath = session.get(key)
    
    if not filepath or not os.path.exists(filepath):
        flash('Файл диаграммы не найден')
        return redirect(url_for('compare'))
    
    filename = f'Диаграмма_{title}_{file1_base}_vs_{file2_base}.png'
    return send_file(filepath, as_attachment=True, download_name=filename, mimetype='image/png')

@app.route('/download_grid/<grid_type>/<file_num>')
def download_grid(grid_type, file_num):
    """Скачивание PNG сетки ТОП-10 (Яндекс/Google × Файл 1/Файл 2)"""
    if 'session_id' not in session:
        flash('Сессия не найдена. Пожалуйста, загрузите файлы заново.')
        return redirect(url_for('index'))
    
    # Проверяем тип сетки
    if grid_type not in ('yandex', 'google'):
        flash(f'Неверный тип сетки: {grid_type}')
        return redirect(url_for('compare'))
    
    # Проверяем номер файла
    if file_num not in ('file1', 'file2'):
        flash(f'Неверный номер файла: {file_num}')
        return redirect(url_for('compare'))
    
    # Получаем имена файлов для формирования названия
    file1_name = session.get('file1_name', 'Файл1')
    file2_name = session.get('file2_name', 'Файл2')
    file1_base = os.path.splitext(file1_name)[0]
    file2_base = os.path.splitext(file2_name)[0]
    
    # Ключ для сессии
    key = f'grid_{grid_type}10_{file_num}'
    filepath = session.get(key)
    
    if not filepath:
        flash(f'Путь к сетке не найден в сессии (ключ: {key})')
        return redirect(url_for('compare'))
    
    if not os.path.exists(filepath):
        flash(f'Файл не существует: {filepath}')
        return redirect(url_for('compare'))
    
    # Формируем название
    search_engine = 'Яндекс' if grid_type == 'yandex' else 'Google'
    file_label = file1_base if file_num == 'file1' else file2_base
    filename = f'Сетка_{search_engine}_ТОП10_{file_label}.png'
    
    return send_file(filepath, as_attachment=True, download_name=filename, mimetype='image/png')

@app.route('/clear')
def clear():
    clear_user_session()
    flash('Все данные очищены')
    return redirect(url_for('index'))

@app.route('/get_baseline_metrics', methods=['GET'])
@login_required
def get_baseline_metrics_ajax():
    """Получение сохраненных стартовых показателей для AJAX запросов"""
    if 'session_id' not in session:
        return jsonify({})
    
    metric_type = request.args.get('metric_type')
    
    baseline = BaselineMetrics.query.filter_by(
        user_id=session['user_id'],
        session_id=session['session_id'],
        metric_type=metric_type
    ).first()
    
    if baseline:
        return jsonify({
            'total_urls': baseline.total_urls,
            'client_site': baseline.client_site_value,
            'positive': baseline.positive_value,
            'neutral': baseline.neutral_value,
            'negative': baseline.negative_value,
            'irrelevant': baseline.irrelevant_value,
            'is_count': baseline.client_site_is_count
        })
    
    return jsonify({})

@app.route('/get_total_urls_from_files', methods=['GET'])
@login_required
def get_total_urls_from_files():
    """Получение реального количества URL из проанализированных файлов"""
    if 'session_id' not in session:
        return jsonify({'success': False, 'error': 'Сессия не найдена'})
    
    metric_type = request.args.get('metric_type')
    
    # Определяем какие stats использовать
    if metric_type == 'yandex_top20':
        stats1 = session.get('stats_y20_1') or {}
    elif metric_type == 'google_top20':
        stats1 = session.get('stats_g20_1') or {}
    elif metric_type == 'yandex_top10':
        stats1 = session.get('stats_y10_1') or {}
    elif metric_type == 'google_top10':
        stats1 = session.get('stats_g10_1') or {}
    else:
        return jsonify({'success': False, 'error': 'Неверный тип метрик'})
    
    # Считаем сумму count всех категорий
    total_urls = 0
    if stats1:
        for cat in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
            if cat in stats1 and isinstance(stats1[cat], dict):
                total_urls += stats1[cat].get('count', 0)
    
    return jsonify({
        'success': True,
        'total_urls': total_urls,
        'metric_type': metric_type
    })

@app.route('/generate_baseline_chart', methods=['POST'])
@login_required
def generate_baseline_chart():
    """Генерация диаграммы со стартовыми показателями"""
    if 'session_id' not in session:
        return jsonify({'success': False, 'error': 'Сессия не найдена'})
    
    try:
        data = request.get_json()
        metric_type = data.get('metric_type')  # yandex_top20, google_top20, yandex_top10, google_top10
        
        # Получаем стартовые показатели
        baseline = BaselineMetrics.query.filter_by(
            user_id=session['user_id'],
            session_id=session['session_id'],
            metric_type=metric_type
        ).first()
        
        if not baseline:
            return jsonify({'success': False, 'error': 'Стартовые показатели не найдены'})
        
        # Определяем какие данные файлов использовать (СНАЧАЛА, чтобы получить реальное количество URL)
        if metric_type == 'yandex_top20':
            stats1 = session.get('stats_y20_1') or {}
            stats2 = session.get('stats_y20_2') or {}
            title_prefix = 'Яндекс ТОП-20'
            filename_prefix = 'chart_start_y20'
            session_key = 'chart_start_y20_path'
        elif metric_type == 'google_top20':
            stats1 = session.get('stats_g20_1') or {}
            stats2 = session.get('stats_g20_2') or {}
            title_prefix = 'Google ТОП-20'
            filename_prefix = 'chart_start_g20'
            session_key = 'chart_start_g20_path'
        elif metric_type == 'yandex_top10':
            stats1 = session.get('stats_y10_1') or {}
            stats2 = session.get('stats_y10_2') or {}
            title_prefix = 'Яндекс ТОП-10'
            filename_prefix = 'chart_start_y10'
            session_key = 'chart_start_y10_path'
        elif metric_type == 'google_top10':
            stats1 = session.get('stats_g10_1') or {}
            stats2 = session.get('stats_g10_2') or {}
            title_prefix = 'Google ТОП-10'
            filename_prefix = 'chart_start_g10'
            session_key = 'chart_start_g10_path'
        else:
            return jsonify({'success': False, 'error': 'Неверный тип метрик'})
        
        # Получаем реальное количество URL из проанализированных файлов
        # Считаем сумму count всех категорий из stats1
        real_total_urls = 0
        if stats1:
            for cat in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
                if cat in stats1 and isinstance(stats1[cat], dict):
                    real_total_urls += stats1[cat].get('count', 0)
        
        # Если не удалось получить из stats1, берем из формы (для обратной совместимости)
        if real_total_urls == 0:
            real_total_urls = baseline.total_urls or 120  # fallback значение
        
        logger.info(f"Real total URLs from analyzed files: {real_total_urls}")
        
        # Формируем stats для стартовых данных
        # Если режим 'count' — конвертируем количество в проценты используя real_total_urls
        # Если режим 'percentage' — используем сохраненные проценты как есть
        is_count_mode = baseline.client_site_is_count  # Все категории имеют одинаковый тип
        
        stats_start = {
            'client_site': {'percentage': 0, 'count': 0},
            'positive': {'percentage': 0, 'count': 0},
            'neutral': {'percentage': 0, 'count': 0},
            'negative': {'percentage': 0, 'count': 0},
            'irrelevant': {'percentage': 0, 'count': 0},
            'total_urls': real_total_urls
        }
        
        values = {
            'client_site': baseline.client_site_value,
            'positive': baseline.positive_value,
            'neutral': baseline.neutral_value,
            'negative': baseline.negative_value,
            'irrelevant': baseline.irrelevant_value
        }
        
        if real_total_urls > 0:
            for cat in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
                val = values[cat]
                if is_count_mode:
                    # Режим 'count': значение — это количество, конвертируем в проценты
                    stats_start[cat]['count'] = int(val)
                    stats_start[cat]['percentage'] = round((val / real_total_urls) * 100, 2)
                else:
                    # Режим 'percentage': значение — это проценты, конвертируем в count
                    stats_start[cat]['percentage'] = val
                    stats_start[cat]['count'] = int(real_total_urls * val / 100)
        else:
            # Если не удалось получить real_total_urls, используем сохраненные значения как есть
            for cat in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
                stats_start[cat]['percentage'] = values[cat]
        
        logger.info(f"Stats1 keys: {list(stats1.keys()) if stats1 else 'None'}")
        logger.info(f"Stats2 keys: {list(stats2.keys()) if stats2 else 'None'}")
        
        label1 = session.get('file1_name', 'Файл 1')[:15]
        label2 = session.get('file2_name', 'Файл 2')[:15]
        
        # Генерируем диаграмму (без title как на обычных диаграммах)
        chart_buffer = comparator.create_three_column_chart(
            stats_start, stats1, stats2,
            'Старт', label1, label2,
            ''  # пустой title
        )
        
        if chart_buffer:
            # Сохраняем файл
            session_path = os.path.join(app.config['UPLOAD_FOLDER'], session['session_id'])
            chart_filename = f'{filename_prefix}.png'
            chart_path = os.path.join(session_path, chart_filename)
            
            with open(chart_path, 'wb') as f:
                f.write(chart_buffer.getvalue())
            
            session[session_key] = chart_path
            
            # URL для отображения
            chart_url = f"/uploads/{session['session_id']}/{chart_filename}"
            
            return jsonify({
                'success': True,
                'chart_url': chart_url,
                'chart_path': chart_path,
                'message': 'Диаграмма создана'
            })
        else:
            return jsonify({'success': False, 'error': 'Не удалось создать диаграмму'})
            
    except Exception as e:
        logger.error(f"Ошибка генерации стартовой диаграммы: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/save_baseline_metrics_comparison/<int:comparison_id>', methods=['POST'])
@login_required
def save_baseline_metrics_comparison(comparison_id):
    """Сохранение baseline метрик для сохраненного сравнения"""
    try:
        data = request.get_json()
        logger.info(f"Сохранение baseline для comparison_id={comparison_id}: {data}")
        
        # Проверяем, что comparison принадлежит текущему пользователю
        comparison = Comparison.query.filter_by(id=comparison_id, user_id=session['user_id']).first()
        if not comparison:
            return jsonify({'success': False, 'error': 'Сравнение не найдено'})
        
        metric_type = data.get('metric_type')
        
        # Ищем существующую запись
        metric = BaselineMetrics.query.filter_by(
            comparison_id=comparison_id,
            metric_type=metric_type
        ).first()
        
        if not metric:
            metric = BaselineMetrics(
                user_id=session['user_id'],
                session_id=session.get('session_id', ''),
                comparison_id=comparison_id,
                metric_type=metric_type
            )
            db.session.add(metric)
        
        # Обновляем значения
        metric.input_type = data.get('input_type', 'percentage')
        metric.total_urls = int(data.get('total_urls', 0)) if data.get('total_urls') else None
        metric.client_site_value = data.get('client_site_value', '')
        metric.positive_value = data.get('positive_value', '')
        metric.neutral_value = data.get('neutral_value', '')
        metric.negative_value = data.get('negative_value', '')
        metric.irrelevant_value = data.get('irrelevant_value', '')
        
        db.session.commit()
        logger.info(f"Baseline метрики сохранены для comparison_id={comparison_id}")
        
        return jsonify({'success': True, 'message': 'Стартовые показатели сохранены'})
        
    except Exception as e:
        logger.error(f"Ошибка сохранения baseline метрик: {e}")
        db.session.rollback()
        return jsonify({'success': False, 'error': str(e)})

@app.route('/generate_baseline_chart_comparison/<int:comparison_id>', methods=['POST'])
@login_required
def generate_baseline_chart_comparison(comparison_id):
    """Генерация baseline диаграммы для сохраненного сравнения"""
    try:
        data = request.get_json()
        metric_type = data.get('metric_type')
        
        logger.info(f"Генерация baseline диаграммы для comparison_id={comparison_id}, metric_type={metric_type}")
        
        # Проверяем, что comparison принадлежит текущему пользователю
        comparison = Comparison.query.filter_by(id=comparison_id, user_id=session['user_id']).first()
        if not comparison:
            return jsonify({'success': False, 'error': 'Сравнение не найдено'})
        
        # Получаем baseline метрики
        baseline = BaselineMetrics.query.filter_by(
            comparison_id=comparison_id,
            metric_type=metric_type
        ).first()
        
        if not baseline:
            return jsonify({'success': False, 'error': 'Сначала сохраните стартовые показатели'})
        
        # Определяем тип метрики
        is_yandex = metric_type.startswith('yandex')
        is_top20 = metric_type.endswith('20')
        
        # Получаем реальное количество URL из файлов сравнения (из сессии)
        def get_stats(key):
            val = session.get(key, '{}')
            if isinstance(val, dict):
                return val
            return json.loads(val) if val else {}
        
        if is_yandex:
            stats1 = get_stats('stats_y20_1' if is_top20 else 'stats_y10_1')
            stats2 = get_stats('stats_y20_2' if is_top20 else 'stats_y10_2')
        else:
            stats1 = get_stats('stats_g20_1' if is_top20 else 'stats_g10_1')
            stats2 = get_stats('stats_g20_2' if is_top20 else 'stats_g10_2')
        
        real_total_urls = stats1.get('total_urls', baseline.total_urls or 100)
        
        # Формируем stats_start
        is_count_mode = baseline.input_type == 'count'
        stats_start = {}
        
        for category in ['client_site', 'positive', 'neutral', 'negative', 'irrelevant']:
            value = getattr(baseline, f'{category}_value', '') or ''
            # Приводим к строке, так как в БД может быть float
            value_str = str(value) if value else ''
            
            if is_count_mode:
                count = int(float(value_str)) if value_str else 0
                percentage = round((count / real_total_urls) * 100, 1) if real_total_urls > 0 else 0
            else:
                percentage = float(value_str.replace('%', '').strip()) if value_str else 0
                count = round((percentage / 100) * real_total_urls)
            
            stats_start[category] = {'count': count, 'percentage': percentage}
        
        # Определяем session_id для сохранения диаграммы
        session_id = session.get('session_id') or str(uuid.uuid4())
        if 'session_id' not in session:
            session['session_id'] = session_id
        
        session_folder = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        os.makedirs(session_folder, exist_ok=True)
        
        # Генерируем диаграмму
        chart_buffer = comparator.create_baseline_chart(
            stats1, stats2, stats_start,
            title='',
            file1_name=comparison.file1_name,
            file2_name=comparison.file2_name
        )
        
        if chart_buffer:
            chart_filename = f"baseline_{metric_type}_{comparison_id}.png"
            chart_path = os.path.join(session_folder, chart_filename)
            
            with open(chart_path, 'wb') as f:
                f.write(chart_buffer.getvalue())
            
            # Обновляем путь к диаграмме в comparison
            if metric_type == 'yandex_top20':
                comparison.chart_start_y20_path = chart_path
            elif metric_type == 'yandex_top10':
                comparison.chart_start_y10_path = chart_path
            elif metric_type == 'google_top20':
                comparison.chart_start_g20_path = chart_path
            elif metric_type == 'google_top10':
                comparison.chart_start_g10_path = chart_path
            
            db.session.commit()
            
            chart_url = f"/uploads/{session_id}/{chart_filename}"
            
            return jsonify({
                'success': True,
                'chart_url': chart_url,
                'chart_path': chart_path,
                'message': 'Диаграмма создана'
            })
        else:
            return jsonify({'success': False, 'error': 'Не удалось создать диаграмму'})
            
    except Exception as e:
        logger.error(f"Ошибка генерации baseline диаграммы: {e}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/download_report/')
def download_report():
    path = request.args.get('path')
    logger.info(f"Download request: path={path}")
    logger.info(f"PERMANENT_FOLDER={app.config['PERMANENT_FOLDER']}")
    safe_path = validate_file_path(path, app.config['PERMANENT_FOLDER'])
    logger.info(f"Safe path result: {safe_path}")
    if safe_path:
        logger.info(f"File exists: {os.path.exists(safe_path)}")
    if not safe_path or not os.path.exists(safe_path):
        flash('Файл не найден', 'error')
        return redirect(url_for('profile'))
    mimetype = 'application/octet-stream'
    if safe_path.endswith('.png'):
        mimetype = 'image/png'
    elif safe_path.endswith('.txt'):
        mimetype = 'text/plain'
    elif safe_path.endswith('.xlsx'):
        mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    elif safe_path.endswith('.pptx'):
        mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    return send_file(safe_path, as_attachment=True, download_name=os.path.basename(safe_path), mimetype=mimetype)

@app.route('/delete_comparison/<int:comp_id>', methods=['POST'])
def delete_comparison(comp_id):
    if 'user_id' not in session:
        flash('Необходимо войти')
        return redirect(url_for('login'))
    comp = Comparison.query.get_or_404(comp_id)
    if comp.user_id != session['user_id']:
        flash('Нет прав на удаление')
        return redirect(url_for('profile'))

    user_perm_dir = os.path.join(app.config['PERMANENT_FOLDER'], str(session['user_id']))
    comp_dir = os.path.join(user_perm_dir, str(comp.id))
    if os.path.exists(comp_dir):
        shutil.rmtree(comp_dir)

    db.session.delete(comp)
    db.session.commit()
    flash('Анализ удалён')
    return redirect(url_for('profile'))

@app.route('/clear_history', methods=['POST'])
def clear_history():
    if 'user_id' not in session:
        flash('Необходимо войти')
        return redirect(url_for('login'))
    user_id = session['user_id']
    comparisons = Comparison.query.filter_by(user_id=user_id).all()
    for comp in comparisons:
        user_perm_dir = os.path.join(app.config['PERMANENT_FOLDER'], str(user_id))
        comp_dir = os.path.join(user_perm_dir, str(comp.id))
        if os.path.exists(comp_dir):
            shutil.rmtree(comp_dir)
        db.session.delete(comp)
    db.session.commit()
    flash('Вся история очищена')
    return redirect(url_for('profile'))

# ---------- ОБРАТНАЯ СВЯЗЬ ----------
@app.route('/feedback', methods=['GET', 'POST'])
def feedback():
    if request.method == 'POST':
        message = request.form.get('message', '').strip()
        email = request.form.get('email', '').strip() if 'email' in request.form else session.get('user_email', '')
        if not message:
            flash('Введите сообщение', 'error')
            return redirect(url_for('feedback'))
        fb = Feedback(
            user_id=session.get('user_id'),
            user_email=email,
            message=message
        )
        db.session.add(fb)
        db.session.commit()
        flash('Спасибо за обратную связь!', 'success')
        return redirect(url_for('index'))
    return render_template('feedback.html')

@app.route('/admin_feedback')
def admin_feedback():
    if 'user_id' not in session:
        flash('Необходимо войти', 'error')
        return redirect(url_for('login'))
    user = db.session.get(User, session['user_id'])
    if not user or not user.is_admin:
        flash('Нет прав доступа', 'error')
        return redirect(url_for('profile'))
    feedbacks = Feedback.query.order_by(Feedback.created_at.desc()).all()
    return render_template('admin_feedback.html', feedbacks=feedbacks)

@app.route('/update_feedback_status/<int:fb_id>/<status>')
def update_feedback_status(fb_id, status):
    if 'user_id' not in session:
        flash('Необходимо войти', 'error')
        return redirect(url_for('login'))
    user = db.session.get(User, session['user_id'])
    if not user or not user.is_admin:
        flash('Нет прав доступа', 'error')
        return redirect(url_for('profile'))
    fb = Feedback.query.get_or_404(fb_id)
    if status in ['new', 'read', 'replied']:
        fb.status = status
        db.session.commit()
        flash('Статус обновлён', 'success')
    else:
        flash('Неверный статус', 'error')
    return redirect(url_for('admin_feedback'))

# ---------- УДАЛЕНИЕ СООБЩЕНИЙ ----------
@app.route('/delete_feedback/<int:fb_id>', methods=['POST'])
def delete_feedback(fb_id):
    if 'user_id' not in session:
        flash('Необходимо войти', 'error')
        return redirect(url_for('login'))
    user = db.session.get(User, session['user_id'])
    if not user or not user.is_admin:
        flash('Нет прав доступа', 'error')
        return redirect(url_for('profile'))
    fb = Feedback.query.get_or_404(fb_id)
    db.session.delete(fb)
    db.session.commit()
    flash('Сообщение удалено', 'success')
    return redirect(url_for('admin_feedback'))

@app.route('/clear_feedback', methods=['POST'])
def clear_feedback():
    if 'user_id' not in session:
        flash('Необходимо войти', 'error')
        return redirect(url_for('login'))
    user = db.session.get(User, session['user_id'])
    if not user or not user.is_admin:
        flash('Нет прав доступа', 'error')
        return redirect(url_for('profile'))
    Feedback.query.delete()
    db.session.commit()
    flash('Все сообщения удалены', 'success')
    return redirect(url_for('admin_feedback'))

# -------------------------------------------------

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)